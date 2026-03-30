#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
api.py
API REST del Motor de Traducció Automàtica Neuronal castellà→valencià (TANEU)
Servei de Llengües i Política Lingüística - Universitat de València

Endpoints:
  POST /tradueix               — Traducció/correcció de text pla
  POST /tradueix-document      — Traducció de fitxers .docx i .pptx
  POST /desa-postedicio        — Desa una postedició humana al corpus
  GET  /estadistiques          — Estadístiques d'ús de la sessió
  GET  /salut                  — Estat del servei

Inici:
  uvicorn api:app --host 0.0.0.0 --port 8000 --reload
"""

import base64
import csv
import hashlib
import io
import json
import logging
import os
import re
import secrets
import shutil
import sys
import tempfile
import unicodedata
import time
import uuid
import zipfile
from datetime import date, datetime
from pathlib import Path
from typing import Any, Dict, List, Literal, Optional

# ─── FastAPI i dependències web ───────────────────────────────────────────────
from fastapi import Body, FastAPI, File, Form, HTTPException, Request, UploadFile
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import Response, StreamingResponse
from pydantic import BaseModel, Field

# ─── Mòduls locals del projecte ──────────────────────────────────────────────
sys.path.insert(0, str(Path(__file__).parent))
from system_prompts import (
    construeix_prompt_correccio,
    construeix_prompt_traduccio_es_va,
    construeix_prompt_traduccio_en_va,
    construeix_prompt_revisio,
)
from format_preserving_editor import (
    DocxFormatPreservingEditor,
    PptxFormatPreservingEditor,
    processa_document_traduccio,
    processa_document_correccio,
)

# ─── python-dotenv: càrrega i persistència de claus API ──────────────────────
try:
    from dotenv import load_dotenv, set_key as _dotenv_set_key
    _DOTENV_OK = True
except ImportError:
    _DOTENV_OK = False

# ─── Rutes del projecte ───────────────────────────────────────────────────────
ARREL_PROJECTE  = Path(__file__).parent.parent
DIR_MODELS      = ARREL_PROJECTE / "model-afinar"
DIR_LOGS        = ARREL_PROJECTE / "logs"
DIR_POSTEDICIONS = ARREL_PROJECTE / "corpus" / "postedicions"
FITXER_LOG      = DIR_LOGS / "api.log"
FITXER_STATS    = DIR_LOGS / "estadistiques.json"

# ─── Fitxer .env: carrega claus API de forma persistent ──────────────────────
import os as _os_env
_ENV_PATH = Path(__file__).resolve().parent.parent / ".env"
_ENV_PATH.touch(exist_ok=True)  # Crea el fitxer si no existeix
if _DOTENV_OK:
    load_dotenv(dotenv_path=str(_ENV_PATH), override=False)

# Neteja preventiva: elimina cometes que python-dotenv pugui haver llegit del .env
import os as _os_clean
for _var_env in ("GEMINI_API_KEY", "ANTHROPIC_API_KEY_CORRECCIO"):
    _val_env = _os_clean.environ.get(_var_env, "")
    _val_net = _val_env.strip().strip("'\"")
    if _val_net != _val_env:
        _os_clean.environ[_var_env] = _val_net

# ─── Funcions auxiliars per obtenir claus API de forma fiable ─────────────────

def _obte_api_key_anthropic() -> str:
    """
    Obté la clau API d'Anthropic de forma fiable.
    Prioritat: os.environ → lectura directa del .env
    Necessari perquè --reload de uvicorn en Windows pot no heretar
    les variables d'entorn correctament al procés fill.
    """
    clau = (
        os.environ.get("ANTHROPIC_API_KEY_CORRECCIO", "").strip().strip("'\"")
        or os.environ.get("ANTHROPIC_API_KEY", "").strip().strip("'\"")
    )
    if clau:
        return clau

    # Fallback: llegeix directament del .env
    env_path = Path(__file__).resolve().parent.parent / ".env"
    if env_path.exists():
        try:
            for linia in env_path.read_text(encoding="utf-8").splitlines():
                linia = linia.strip()
                if linia.startswith("ANTHROPIC_API_KEY_CORRECCIO="):
                    clau = linia.split("=", 1)[1].strip().strip("'\"")
                    if clau:
                        os.environ["ANTHROPIC_API_KEY_CORRECCIO"] = clau
                        return clau
        except Exception:
            pass
    return ""


def _obte_api_key_gemini() -> str:
    """
    Obté la clau API de Gemini de forma fiable.
    Prioritat: os.environ → lectura directa del .env
    """
    clau = os.environ.get("GEMINI_API_KEY", "").strip().strip("'\"")
    if clau:
        return clau

    env_path = Path(__file__).resolve().parent.parent / ".env"
    if env_path.exists():
        try:
            for linia in env_path.read_text(encoding="utf-8").splitlines():
                linia = linia.strip()
                if linia.startswith("GEMINI_API_KEY="):
                    clau = linia.split("=", 1)[1].strip().strip("'\"")
                    if clau:
                        os.environ["GEMINI_API_KEY"] = clau
                        return clau
        except Exception:
            pass
    return ""


# ─── Constants ────────────────────────────────────────────────────────────────
VERSIO          = "1.0"
MODEL_ID        = "projecte-aina/aina-translator-es-ca"
MAX_MIDA_FITXER = 150 * 1024 * 1024         # 150 MB en bytes
EXTENSIONS_OK   = {".docx", ".pptx"}

# ─── Glossaris ────────────────────────────────────────────────────────────────
DIR_GLOSSARIS = ARREL_PROJECTE / "glossaris"
DIR_GLOSSARIS.mkdir(exist_ok=True)

DOMINIS = [
    "Art i Història de l'Art",
    "Astronomia",
    "Altres",
    "Biologia",
    "Ciències Polítiques",
    "Comunicacions institucionals i discursos",
    "Convenis",
    "Convocatòries: cursos, premis, beques, concursos",
    "Dret",
    "Ecologia i Medi Ambient",
    "Economia",
    "Enginyeries",
    "Farmàcia",
    "Filologia i Lingüística",
    "Filosofia",
    "Física",
    "Formació del professorat i Ciències de l'Educació",
    "Formularis",
    "Geografia",
    "Història i Antropologia",
    "Informàtica i Noves Tecnologies",
    "Logopèdia",
    "Matemàtiques i Estadística",
    "Medicina i Infermeria",
    "Música i Arts Escèniques",
    "Notes de Premsa",
    "Odontologia",
    "Pedagogia",
    "Psicologia",
    "Química",
    "Salut Laboral i Prevenció de Riscos",
    "Textos administratius",
]

# Mapa de noms de domini obsolets → nous (per a migració de fitxers TSV)
_MIGRACIONS_DOMINIS = {
    "Medi_Ambient": "Ecologia_i_Medi_Ambient",
    "Música":       "Música_i_Arts_Escèniques",
    "Informàtica":  "Informàtica_i_Noves_Tecnologies",
}


def _migra_glossaris_antics():
    """Renombra fitxers de glossari amb noms de domini obsolets."""
    if not DIR_GLOSSARIS.exists():
        return
    for nom_antic, nom_nou in _MIGRACIONS_DOMINIS.items():
        fitxer_antic = DIR_GLOSSARIS / f"{nom_antic}.tsv"
        fitxer_nou   = DIR_GLOSSARIS / f"{nom_nou}.tsv"
        if fitxer_antic.exists() and not fitxer_nou.exists():
            fitxer_antic.rename(fitxer_nou)
            log.info("Glossari migrat: %s → %s", fitxer_antic.name, fitxer_nou.name)


_migra_glossaris_antics()


def nom_fitxer_glossari(domini: str) -> Path:
    """Retorna la ruta del fitxer TSV per a un domini."""
    nom = re.sub(r'[^\w\s-]', '', domini, flags=re.UNICODE)
    nom = re.sub(r'\s+', '_', nom.strip())
    return DIR_GLOSSARIS / f"{nom}.tsv"


def carrega_glossari(domini: str) -> list[dict]:
    """Carrega les entrades d'un glossari des del disc."""
    path = nom_fitxer_glossari(domini)
    if not path.exists():
        return []
    entrades = []
    with open(path, encoding='utf-8', newline='') as f:
        reader = csv.DictReader(f, delimiter='\t')
        for row in reader:
            if row.get('es') and row.get('ca'):
                entrades.append({
                    'es':     row['es'].strip(),
                    'ca':     row['ca'].strip(),
                    'tecnic': row.get('tecnic', '').strip(),
                    'data':   row.get('data', '').strip(),
                    'domini': domini,
                })
    return entrades


def carrega_glossari_com_diccionari(domini: str) -> dict[str, str]:
    """
    Retorna el glossari d'un domini com a diccionari {terme_es: terme_ca}.
    Ordena els termes de més llarg a més curt per evitar substitucions parcials.
    """
    if not domini or domini not in DOMINIS:
        return {}
    entrades = carrega_glossari(domini)
    # Ordena de més llarg a més curt per prioritzar expressions multiparaula
    entrades_ordenades = sorted(entrades, key=lambda e: len(e["es"]), reverse=True)
    return {e["es"].strip(): e["ca"].strip() for e in entrades_ordenades if e["es"] and e["ca"]}


def aplica_glossari_al_text(text: str, glossari: dict[str, str]) -> str:
    """
    Aplica les substitucions del glossari a un text.
    Fa substitucions case-insensitive però preserva la capitalització original.
    """
    import re as _re_gl
    if not glossari or not text:
        return text

    resultat = text
    for terme_es, terme_ca in glossari.items():
        if not terme_es:
            continue
        # Substitueix respectant límits de paraula, case-insensitive
        patró = _re_gl.compile(_re_gl.escape(terme_es), _re_gl.IGNORECASE)
        def substitueix_preservant_capitalitzacio(match, _ca=terme_ca):
            original = match.group(0)
            # Si l'original comença en majúscula, capitalitza la substitució
            if original and original[0].isupper():
                return _ca[0].upper() + _ca[1:] if _ca else _ca
            return _ca
        resultat = patró.sub(substitueix_preservant_capitalitzacio, resultat)
    return resultat


def desa_glossari(domini: str, entrades: list[dict]) -> None:
    """Desa les entrades d'un glossari al disc."""
    path = nom_fitxer_glossari(domini)
    with open(path, 'w', encoding='utf-8', newline='') as f:
        writer = csv.DictWriter(
            f, fieldnames=['es', 'ca', 'tecnic', 'data', 'domini'],
            delimiter='\t',
        )
        writer.writeheader()
        writer.writerows(entrades)


def genera_nom_traduit(nom_original: str) -> str:
    """
    Genera el nom del fitxer traduït afegint el sufix _VAL.
    (Funció antiga conservada per compatibilitat. Usa genera_nom_arxiu().)
    """
    return genera_nom_arxiu(nom_original, sufix="VAL")


def genera_nom_arxiu(nom_original: str, domini: str = "", sufix: str = "VAL") -> str:
    """
    Genera el nom de l'arxiu traduït/corregit seguint la regla:
    nom_sense_accents_DOMINI_SUFIX.ext

    - Elimina accents i caràcters especials
    - Substitueix espais per guions baixos
    - Afig el domini en majúscules (si n'hi ha)
    - Afig el sufix (VAL per defecte)

    Exemples:
      genera_nom_arxiu("Hàbitats de la papallona.pptx", "Biologia", "VAL")
        → "Habitats_de_la_papallona_BIOLOGIA_VAL.pptx"
      genera_nom_arxiu("Acta reunió.docx", "Textos administratius", "CORR_VAL")
        → "Acta_reunio_TEXTOS_ADMINISTRATIUS_CORR_VAL.docx"
      genera_nom_arxiu("Conferencia María.docx", "", "VAL")
        → "Conferencia_Maria_VAL.docx"
    """
    path = Path(nom_original)
    nom_base = path.stem
    extensio = path.suffix

    # Elimina accents: à→a, é→e, ñ→n, ç→c, etc.
    nom_net = unicodedata.normalize('NFKD', nom_base)
    nom_net = ''.join(c for c in nom_net if not unicodedata.combining(c))

    # Substitueix caràcters especials per guió baix
    nom_net = re.sub(r'[^a-zA-Z0-9]', '_', nom_net)

    # Elimina guions baixos múltiples
    nom_net = re.sub(r'_+', '_', nom_net)

    # Elimina guions baixos al principi i al final
    nom_net = nom_net.strip('_')

    # Construeix el nom final
    parts = [nom_net]
    if domini:
        domini_net = unicodedata.normalize('NFKD', domini.upper())
        domini_net = ''.join(c for c in domini_net if not unicodedata.combining(c))
        domini_net = re.sub(r'[^A-Z0-9]', '_', domini_net)
        domini_net = re.sub(r'_+', '_', domini_net).strip('_')
        parts.append(domini_net)
    if sufix:
        parts.append(sufix)

    return '_'.join(parts) + extensio


# ─── Configuració del logging ─────────────────────────────────────────────────
DIR_LOGS.mkdir(parents=True, exist_ok=True)

logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(name)s — %(message)s",
    datefmt="%Y-%m-%d %H:%M:%S",
    handlers=[
        logging.FileHandler(FITXER_LOG, encoding="utf-8"),
        logging.StreamHandler(),
    ],
)
log = logging.getLogger("taneu.api")

# ─── Import del mòdul de preservació de format (estratègia XML directa) ───────
# Afegim l'arrel del projecte al path perquè 'scripts' siga trobable tant
# quan s'executa des de webapp/ com des de l'arrel (uvicorn webapp.api:app)
_arrel_sys = str(ARREL_PROJECTE)
if _arrel_sys not in sys.path:
    sys.path.insert(0, _arrel_sys)

try:
    from scripts.preserva_xml import tradueix_document as _tradueix_document_xml
    log.info("Mòdul preserva_xml carregat correctament.")
except ImportError as _err_import:
    _tradueix_document_xml = None  # type: ignore[assignment]
    log.warning("No s'ha pogut importar preserva_xml: %s", _err_import)

# ─── Estat global de l'aplicació ──────────────────────────────────────────────
_tokenizer  = None
_model      = None


def _carrega_postedicions_count() -> int:
    """Compta els fitxers de postedició existents al disc."""
    try:
        return len(list(DIR_POSTEDICIONS.glob("*.jsonl")))
    except Exception:
        return 0


_stats: dict = {
    "data":          str(date.today()),
    "paraules_avui": 0,
    "fitxers_avui":  0,
    "postedicions":  _carrega_postedicions_count(),
}


def _reinicia_stats_si_cal() -> None:
    """Reinicia els comptadors diaris si ha canviat el dia."""
    avui = str(date.today())
    if _stats["data"] != avui:
        _stats.update({
            "data":          avui,
            "paraules_avui": 0,
            "fitxers_avui":  0,
        })
        log.info("Comptadors diaris reiniciats per al dia %s", avui)


# NOTA: _get_model() ja no s'utilitza perquè _tradueix_text() crida el servidor CTranslate2 (port 5001).
def _get_model():
    """Carrega el model i el tokenitzador de forma mandrosa (lazy loading)."""
    global _tokenizer, _model
    if _tokenizer is None or _model is None:
        try:
            from transformers import MarianMTModel, MarianTokenizer
            origen = str(DIR_MODELS) if DIR_MODELS.exists() else MODEL_ID
            log.info("Carregant model des de: %s", origen)
            _tokenizer = MarianTokenizer.from_pretrained(origen)
            _model     = MarianMTModel.from_pretrained(origen)
            log.info("Model carregat correctament.")
        except ImportError:
            log.error("Paquet 'transformers' no instal·lat.")
            raise HTTPException(
                status_code=503,
                detail="El servei de traducció no està disponible: "
                       "instal·la les dependències (pip install -r requirements.txt).",
            )
        except Exception as exc:
            log.exception("Error carregant el model: %s", exc)
            raise HTTPException(
                status_code=503,
                detail=f"No s'ha pogut carregar el model de traducció: {exc}",
            )
    return _tokenizer, _model


def _tradueix_text(text: str) -> str:
    """Tradueix un text castellà→català usant el servidor CTranslate2 (port 5001)."""
    import requests as _req
    paragrafs = text.split("\n")
    resultats = []
    for paragraf in paragrafs:
        paragraf = paragraf.strip()
        if not paragraf:
            resultats.append("")
            continue
        try:
            resp = _req.post(
                "http://127.0.0.1:5001/translate",
                json={"text": paragraf, "src": "es", "tgt": "ca"},
                timeout=60,
            )
            resp.raise_for_status()
            resultats.append(resp.json().get("translation", paragraf))
        except Exception:
            resultats.append(paragraf)
    return "\n".join(resultats)

def _compta_paraules(text: str) -> int:
    """Retorna el nombre de paraules d'un text."""
    return len(text.split())


# ─── Creació de l'aplicació FastAPI ───────────────────────────────────────────
app = FastAPI(
    title       = "TANEU — Motor de Traducció Automàtica Neuronal",
    description = (
        "API REST del Motor de Traducció Automàtica Neuronal castellà→valencià "
        "del Servei de Llengües i Política Lingüística de la Universitat de València."
    ),
    version     = VERSIO,
    docs_url    = "/docs",
    redoc_url   = "/redoc",
)

# CORS: permet accés des de Netlify i qualsevol origen
# expose_headers és IMPRESCINDIBLE perquè el navegador puga llegir
# Content-Disposition (nom del fitxer) en respostes fetch() cross-origin.
app.add_middleware(
    CORSMiddleware,
    allow_origins     = ["*"],
    allow_credentials = False,          # ha de ser False quan allow_origins=["*"]
    allow_methods     = ["*"],
    allow_headers     = ["*"],
    expose_headers    = [               # capçaleres llegibles per JavaScript
        "Content-Disposition",
        "X-Paraules-Traduides",
        "X-Temps-Ms",
        "X-Num-Imatges",
        "X-Motor",
        "X-Direccio",
    ],
)


# ─── Esquemes Pydantic ────────────────────────────────────────────────────────

class PeticioTradueix(BaseModel):
    text: str = Field(
        ...,
        min_length  = 1,
        max_length  = 50_000,
        description = "Text en castellà a traduir o corregir.",
        examples    = ["El sistema universitario valenciano necesita mejoras."],
    )
    mode: Literal["traduccio", "correccio"] = Field(
        default     = "traduccio",
        description = "'traduccio' (es→ca) o 'correccio' (millora del text en valencià).",
    )


class RespostaTradueix(BaseModel):
    traduccio: str
    paraules:  int
    temps_ms:  int


class PeticioPostedicio(BaseModel):
    origen:      str = Field(..., description="Text original en castellà.")
    ta:          str = Field(..., description="Traducció automàtica generada pel sistema.")
    posteditada: str = Field(..., description="Traducció corregida pel tècnic.")
    tecnic:      str = Field(..., description="Identificador o nom del tècnic.")


class RespostaPostedicio(BaseModel):
    estat: str
    id:    str


class RespostaEstadistiques(BaseModel):
    paraules_avui: int
    fitxers_avui:  int
    postedicions:  int


class RespostaSalut(BaseModel):
    estat:  str
    model:  str
    versio: str


# ═══════════════════════════════════════════════════════════════════════════════
# AUTENTICACIÓ D'USUARIS
# ═══════════════════════════════════════════════════════════════════════════════

# Emmagatzematge de sessions (en memòria, es perd al reiniciar)
_sessions: dict[str, str] = {}  # token → username

_USERS_PATH = Path(__file__).parent.parent / "users.json"


def _carrega_usuaris() -> dict:
    if _USERS_PATH.exists():
        with open(_USERS_PATH, "r", encoding="utf-8") as f:
            return json.load(f)
    return {}


def _desa_usuaris(usuaris: dict) -> None:
    with open(_USERS_PATH, "w", encoding="utf-8") as f:
        json.dump(usuaris, f, ensure_ascii=False, indent=2)


def _valida_sessio(request: Request) -> str | None:
    """Retorna el username si la sessió és vàlida, None si no."""
    token = request.headers.get("Authorization", "").replace("Bearer ", "")
    if not token:
        token = request.query_params.get("token", "")
    return _sessions.get(token)


def _requereix_admin(request: Request) -> None:
    """Llança 403 si l'usuari no és admin."""
    username = _valida_sessio(request)
    if not username:
        raise HTTPException(status_code=401, detail="Sessió no vàlida. Identifiqueu-vos.")
    usuaris = _carrega_usuaris()
    if username not in usuaris or usuaris[username].get("rol") != "admin":
        raise HTTPException(status_code=403, detail="Accés restringit a l'administrador.")


# ─── Login ───────────────────────────────────────────────────────────────────

class PeticioLogin(BaseModel):
    username: str
    password: str


@app.post("/auth/login", tags=["Autenticació"])
async def login(peticio: PeticioLogin):
    usuaris = _carrega_usuaris()
    user = usuaris.get(peticio.username)
    if not user or not user.get("actiu", False):
        raise HTTPException(status_code=401, detail="Usuari o contrasenya incorrectes.")
    password_hash = hashlib.sha256(peticio.password.encode()).hexdigest()
    if user["password_hash"] != password_hash:
        raise HTTPException(status_code=401, detail="Usuari o contrasenya incorrectes.")
    token = secrets.token_hex(32)
    _sessions[token] = peticio.username
    return {
        "token": token,
        "username": peticio.username,
        "nom": user.get("nom", ""),
        "rol": user.get("rol", "user"),
    }


# ─── Validar sessió ─────────────────────────────────────────────────────────

@app.get("/auth/validar", tags=["Autenticació"])
async def validar_sessio(request: Request):
    username = _valida_sessio(request)
    if not username:
        raise HTTPException(status_code=401, detail="Sessió no vàlida.")
    usuaris = _carrega_usuaris()
    user = usuaris.get(username, {})
    return {
        "username": username,
        "nom": user.get("nom", ""),
        "rol": user.get("rol", "user"),
    }


# ─── Logout ──────────────────────────────────────────────────────────────────

@app.post("/auth/logout", tags=["Autenticació"])
async def logout(request: Request):
    token = request.headers.get("Authorization", "").replace("Bearer ", "")
    _sessions.pop(token, None)
    return {"ok": True}


# ─── Admin: llistar usuaris ──────────────────────────────────────────────────

@app.get("/auth/usuaris", tags=["Administració"])
async def llistar_usuaris(request: Request):
    _requereix_admin(request)
    usuaris = _carrega_usuaris()
    return [
        {
            "username": u,
            "nom": d.get("nom", ""),
            "rol": d.get("rol", "user"),
            "actiu": d.get("actiu", False),
        }
        for u, d in usuaris.items()
    ]


# ─── Admin: crear usuari ────────────────────────────────────────────────────

class NouUsuari(BaseModel):
    username: str
    password: str
    nom: str = ""
    rol: str = "user"


@app.post("/auth/usuaris", tags=["Administració"])
async def crear_usuari(peticio: NouUsuari, request: Request):
    _requereix_admin(request)
    usuaris = _carrega_usuaris()
    if peticio.username in usuaris:
        raise HTTPException(status_code=409, detail="L'usuari ja existeix.")
    usuaris[peticio.username] = {
        "password_hash": hashlib.sha256(peticio.password.encode()).hexdigest(),
        "nom": peticio.nom,
        "rol": peticio.rol,
        "actiu": True,
    }
    _desa_usuaris(usuaris)
    return {"ok": True, "username": peticio.username}


# ─── Admin: eliminar usuari ─────────────────────────────────────────────────

@app.delete("/auth/usuaris/{username}", tags=["Administració"])
async def eliminar_usuari(username: str, request: Request):
    _requereix_admin(request)
    usuaris = _carrega_usuaris()
    if username not in usuaris:
        raise HTTPException(status_code=404, detail="Usuari no trobat.")
    if username == "coitor":
        raise HTTPException(status_code=403, detail="No es pot eliminar l'administrador.")
    del usuaris[username]
    _desa_usuaris(usuaris)
    return {"ok": True}


# ─── Admin: canviar contrasenya d'un usuari ──────────────────────────────────

class CanviPassword(BaseModel):
    nova_password: str


@app.put("/auth/usuaris/{username}/password", tags=["Administració"])
async def canviar_password(username: str, peticio: CanviPassword, request: Request):
    _requereix_admin(request)
    usuaris = _carrega_usuaris()
    if username not in usuaris:
        raise HTTPException(status_code=404, detail="Usuari no trobat.")
    usuaris[username]["password_hash"] = hashlib.sha256(peticio.nova_password.encode()).hexdigest()
    _desa_usuaris(usuaris)
    return {"ok": True}


# ─── Canvi de contrasenya propi (qualsevol usuari autenticat) ─────────────────

class CanviPasswordPropi(BaseModel):
    password_actual: str
    password_nova: str


@app.put("/auth/canvi-password", tags=["Autenticació"])
async def canvi_password_propi(peticio: CanviPasswordPropi, request: Request):
    """Permet a qualsevol usuari autenticat canviar la seua pròpia contrasenya."""
    username = _valida_sessio(request)
    if not username:
        raise HTTPException(status_code=401, detail="Sessió no vàlida.")

    usuaris = _carrega_usuaris()
    user = usuaris.get(username)
    if not user:
        raise HTTPException(status_code=404, detail="Usuari no trobat.")

    # Verifica la contrasenya actual
    hash_actual = hashlib.sha256(peticio.password_actual.encode()).hexdigest()
    if user["password_hash"] != hash_actual:
        raise HTTPException(status_code=403, detail="La contrasenya actual és incorrecta.")

    # Actualitza la contrasenya
    usuaris[username]["password_hash"] = hashlib.sha256(peticio.password_nova.encode()).hexdigest()
    _desa_usuaris(usuaris)

    return {"ok": True, "missatge": "Contrasenya actualitzada correctament."}


# ─── Endpoint: GET /salut ─────────────────────────────────────────────────────

@app.get(
    "/salut",
    response_model = RespostaSalut,
    summary        = "Comprova l'estat del servei",
    tags           = ["Sistema"],
)
async def salut() -> RespostaSalut:
    """Retorna l'estat operatiu de l'API i la versió."""
    log.debug("GET /salut — petició rebuda")
    origen_model = str(DIR_MODELS) if DIR_MODELS.exists() else MODEL_ID
    return RespostaSalut(
        estat  = "actiu",
        model  = origen_model,
        versio = VERSIO,
    )


@app.get("/health")
async def health():
    return await salut()


# ─── Endpoint: POST /translate (àlies per al frontend) ───────────────────────
# Accepta el format del frontend: {"text": "...", "src": "es", "tgt": "ca"}
# Retorna el format del frontend:  {"translation": "..."}

@app.post("/translate", tags=["Traducció"],
          summary="Àlies /tradueix compatible amb el frontend")
async def translate(peticio: dict):
    """
    Àlies de /tradueix amb format compatible amb config.js:
      Entrada: {"text": "...", "src": "es", "tgt": "ca"}
      Eixida:  {"translation": "...", "temps_ms": ...}
    """
    text = (peticio.get("text") or "").strip()
    if not text:
        raise HTTPException(status_code=400,
                            detail="El camp 'text' és obligatori.")
    _reinicia_stats_si_cal()
    inici = time.perf_counter()
    try:
        resultat = _tradueix_text(text)
    except Exception as exc:
        log.exception("Error durant la traducció (/translate): %s", exc)
        raise HTTPException(status_code=500,
                            detail=f"Error intern: {exc}")
    temps_ms = int((time.perf_counter() - inici) * 1000)
    n_paraules = _compta_paraules(text)
    _stats["paraules_avui"] += n_paraules
    log.info("POST /translate — %d ms, %d paraules", temps_ms, n_paraules)
    return {"translation": resultat, "temps_ms": temps_ms}


# ─── Endpoint: POST /tradueix ─────────────────────────────────────────────────

@app.post(
    "/tradueix",
    response_model = RespostaTradueix,
    summary        = "Tradueix o corregeix text pla",
    tags           = ["Traducció"],
)
async def tradueix(peticio: PeticioTradueix) -> RespostaTradueix:
    """
    Tradueix text castellà→català (mode *traduccio*) o millora
    un text ja en valencià (mode *correccio*).
    """
    _reinicia_stats_si_cal()
    log.info("POST /tradueix — mode=%s paraules=%d", peticio.mode, _compta_paraules(peticio.text))

    if not peticio.text.strip():
        raise HTTPException(status_code=422, detail="El camp 'text' no pot estar buit.")

    inici = time.perf_counter()
    try:
        resultat = _tradueix_text(peticio.text)
    except HTTPException:
        raise
    except Exception as exc:
        log.exception("Error durant la traducció: %s", exc)
        raise HTTPException(
            status_code=500,
            detail=f"Error intern durant la traducció: {exc}",
        )
    temps_ms = int((time.perf_counter() - inici) * 1000)

    n_paraules = _compta_paraules(peticio.text)
    _stats["paraules_avui"] += n_paraules
    log.info("Traducció completada en %d ms (%d paraules)", temps_ms, n_paraules)

    return RespostaTradueix(
        traduccio = resultat,
        paraules  = n_paraules,
        temps_ms  = temps_ms,
    )


# ─── Endpoint: POST /tradueix-document ───────────────────────────────────────

@app.post(
    "/tradueix-document",
    summary = "Tradueix un document .docx o .pptx",
    tags    = ["Traducció"],
)
async def tradueix_document(
    fitxer:  UploadFile = File(..., description="Fitxer .docx o .pptx (màx. 150 MB)."),
    mode:    str        = Form(default="traduccio", description="'traduccio' o 'correccio'"),
    domini:  str        = Form(default="", description="Domini lingüístic per aplicar el glossari."),
) -> StreamingResponse:
    """
    Rep un fitxer .docx o .pptx, tradueix tot el text que conté
    i retorna el fitxer amb el **mateix nom i extensió** que l'original.
    """
    _reinicia_stats_si_cal()

    # ── Validació del fitxer ──────────────────────────────────────
    nom_original = fitxer.filename or "document"
    extensio     = Path(nom_original).suffix.lower()

    if extensio not in EXTENSIONS_OK:
        raise HTTPException(
            status_code = 415,
            detail      = (
                f"Extensió '{extensio}' no admesa. "
                f"Els formats acceptats són: {', '.join(sorted(EXTENSIONS_OK))}."
            ),
        )

    contingut = await fitxer.read()
    if len(contingut) > MAX_MIDA_FITXER:
        raise HTTPException(
            status_code = 413,
            detail      = (
                f"El fitxer supera la mida màxima de 150 MB "
                f"({len(contingut) / 1_048_576:.1f} MB rebuts)."
            ),
        )

    log.info(
        "POST /tradueix-document — fitxer='%s' mida=%.1f KB mode=%s",
        nom_original, len(contingut) / 1024, mode,
    )

    inici       = time.perf_counter()
    total_par   = 0
    buffer_eixida = io.BytesIO()

    mime_types = {
        ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    }
    mime_type = mime_types[extensio]

    # Carrega el glossari del domini si s'ha especificat
    glossari_domini = carrega_glossari_com_diccionari(domini) if domini else {}
    if glossari_domini:
        log.info(
            "Aplicant glossari '%s' (%d termes) al document '%s'",
            domini, len(glossari_domini), nom_original,
        )

    try:
        total_par = _tradueix_fitxer_xml(
            contingut, extensio.lstrip('.'), buffer_eixida, glossari_domini
        )
    except HTTPException:
        raise
    except Exception as exc:
        log.exception("Error traduint el document '%s': %s", nom_original, exc)
        raise HTTPException(
            status_code = 500,
            detail      = f"Error intern processant el document: {exc}",
        )

    temps_ms = int((time.perf_counter() - inici) * 1000)
    _stats["paraules_avui"] += total_par
    _stats["fitxers_avui"]  += 1

    nom_sortida = genera_nom_arxiu(nom_original, domini=domini, sufix="VAL")
    log.info(
        "Document traduït en %d ms (%d paraules) → '%s'",
        temps_ms, total_par, nom_sortida,
    )

    buffer_eixida.seek(0)
    return StreamingResponse(
        buffer_eixida,
        media_type = mime_type,
        headers    = {
            "Content-Disposition": f'attachment; filename="{nom_sortida}"',
            "X-Paraules-Traduides": str(total_par),
            "X-Temps-Ms":          str(temps_ms),
        },
    )


def _tradueix_fitxer_xml(
    contingut: bytes,
    extensio:  str,
    eixida:    io.BytesIO,
    glossari:  dict[str, str] | None = None,
) -> int:
    """
    Tradueix un .docx o .pptx usant manipulació directa de l'XML intern
    (scripts/preserva_xml.py). Retorna el nombre de paraules de l'original.

    Si s'especifica un glossari, aplica les substitucions terme_es→terme_ca
    a cada segment ABANS de passar-lo al motor de traducció.

    El recompte de paraules es fa llegint el ZIP directament amb lxml,
    sense necessitat de python-docx ni python-pptx.
    """
    if _tradueix_document_xml is None:
        raise HTTPException(
            status_code=503,
            detail="El mòdul preserva_xml no està disponible. "
                   "Comprova que scripts/preserva_xml.py existeix i que "
                   "lxml està instal·lat (pip install lxml).",
        )

    # ── Recompte de paraules de l'original via lxml ──────────────────────────
    import zipfile as _zf
    from lxml import etree as _et

    _NS_W = 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'
    _NS_A = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    # DOCX_REGEX / PPTX_REGEX es van moure dins de tradueix_document() (C4, v3).
    # Es defineixen localment aquí per al recompte de paraules.
    import re as _re
    _DOCX_REGEX = _re.compile(
        r'^word/(document|header\d+|footer\d+|footnotes|endnotes|comments)\.xml$'
    )
    _PPTX_REGEX = _re.compile(
        r'^ppt/(slides/slide\d+|notesSlides/notesSlide\d+)\.xml$'
    )

    regex  = _DOCX_REGEX if extensio == 'docx' else _PPTX_REGEX
    ns_t   = _NS_W      if extensio == 'docx' else _NS_A
    tag_t  = f'{{{ns_t}}}t'

    total_par = 0
    with _zf.ZipFile(io.BytesIO(contingut)) as z:
        for nom in z.namelist():
            if regex.match(nom):
                try:
                    arbre = _et.fromstring(z.read(nom))
                    for t in arbre.iter(tag_t):
                        if t.text:
                            total_par += _compta_paraules(t.text)
                except Exception:
                    pass

    # ── Funció de traducció (amb glossari opcional) ───────────────────────────
    if glossari:
        def _fn_tradueix_amb_glossari(text: str) -> str:
            # Pre-traducció: substitueix termes castellans pels valencians
            text_preprocessat = aplica_glossari_al_text(text, glossari)
            return _tradueix_text(text_preprocessat)
        fn_tradueix = _fn_tradueix_amb_glossari
    else:
        fn_tradueix = _tradueix_text

    # ── Traducció preservant format ──────────────────────────────────────────
    resultat = _tradueix_document_xml(contingut, extensio, fn_tradueix)
    eixida.write(resultat)
    eixida.seek(0)
    return total_par


# ─── OCR per a filtratge d'imatges amb text ─────────────────────────────────
_TESSERACT_OK = False
try:
    import pytesseract
    from PIL import Image as _PILImage
    # Configura la ruta de Tesseract per a Windows
    import platform as _platform
    if _platform.system() == 'Windows':
        _tesseract_paths = [
            r'C:\Program Files\Tesseract-OCR\tesseract.exe',
            r'C:\Program Files (x86)\Tesseract-OCR\tesseract.exe',
            r'C:\Users\santi\AppData\Local\Programs\Tesseract-OCR\tesseract.exe',
        ]
        for _tp in _tesseract_paths:
            if os.path.isfile(_tp):
                pytesseract.pytesseract.tesseract_cmd = _tp
                break
    # Prova que funciona
    pytesseract.get_tesseract_version()
    _TESSERACT_OK = True
    log.info("Tesseract OCR disponible: %s", pytesseract.get_tesseract_version())
except Exception as _e:
    log.warning("Tesseract OCR NO disponible: %s. El filtre d'imatges amb text NO funcionarà.", _e)
    _TESSERACT_OK = False


def _imatge_conte_text_real(dades_imatge: bytes, min_paraules: int = 1) -> tuple[bool, str]:
    """
    Analitza una imatge amb Tesseract OCR per determinar si conté text real.

    Usa image_to_string() amb --psm 3 (detecció automàtica), que és l'enfocament
    provat i verificat que funciona correctament amb imatges de documents.

    Una paraula vàlida és una seqüència d'almenys 3 caràcters exclusivament
    alfabètics (a-z, A-Z, amb o sense accents). Es descarten números,
    símbols, abreviatures d'1-2 lletres i soroll OCR.

    Retorna (True, resum_text) si la imatge conté almenys `min_paraules` paraules vàlides.
    Retorna (False, '') en cas contrari.
    """
    if not _TESSERACT_OK:
        log.warning("  Tesseract no disponible → imatge DESCARTADA per precaució")
        return False, ''

    try:
        img = _PILImage.open(io.BytesIO(dades_imatge))

        # Converteix RGBA/P a RGB amb fons blanc (evita problemes amb transparència)
        if img.mode == 'RGBA':
            fons = _PILImage.new('RGB', img.size, (255, 255, 255))
            fons.paste(img, mask=img.split()[3])
            img = fons
        elif img.mode not in ('RGB', 'L'):
            img = img.convert('RGB')

        # Redimensiona imatges molt grans per accelerar l'OCR
        max_dim = 2500
        if max(img.size) > max_dim:
            ratio = max_dim / max(img.size)
            img = img.resize((int(img.width * ratio), int(img.height * ratio)))

        # OCR amb detecció automàtica (--psm 3), que és el mode que funciona
        text_ocr = pytesseract.image_to_string(
            img,
            lang='spa+eng',
            timeout=15,
            config='--psm 3',
        )

    except Exception as exc:
        log.warning("  Error OCR: %s → imatge DESCARTADA", exc)
        return False, ''

    if not text_ocr or not text_ocr.strip():
        return False, ''

    # Extrau paraules candidates: 3+ caràcters exclusivament alfabètics
    import re as _re_ocr
    paraules_candidates = [w for w in text_ocr.split() if len(w) >= 3 and w.isalpha()]

    # Filtra paraules que NO són text real per a traduir
    _PARAULES_DESCARTAR = {
        # Copyright i metadades
        'copyright', 'rights', 'reserved', 'all', 'inc', 'ltd', 'corp',
        'publishing', 'published', 'pearson', 'elsevier', 'springer',
        'wiley', 'mcgraw', 'hill', 'benjamin', 'cummings',
        # Soroll OCR comú
        'the', 'and', 'for', 'that', 'this', 'with', 'from', 'are', 'was',
        'not', 'but', 'have', 'has', 'had', 'been', 'will', 'can', 'may',
    }

    def _es_paraula_traduible(p):
        p_lower = p.lower()
        # Descarta paraules de la llista de descart
        if p_lower in _PARAULES_DESCARTAR:
            return False
        # Descarta si totes les lletres són iguals (aaa, bbb)
        if len(set(p_lower)) <= 1:
            return False
        # Descarta si és probablement una sigla (3-4 lletres totes majúscules)
        if len(p) >= 3 and p.isupper() and len(p) <= 4:
            return False
        return True

    paraules_valides = [p for p in paraules_candidates if _es_paraula_traduible(p)]

    # Comprova si el text és NOMÉS copyright
    text_lower = text_ocr.lower()
    es_copyright = ('copyright' in text_lower or '©' in text_lower or 'rights reserved' in text_lower)
    if es_copyright and len(paraules_valides) < 5:
        return False, ''

    if len(paraules_valides) >= min_paraules:
        resum = ' '.join(paraules_valides[:10])
        return True, resum
    else:
        return False, ''


# ─── Endpoint: POST /extreu-imatges-document ─────────────────────────────────

@app.post(
    "/extreu-imatges-document",
    summary="Extreu les imatges incrustades d'un document .docx o .pptx",
    tags=["Traducció"],
)
async def extreu_imatges_document(
    fitxer: UploadFile = File(..., description="Fitxer .docx o .pptx"),
) -> Response:
    """
    Rep un document .docx o .pptx, n'extreu totes les imatges incrustades
    (.png, .jpg, .jpeg, .gif, .bmp, .tiff) i les retorna com un fitxer ZIP.
    """
    nom = fitxer.filename or "document"
    extensio = Path(nom).suffix.lower()
    if extensio not in (".docx", ".pptx"):
        raise HTTPException(status_code=415, detail="Només .docx i .pptx.")

    contingut = await fitxer.read()
    if len(contingut) > MAX_MIDA_FITXER:
        raise HTTPException(status_code=413, detail="Fitxer massa gran.")

    log.info("POST /extreu-imatges-document — '%s' %.1f KB", nom, len(contingut) / 1024)

    import zipfile as _zf

    extensions_imatge = ('.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff', '.tif')
    imatges = []
    total_imatges_analitzades = 0

    try:
        with _zf.ZipFile(io.BytesIO(contingut)) as z:
            for entry in z.namelist():
                entry_lower = entry.lower()
                # En .docx les imatges estan a word/media/
                # En .pptx les imatges estan a ppt/media/
                if '/media/' in entry_lower and any(entry_lower.endswith(ext) for ext in extensions_imatge):
                    total_imatges_analitzades += 1
                    nom_imatge = Path(entry).name
                    dades_imatge = z.read(entry)

                    # Filtre OCR: només inclou imatges amb text real
                    te_text, text_detectat = _imatge_conte_text_real(dades_imatge)
                    if te_text:
                        imatges.append((nom_imatge, dades_imatge))
                        log.info(
                            "  ✓ Imatge '%s' INCLOSA (text: %s)",
                            nom_imatge,
                            text_detectat[:80] if text_detectat else '(OCR no disponible)',
                        )
                    else:
                        log.info("  ✗ Imatge '%s' DESCARTADA (sense text real)", nom_imatge)
    except Exception as exc:
        log.exception("Error extraient imatges de '%s': %s", nom, exc)
        raise HTTPException(status_code=500, detail=f"Error extraient imatges: {exc}")

    if not imatges:
        if not _TESSERACT_OK:
            raise HTTPException(
                status_code=503,
                detail="Tesseract OCR no està instal·lat al servidor. "
                       "No es poden filtrar les imatges amb text. "
                       "Instal·la Tesseract: https://github.com/UB-Mannheim/tesseract/wiki",
            )
        raise HTTPException(
            status_code=404,
            detail="No s'han trobat imatges amb text real en aquest document. "
                   "S'han analitzat totes les imatges incrustades i cap conté "
                   "paraules completes en castellà, anglès o francès.",
        )

    log.info(
        "Anàlisi completada: %d imatges amb text de %d totals analitzades en '%s'",
        len(imatges), nom, total_imatges_analitzades,
    )

    # Crea un ZIP amb totes les imatges
    buffer_zip = io.BytesIO()
    nom_base = Path(nom).stem
    with _zf.ZipFile(buffer_zip, 'w', _zf.ZIP_DEFLATED) as zout:
        for i, (nom_img, dades_img) in enumerate(imatges, 1):
            # Renombra per claredat: imatge_01_nom.png
            ext_img = Path(nom_img).suffix
            nom_net = f"imatge_{i:02d}{ext_img}"
            zout.writestr(nom_net, dades_img)

        # Inclou un manifest JSON amb metadades
        import json as _json_manifest
        manifest = _json_manifest.dumps({
            "num_imatges": len(imatges),
            "total_analitzades": total_imatges_analitzades,
            "fitxer_original": nom,
        }, ensure_ascii=False)
        zout.writestr("_manifest.json", manifest)

    buffer_zip.seek(0)
    nom_zip = f"Imatges amb text_{nom_base}.zip"

    return Response(
        content=buffer_zip.getvalue(),
        media_type="application/zip",
        headers={
            "Content-Disposition": f'attachment; filename="{nom_zip}"',
            "Content-Length": str(len(buffer_zip.getvalue())),
            "X-Num-Imatges": str(len(imatges)),
        },
    )


# ─── Endpoint: GET /dominis-amb-glossari ─────────────────────────────────────

@app.get(
    "/dominis-amb-glossari",
    summary = "Llista de dominis amb informació del glossari",
    tags    = ["Glossari"],
)
async def dominis_amb_glossari():
    """Retorna la llista de dominis amb indicació de si tenen glossari."""
    resultat = []
    for domini in DOMINIS:
        path    = nom_fitxer_glossari(domini)
        entrades = carrega_glossari(domini) if path.exists() else []
        resultat.append({
            "domini":       domini,
            "te_glossari":  len(entrades) > 0,
            "num_entrades": len(entrades),
        })
    return {"dominis": resultat}


# ─── Endpoint: POST /recompte-paraules ───────────────────────────────────────

@app.post(
    "/recompte-paraules",
    summary = "Compta les paraules d'un document .docx o .pptx",
    tags    = ["Traducció"],
)
async def recompte_paraules(fitxer: UploadFile = File(...)):
    nom  = fitxer.filename or ""
    ext  = nom.rsplit(".", 1)[-1].lower() if "." in nom else ""
    contingut = await fitxer.read()
    buf  = io.BytesIO(contingut)
    buf.seek(0)
    total = 0
    if ext == "docx":
        from docx import Document
        doc = Document(buf)
        for para in doc.paragraphs:
            total += len(para.text.split())
        for taula in doc.tables:
            for fila in taula.rows:
                for cel in fila.cells:
                    for para in cel.paragraphs:
                        total += len(para.text.split())
    elif ext == "pptx":
        from pptx import Presentation
        prs = Presentation(buf)
        for diap in prs.slides:
            for shape in diap.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        total += len(para.text.split())
    return {"paraules": total, "fitxer": nom}


# ─── Endpoint: POST /desa-postedicio ─────────────────────────────────────────

@app.post(
    "/desa-postedicio",
    response_model = RespostaPostedicio,
    summary        = "Desa una postedició humana al corpus",
    tags           = ["Corpus"],
)
async def desa_postedicio(peticio: PeticioPostedicio) -> RespostaPostedicio:
    """
    Desa un parell (traducció automàtica, traducció corregida) al corpus
    de postedicions per a l'entrenament i l'avaluació del sistema.
    """
    _reinicia_stats_si_cal()

    if not peticio.origen.strip() or not peticio.posteditada.strip():
        raise HTTPException(
            status_code=422,
            detail="Els camps 'origen' i 'posteditada' no poden estar buits.",
        )

    id_entrada = str(uuid.uuid4())
    avui       = str(date.today())

    entrada = {
        "id":          id_entrada,
        "data":        avui,
        "tecnic":      peticio.tecnic.strip(),
        "origen":      peticio.origen.strip(),
        "ta":          peticio.ta.strip(),
        "posteditada": peticio.posteditada.strip(),
    }

    # Desa al fitxer JSONL del dia actual
    DIR_POSTEDICIONS.mkdir(parents=True, exist_ok=True)
    fitxer_dia = DIR_POSTEDICIONS / f"postedicions_{avui}.jsonl"
    try:
        with open(fitxer_dia, "a", encoding="utf-8") as f:
            f.write(json.dumps(entrada, ensure_ascii=False) + "\n")
    except OSError as exc:
        log.error("No s'ha pogut desar la postedició: %s", exc)
        raise HTTPException(
            status_code=500,
            detail=f"Error desant la postedició al disc: {exc}",
        )

    _stats["postedicions"] += 1
    log.info(
        "Postedició desada — id=%s tecnic='%s' fitxer=%s",
        id_entrada, peticio.tecnic, fitxer_dia.name,
    )

    return RespostaPostedicio(estat="desat", id=id_entrada)


# ─── Endpoint: GET /estadistiques ────────────────────────────────────────────

@app.get(
    "/estadistiques",
    response_model = RespostaEstadistiques,
    summary        = "Estadístiques d'ús de la sessió actual",
    tags           = ["Sistema"],
)
async def estadistiques() -> RespostaEstadistiques:
    """Retorna el recompte de paraules traduïdes, fitxers processats i postedicions del dia."""
    _reinicia_stats_si_cal()
    return RespostaEstadistiques(
        paraules_avui = _stats["paraules_avui"],
        fitxers_avui  = _stats["fitxers_avui"],
        postedicions  = _stats["postedicions"],
    )


# ─── Esquemes Pydantic — Glossari ────────────────────────────────────────────

class PeticioEntradaGlossari(BaseModel):
    es:     str
    ca:     str
    tecnic: str = ''
    domini: str


class RespostaGlossari(BaseModel):
    domini:   str
    entrades: list[dict]
    total:    int


# ─── Endpoints: Glossaris ─────────────────────────────────────────────────────

@app.get("/glossaris", tags=["Glossari"],
         summary="Llista els dominis disponibles")
async def llista_dominis():
    """Retorna la llista de dominis temàtics dels glossaris."""
    return {"dominis": DOMINIS}


# ─── Desar termes massius al glossari (IMPORTANT: ha d'anar ABANS de /glossari/{domini}) ──

class PeticioDesaGlossariMassiu(BaseModel):
    domini: str
    termes: List[Dict[str, Any]]  # [{"es": "...", "va": "..."}]


@app.post("/glossari/desa-massiu", tags=["Glossari"],
          summary="Afig una llista de termes al glossari d'un domini")
async def desa_glossari_massiu(peticio: PeticioDesaGlossariMassiu):
    """Afig una llista de termes al glossari d'un domini evitant duplicats."""
    log.info("DESA GLOSSARI: domini='%s', num_termes=%d, primer_terme=%s",
             peticio.domini, len(peticio.termes),
             peticio.termes[0] if peticio.termes else 'cap')

    if peticio.domini not in DOMINIS:
        raise HTTPException(status_code=404,
                            detail=f"Domini no trobat: {peticio.domini}")

    path = nom_fitxer_glossari(peticio.domini)
    entrades = carrega_glossari(peticio.domini) if path.exists() else []

    # Conjunt de termes existents per evitar duplicats (camp 'es' + 'ca')
    existents = {(e.get("es", "").lower(), e.get("ca", "").lower()) for e in entrades}

    nous = 0
    for terme in peticio.termes:
        es = terme.get("es", "").strip()
        # Accepta tant 'va' com 'ca' del frontend
        va = terme.get("va", terme.get("ca", "")).strip()
        if es and va and (es.lower(), va.lower()) not in existents:
            entrades.append({
                "es":     es,
                "ca":     va,           # El camp del TSV és 'ca', no 'va'
                "tecnic": "extracció automàtica",
                "data":   date.today().isoformat(),
                "domini": peticio.domini,
            })
            existents.add((es.lower(), va.lower()))
            nous += 1

    # Desa el glossari actualitzat
    desa_glossari(peticio.domini, entrades)

    log.info("Glossari '%s': %d termes nous afegits (%d totals)",
             peticio.domini, nous, len(entrades))
    return {"ok": True, "nous": nous, "total": len(entrades)}


@app.get("/glossari/{domini}", response_model=RespostaGlossari,
         tags=["Glossari"], summary="Obté les entrades d'un glossari")
async def obte_glossari(domini: str):
    """Retorna totes les entrades del glossari d'un domini."""
    if domini not in DOMINIS:
        raise HTTPException(status_code=404,
                            detail=f"Domini no trobat: {domini}")
    entrades = carrega_glossari(domini)
    return RespostaGlossari(domini=domini, entrades=entrades, total=len(entrades))


@app.post("/glossari/{domini}", tags=["Glossari"],
          summary="Afegeix o actualitza una entrada del glossari")
async def afegeix_entrada(domini: str, peticio: PeticioEntradaGlossari):
    """
    Afegeix una nova entrada o actualitza una existent (per terme en castellà).
    Si el terme ja existeix, n'actualitza la traducció i el tècnic.
    """
    if domini not in DOMINIS:
        raise HTTPException(status_code=404,
                            detail=f"Domini no trobat: {domini}")
    entrades = carrega_glossari(domini)
    for e in entrades:
        if e['es'].lower() == peticio.es.lower():
            e['ca']     = peticio.ca
            e['tecnic'] = peticio.tecnic
            e['data']   = datetime.now().strftime('%Y-%m-%d')
            desa_glossari(domini, entrades)
            log.info("Glossari '%s' — terme actualitzat: '%s'", domini, peticio.es)
            return {"estat": "actualitzat", "entrada": e}
    nova = {
        'es':     peticio.es.strip(),
        'ca':     peticio.ca.strip(),
        'tecnic': peticio.tecnic.strip(),
        'data':   datetime.now().strftime('%Y-%m-%d'),
        'domini': domini,
    }
    entrades.append(nova)
    desa_glossari(domini, entrades)
    log.info("Glossari '%s' — terme afegit: '%s'", domini, peticio.es)
    return {"estat": "afegit", "entrada": nova}


@app.delete("/glossari/{domini}/{terme_es}", tags=["Glossari"],
            summary="Elimina una entrada del glossari")
async def elimina_entrada(domini: str, terme_es: str):
    """Elimina una entrada del glossari pel terme en castellà."""
    if domini not in DOMINIS:
        raise HTTPException(status_code=404,
                            detail=f"Domini no trobat: {domini}")
    entrades       = carrega_glossari(domini)
    entrades_noves = [e for e in entrades if e['es'].lower() != terme_es.lower()]
    if len(entrades_noves) == len(entrades):
        raise HTTPException(status_code=404,
                            detail=f"Terme no trobat: {terme_es}")
    desa_glossari(domini, entrades_noves)
    log.info("Glossari '%s' — terme eliminat: '%s'", domini, terme_es)
    return {"estat": "eliminat", "terme": terme_es}


@app.get("/glossari/{domini}/exporta", tags=["Glossari"],
         summary="Exporta el glossari d'un domini com a fitxer TSV")
async def exporta_glossari(domini: str):
    """Exporta el glossari d'un domini com a fitxer TSV descarregable."""
    from fastapi.responses import FileResponse

    if domini not in DOMINIS:
        raise HTTPException(status_code=404, detail=f"Domini no trobat: {domini}")

    path = nom_fitxer_glossari(domini)
    if not path.exists() or path.stat().st_size == 0:
        raise HTTPException(status_code=404,
                            detail="El glossari és buit, no hi ha res a exportar.")

    # Nom del fitxer: glossari_NomDomini_ddmmaaaa.tsv
    data_avui      = datetime.now().strftime('%d%m%Y')
    nom_domini_net = re.sub(r'[^\w]', '_', domini).strip('_')
    nom_fitxer     = f"glossari_{nom_domini_net}_{data_avui}.tsv"

    log.info("GET /glossari/%s/exporta → %s", domini, nom_fitxer)
    return FileResponse(
        path        = str(path),
        media_type  = 'text/tab-separated-values; charset=utf-8',
        filename    = nom_fitxer,
        headers     = {"Content-Disposition": f'attachment; filename="{nom_fitxer}"'},
    )


# ─── Extracció automàtica de glossari bilingüe amb Claude ─────────────────────

class PeticioExtreuGlossari(BaseModel):
    text_original:  str = Field(..., min_length=10, max_length=500_000,
                                description="Text original en castellà")
    text_traduccio: str = Field(..., min_length=10, max_length=500_000,
                                description="Traducció al valencià revisada")
    domini:         str = Field(..., min_length=1,
                                description="Domini lingüístic d'especialitat")


@app.post(
    "/extreu-glossari",
    summary="Extreu glossari bilingüe es↔va d'un parell de textos",
    tags=["Glossari"],
)
async def extreu_glossari(peticio: PeticioExtreuGlossari):
    """
    Analitza un text original en castellà i la seua traducció al valencià
    per a extraure'n el vocabulari especialitzat bilingüe.
    Usa Claude Sonnet per a la detecció de termes.
    """
    api_key = _obte_api_key_anthropic()
    if not api_key:
        raise HTTPException(status_code=503, detail="Clau API d'Anthropic no configurada.")

    prompt_extraccio = (
        "Ets un terminòleg expert en extracció de vocabulari especialitzat bilingüe "
        "castellà-valencià/català.\n\n"
        "Analitza els dos textos següents: un ORIGINAL en castellà i la seua TRADUCCIÓ "
        "al valencià. Extrau TOTS els termes especialitzats, tècnics o rellevants per "
        "al domini indicat.\n\n"
        f"Domini lingüístic: {peticio.domini}\n\n"
        "INSTRUCCIONS:\n"
        "1. Identifica substantius, adjectius, verbs i expressions tècniques del domini.\n"
        "2. NO incloure paraules comunes o genèriques (articles, preposicions, conjuncions, pronoms).\n"
        "3. NO incloure noms propis de persones.\n"
        "4. SÍ incloure: termes tècnics, expressions fixades, locucions, "
        "unitats terminològiques complexes (2-4 paraules).\n"
        "5. Per a cada terme, proporciona la parella castellà → valencià.\n"
        "6. Ordena els termes alfabèticament pel castellà.\n"
        "7. Si un terme castellà té múltiples traduccions, tria la que apareix al text traduït.\n\n"
        f"ORIGINAL (castellà):\n{peticio.text_original[:50000]}\n\n"
        f"TRADUCCIÓ (valencià):\n{peticio.text_traduccio[:50000]}\n\n"
        "Respon EXCLUSIVAMENT amb un JSON vàlid (sense blocs markdown, sense ```json). "
        "Format exacte:\n\n"
        "{\n"
        f'  "domini": "{peticio.domini}",\n'
        '  "termes": [\n'
        '    {"es": "terme en castellà", "va": "terme en valencià"},\n'
        '    {"es": "altre terme", "va": "altra traducció"}\n'
        "  ],\n"
        '  "num_termes": 0\n'
        "}"
    )

    try:
        system_blocks = [
            {
                "type": "text",
                "text": (
                    "Ets un terminòleg expert en extracció de vocabulari bilingüe "
                    "castellà-valencià.\n"
                    "REGLA ABSOLUTA: la teua resposta ha de contindre EXCLUSIVAMENT "
                    "un objecte JSON vàlid. Cap text abans, cap text després, cap "
                    "bloc markdown (```), cap comentari. NOMÉS JSON pur."
                ),
                "cache_control": {"type": "ephemeral"},
            }
        ]

        resposta = await _crida_claude_amb_cache(
            system_blocks=system_blocks,
            missatge_usuari=prompt_extraccio,
            api_key=api_key,
            max_tokens=4096,
        )

        # ── Parseig robust de la resposta ──────────────────────────────────
        resposta_crua = resposta  # guardem l'original per a diagnòstic
        log.debug("Resposta crua de Claude (extreu-glossari): %.500s", resposta_crua)

        resposta = resposta.strip()

        # 1. Elimina blocs markdown ```json ... ``` o ``` ... ```
        resposta = re.sub(r'^```\w*\s*\n?', '', resposta)
        resposta = re.sub(r'\n?```\s*$', '', resposta)
        resposta = resposta.strip()

        # 2. Intenta parseig directe
        resultat = None
        try:
            resultat = json.loads(resposta)
        except json.JSONDecodeError:
            pass

        # 3. Fallback: busca objecte JSON {...} dins del text
        if resultat is None:
            m = re.search(r'\{[\s\S]*\}', resposta)
            if m:
                try:
                    resultat = json.loads(m.group())
                except json.JSONDecodeError:
                    pass

        # 4. Fallback: busca array JSON [...] i empaqueta-ho
        if resultat is None:
            m = re.search(r'\[[\s\S]*\]', resposta)
            if m:
                try:
                    termes_array = json.loads(m.group())
                    if isinstance(termes_array, list):
                        resultat = {"domini": peticio.domini, "termes": termes_array}
                except json.JSONDecodeError:
                    pass

        # 5. Si tot falla, error amb diagnòstic
        if resultat is None:
            log.error(
                "No s'ha pogut extraure JSON vàlid de la resposta de Claude. "
                "Primers 500 chars: %.500s", resposta_crua,
            )
            raise HTTPException(
                status_code=500,
                detail="El model no ha retornat un JSON vàlid. Torneu a intentar-ho.",
            )

        # ── Normalització de camps ─────────────────────────────────────────
        if "termes" in resultat and isinstance(resultat["termes"], list):
            termes_normalitzats = []
            for t in resultat["termes"]:
                if not isinstance(t, dict):
                    continue
                # Acceptem múltiples noms de camp per al castellà
                es = (t.get("es") or t.get("castellà") or t.get("castellano")
                      or t.get("español") or "")
                # Acceptem múltiples noms de camp per al valencià
                va = (t.get("va") or t.get("ca") or t.get("valencià")
                      or t.get("valenciano") or t.get("catalán")
                      or t.get("català") or "")
                if es and va:
                    termes_normalitzats.append({"es": es.strip(), "va": va.strip()})
            resultat["termes"] = termes_normalitzats

        resultat["domini"] = peticio.domini
        resultat["num_termes"] = len(resultat.get("termes", []))

        log.info(
            "POST /extreu-glossari — domini='%s', %d termes extrets",
            peticio.domini, resultat["num_termes"],
        )

        return resultat

    except HTTPException:
        raise
    except Exception as exc:
        log.exception("Error extraient glossari: %s", exc)
        raise HTTPException(status_code=500,
                            detail=f"Error en l'extracció del glossari: {exc}")



# ─── Traducció d'imatges amb Gemini (Nano Banana Pro) ────────────────────────

PROMPT_TRADUCCIO_IMATGE_DEFAULT = (
    "Usa Nano Banana Pro i edita aquesta imatge per a traduir-ne tot el text del castellà "
    "(espanyol), anglès i/o qualsevol altra llengua a la varietat valenciana universitària "
    "del català. Has de retornar exactament la mateixa imatge (mateix format, grandària, "
    "disposició, distribució i separació del text, grafisme, tipografia, fons, disseny, "
    "icones, estructura, colors, tipus i grandària de les fonts, espaiat, imatges, etc.) "
    "però amb tot el text traduït al català valencià i amb el menor pes de fitxer possible, "
    "sempre preservant-ne una resolució òptima i la màxima llegibilitat del text. "
    "Pel que fa al model de llengua a usar, recorda que has de traduir al valencià formal "
    "universitari, d'acord amb els criteris lingüístics per als usos institucionals de les "
    "universitats valencianes, els quals s'adjunten: demostratius reforçats (aquest, aquesta, "
    "aquests i aquestes), terminacions de verbs incoatius en -eix (divideix, parteix, segueix, "
    "etc.), el futur simple NO es pot usar per a expressar obligatorietat, sinó que per a "
    "aquests casos s'usa el present simple (per exemple, \"el tribunal està format pel "
    "president i cinc vocals\", i no \"el tribunal *estarà format pel president i cinc "
    "vocals\") o la perífrasi d'obligació \"haver de + infinitiu\" (per exemple, \"els equips "
    "han d'estar formats per cinc membres\", i no \"els equips *estaran formats per cinc "
    "membres\"), i tota la resta de criteris normatius que conformen els esmentats criteris "
    "lingüístics. En casos de doblets lèxics o geosinònims (com ara \"bresquilla/préssec\", "
    "\"xic/noi\", \"tomaca/tomàquet\", \"eixir/sortir\", \"espill/mirall\" o \"redó/rodó\", "
    "entre molts d'altres), has de triar sempre la primera opció, l'opció valenciana. "
    "Particularment important: el participi preferent del verb \"ser\" és \"sigut\" (mai "
    "\"estat\"); per a expressar finalitat, la locució preferent és \"per a + infinitiu\" "
    "(no \"per + infinitiu\"); la preposició \"de\" s'apostrofa sempre davant de vocal o h "
    "muda (d'art, d'història, d'eines); cal evitar l'ús abusiu del possessiu (seu/seua) "
    "substituint-lo pel pronom feble \"en\" quan el posseïdor ja és clar pel context; els "
    "plurals dels mots acabats en -sc, -st, -xt i -ig es fan preferentment en -os (textos, "
    "gustos, discos, rojos); els infinitius preferents són tenir, venir, cabre; les formes "
    "de l'imperfet de subjuntiu preferents són les que contenen -ra (cantara, perdera, "
    "servira); els possessius preferents són els que contenen -u- (meua, teua, seua); i els "
    "numerals preferents són vuit, disset, divuit, dinou, cinquè, sisè, desè. En cas de "
    "dubte, aplica sempre les formes recomanades pels Criteris lingüístics per als usos "
    "institucionals de les universitats valencianes. Sempre que pugues, prefereix servei (no "
    "servici), ordre (no orde), vacances (no vacacions), veure (no vore), desenvolupar (no "
    "desenrotllar), eina (no ferramenta), mentre (no mentres), endemà (no sendemà), meitat "
    "(no mitat), avui (no hui), aprendre (no dependre), judici (no juí), defensar (no "
    "defendre) i petit (no xicotet). "
    "Quan trobes terminologia especialitzada de qualsevol àmbit (biologia, medicina, física, "
    "química, dret, informàtica, enginyeria, etc.), NO inventes ni al·lucines la traducció. "
    "Busca sempre el terme correcte en català/valencià al Termcat "
    "(https://www.termcat.cat/ca/cercaterm), que és el centre de terminologia de la llengua "
    "catalana i l'autoritat en neologia i terminologia especialitzada. Exemples: "
    "\"cell membrane\" → \"membrana cel·lular\" (no \"membrana de la cèl·lula\"); "
    "\"machine learning\" → \"aprenentatge automàtic\" (no \"aprenentatge de màquina\"); "
    "\"hard drive\" → \"disc dur\" (no \"unitat de disc dur\"); "
    "\"cloud computing\" → \"informàtica en el núvol\" (no \"computació al núvol\"). "
    "Si no estàs segur de la traducció d'un terme especialitzat, usa la forma més establida "
    "en l'àmbit acadèmic i científic català/valencià, sense inventar neologismes."
)


class PeticioTraduccioImatge(BaseModel):
    imatge_base64:     str
    tipus_mime:        str = "image/png"
    prompt_addicional: str = ""
    mode:              str = "traduccio"  # 'traduccio' o 'refinament'


@app.post("/tradueix-imatge", tags=["Traducció"],
          summary="Tradueix el text d'una imatge del castellà/anglès al valencià (Gemini)")
async def tradueix_imatge(peticio: PeticioTraduccioImatge):
    """
    Tradueix el text inserit en una imatge del castellà/anglès al valencià
    usant l'API de Gemini (Nano Banana Pro / gemini-3-pro-image-preview).
    Usa el nou paquet google-genai (>=1.0.0).
    """
    import os
    from google import genai
    from google.genai import types

    api_key = _obte_api_key_gemini()
    if not api_key:
        raise HTTPException(
            status_code=503,
            detail="GEMINI_API_KEY no configurada. Introdueix-la al panell de configuració.",
        )

    try:
        client = genai.Client(api_key=api_key)
        imatge_bytes = base64.b64decode(peticio.imatge_base64)

        # Construeix el prompt posant les instruccions addicionals AL PRINCIPI
        # si n'hi ha, perquè Gemini les prioritze sobre el prompt general
        if peticio.prompt_addicional.strip():
            if peticio.mode == 'refinament':
                prompt_final = (
                    f"INSTRUCCIONS PRIORITÀRIES DEL TÈCNIC LINGÜÍSTIC"
                    f" (cal aplicar-les OBLIGATÒRIAMENT i EXACTAMENT):\n"
                    f"{peticio.prompt_addicional.strip()}\n\n"
                    f"Aplica aquestes instruccions sobre el text de la imatge."
                    f" No canvies cap element visual (colors, fonts, disseny, layout,"
                    f" fons, icones, estructura) excepte el text especificat."
                    f" Retorna exactament la mateixa imatge amb únicament els canvis"
                    f" de text sol·licitats."
                )
            else:
                prompt_final = (
                    f"INSTRUCCIONS PRIORITÀRIES DEL TÈCNIC LINGÜÍSTIC"
                    f" (cal aplicar-les OBLIGATÒRIAMENT i amb MÀXIMA PRIORITAT,"
                    f" per damunt de qualsevol altra consideració):\n"
                    f"{peticio.prompt_addicional.strip()}\n\n"
                    f"A més de les instruccions anteriors, aplica també les normes"
                    f" generals següents:\n{PROMPT_TRADUCCIO_IMATGE_DEFAULT}"
                )
        else:
            if peticio.mode == 'refinament':
                prompt_final = (
                    "Edita aquesta imatge aplicant les modificacions indicades sobre el text. "
                    "No canvies cap element visual (colors, fonts, disseny, layout, fons,"
                    " icones, estructura) excepte el text especificat. Retorna exactament"
                    " la mateixa imatge amb únicament els canvis de text sol·licitats."
                )
            else:
                prompt_final = PROMPT_TRADUCCIO_IMATGE_DEFAULT

        resposta = client.models.generate_content(
            model="gemini-3-pro-image-preview",
            contents=[
                types.Part.from_bytes(
                    data=imatge_bytes,
                    mime_type=peticio.tipus_mime,
                ),
                prompt_final,
            ],
            config=types.GenerateContentConfig(
                response_modalities=["IMAGE", "TEXT"],
            ),
        )

        for part in resposta.candidates[0].content.parts:
            if part.inline_data is not None:
                imatge_traduit_b64 = base64.b64encode(
                    part.inline_data.data
                ).decode('utf-8')
                log.info("POST /tradueix-imatge — OK tipus=%s", part.inline_data.mime_type)
                return {
                    "imatge_base64": imatge_traduit_b64,
                    "tipus_mime":    part.inline_data.mime_type,
                    "estat":         "ok",
                }

        raise HTTPException(
            status_code=500,
            detail="Gemini no ha retornat cap imatge. Prova amb un prompt diferent.",
        )

    except HTTPException:
        raise
    except Exception as e:
        log.exception("Error en la traducció d'imatge: %s", e)
        raise HTTPException(
            status_code=500,
            detail=f"Error en la traducció d'imatge: {e}",
        )


@app.post("/configura-gemini", tags=["Configuració"],
          summary="[DEPRECAT] Usa /api-keys/desa en comptes d'aquest endpoint")
async def configura_gemini(api_key: str = Body(..., embed=True)):
    """Mantingut per compatibilitat. Usa POST /api-keys/desa."""
    import os
    if not api_key.startswith("AIza"):
        raise HTTPException(status_code=400,
                            detail="Clau API de Gemini no vàlida (ha de començar per 'AIza').")
    api_key_neta = api_key.strip().strip("'\"")
    os.environ["GEMINI_API_KEY"] = api_key_neta
    if _DOTENV_OK:
        _dotenv_set_key(str(_ENV_PATH), "GEMINI_API_KEY", api_key_neta, quote_mode="never")
    log.info("Clau API de Gemini configurada (via endpoint antic).")
    return {"estat": "ok", "missatge": "Clau API configurada correctament."}


# ─── Correcció/postedició de documents en valencià ────────────────────────────

# Clau API d'Anthropic (es pot establir per entorn o via /configura-anthropic)
ANTHROPIC_API_KEY_CORRECCIO: str = ""

PROMPT_CORRECCIO_SISTEMA = """Ets un corrector ortogràfic, gramatical i estilístic especialitzat en **valencià normatiu universitari**.
La teua funció és revisar textos en català/valencià i retornar-los completament corregits, aplicant estrictament les normes que es detallen a continuació.

════════════════════════════════════════════════════════════════
BLOC A — CRITERIS LINGÜÍSTICS DE LES UNIVERSITATS VALENCIANES
════════════════════════════════════════════════════════════════

A1. DEMOSTRATIUS. Usa exclusivament el sistema binari reforçat en registre formal:
    · Proximitat: aquest/aquesta/aquests/aquestes (NO este/esta/estos/estes)
    · Llunyania: aquell/aquella/aquells/aquelles
    · Elimina el sistema ternari: este/eixe/aquell → aquest/aquell

A2. VERBS INCOATIUS. Forma culta amb -esc/-eix (no -isc/-ix):
    · servesc/serveixes/serveix/servim/serviu/serveixen
    · patesc/pateixes/pateix/patim/patiu/pateixen
    · oferesc/ofereixes/ofereix/oferim/oferiu/ofereixen
    · establesc, aparesc, meresc, conec, aparesc, pertanys

A3. POSSESSIUS. Forma amb -u- (tònica):
    · meua/teua/seua (NO meva/teva/seva)
    · meus/teus/seus, meues/teues/seues (NO meves/teves/seves)

A4. ACCENTUACIÓ. Segueix la norma del català central excepte:
    · anglès/anglesa (NO anglés/anglesa)
    · francès/francesa (NO francés/francesa)
    · cortès/cortesa (NO cortés/cortesa)
    · Al remat de paraula: café → cafè, però mantén: perquè, però

A5. LÈXIC. Prefereix les formes cultes i generals:
    · avui (NO hui), però → però/ara bé
    · menut/menuda (NO xicotet/xicoteta en registre formal)
    · aprendre (NO dependre en el sentit cognitiu)
    · vespre (NO vesprada com a sinònim de tarda en formal)

A6. PLURALS DELS MOTS EN -C/-G/-X:
    · discos, textos, èxits (NO discs, texts, èxits)
    · Però: aspectes, projectes (correctes)

A7. PARTICIPIS. Formes regulars en -it (no irregulars afavorides en col·loquial):
    · complit (NO complert), oferit (NO ofert)
    · establit (NO establert), omplit (NO omplert)
    · Excepcions lexicalitzades: obert, escrit, vist, dit, fet, mort (formes fortes acceptades)

A8. CONNECTORS DISCURSIUS. Usa connectors propis del registre escrit formal:
    · Per tant, per consegüent, de manera que, tanmateix, malgrat això
    · Evita: llavors (com a connector causal), aleshores (excepte en sentit temporal)

A9. TRACTAMENT PERSONAL. En comunicació institucional:
    · Vós/vosaltres per al tractament de cortesia col·lectiu
    · Vostè/vostès en comunicació molt formal i protocol·lària

A10. FORMES VERBALS PERIFRÀSTIQUES. Pretèrit perfet perifràstic:
     · va fer, va dir, van anar (formes acceptades)
     · NB: el perfet sintètic (feu, digué, anaren) és preferible en estil literari formal

════════════════════════════════════════════════════════════════
BLOC B — GRAMÀTICA: NORMES SINTÀCTIQUES I MORFOLÒGIQUES (60 REGLES)
════════════════════════════════════════════════════════════════

B1.  GÈNERE: masculí per defecte en noms epicens institucionals quan no hi ha referent concret.
B2.  NOMBRE: concordança estricta subjecte–verb fins i tot en construccions invertides.
B3.  ARTICLE DEFINIT: el/la/els/les davant noms propis geogràfics catalans (el País Valencià).
B4.  ARTICLE NEUTRE: ho (no el) com a pronom neutre (m'agrada ← ho faig bé).
B5.  APOSTROFACIÓ: l'home, l'hora, l'IVAM; però: la universitat (vocal feble no s'apostrofa si és u àtona).
B6.  PREPOSICIONS. De/del/de la: no contraure si és nom propi femení (de la Maria).
B7.  QUE/QUÈ: qué interrogatiu/exclamatiu → qué; relatiu sense pausa → que.
B8.  RELATIU QUI: reservat a persones (la persona qui ho fa); QUE per a coses.
B9.  ON/ON QUE: on (lloc); on que (incorrecte → on / en el qual).
B10. PER A / PER: finalitat → per a; causa/agent → per.
B11. EN / AMB: mitjà instrumental → amb; lloc dins → en.
B12. INFINITIU: no usar -r final en contacte amb pronom (fer-lo, dir-li; NO fè-lo).
B13. GERUNDI. Evita gerundis no concurrents o de posterioritat.
B14. PASIVA REFLEXA: es construeix amb se (es publica, se celebrarà) en registre formal.
B15. PRONOM SE: no confondre se reflexiu i se passiu; no usar en registre formal.
B16. PRONOMS FEBLES. Ordre: reflexiu > datiu > acusatiu > hi > en (se li'n dóna).
B17. CLÍTICS. No duplicar pronoms si l'objecte és tònic (ho fa ell, no *el fa ell).
B18. EN ENANTIOPOSICIÓ: en davant de numerals partitius (en tinc tres, no *tinc tres).
B19. HI/HI HA: hi ha (existencial); hi ha d'haver (obligació); no *hi ha de que.
B20. ARTICLE + ADJECTIU POSSESSIU: la meua feina (article + possessiu obligatori en català).
B21. DEMOSTRATIU + POSSESSIU: aquesta feina meua (ordre fix).
B22. NEGACIÓ: no + verb; ni … ni (correlativa); tampoc (no *tampoc no en registre no marcat).
B23. DOBLE NEGACIÓ: en registre formal, evitar «no … pas» (dialectalisme).
B24. SUBJECTIU VALOR MODAL: cal que + subjuntiu; és necessari que + subjuntiu.
B25. CONDICIONAL HIPOTÈTIC: si + imperfet de subjuntiu, condicional simple (si tinguera, faria).
B26. CONCORDANÇA TEMPORAL: mantén la seqüència temporal (narració al passat → imperfet/plusquamperfet).
B27. VOZ ACTIVA preferible a veu passiva perifràstica en registre administratiu.
B28. QUEISME: evitar *de que quan la subordinada és subjecte o objecte directe.
B29. DEQUEISME: usar de que quan el verb regeix preposició de (estic segur que → estic segur que ← correcte).
B30. COMPLEMENT PREDICATIU: concorda amb el subjecte (van arribar contents, no *content).
B31. ADJECTIU ATRIBUTIU: posició preferent postnominal en registre formal (informe detallat).
B32. ADVERBI DE MANERA: forma en -ment (clarament, ràpidament); evita perífrasis innecessàries.
B33. COMPARATIVES. Igualtat: tan … com; superioritat: més … que; inferioritat: menys … que.
B34. SUPERLATIUS. Absolut: molt + adj. o -íssim (utilíssim); relatiu: el més + adj.
B35. CONSTRUCCIÓ ABSOLUTA: havent acabat la reunió, … (correcte); *Acabant la reunió, … (evitar).
B36. ORACIONS DE RELATIU EXPLICATIVES: amb comes (la proposta, que fou aprovada, …).
B37. ORACIONS DE RELATIU ESPECIFICATIVES: sense comes (la proposta que fou aprovada …).
B38. MAJÚSCULES: noms propis, institucions, càrrecs en protocol, però no adjectius de gentilici.
B39. NUMERALS. Escriu en lletra els nombres de l'u al nou en prosa; xifres des del 10.
B40. FRACCIONS: la meitat, un terç, tres quartes parts (no *tres cuartos en calc del castellà).
B41. PERCENTATGE: el 35 % (amb espai i sense punt entre xifra i símbol).
B42. DATA: dia mes any (15 de març de 2026); no usar el format anglès ni abreviatures de mesos.
B43. HORA: les 9.30 h (punts, no dos punts); les 12 del migdia, les 0 hores.
B44. DIVISES: 1.500 € (punt per a milers, coma per a decimals: 1.500,75 €).
B45. ABREVIATURES. Usos normatius: pàg./pàgs., núm./núms., Sr./Sra., Dra./Dr.
B46. SIGLES: sense punts (UV, UPV, IEC); article concorda amb el substantiu elidit (l'IEC, la UV).
B47. TOPÒNIMS: forma oficial catalana (el País Valencià, Alacant, Castelló de la Plana).
B48. ANTROPÒNIMS: respecta la forma oficial de cada persona; en documentació usa nom complet en primera menció.
B49. ESTRANGERISMES: en cursiva si no adaptats; adapta els que tenen forma catalana (internet → internet, sense cursiva).
B50. LLATINISMES: en cursiva si no lexicalitzats (in situ, ex aequo, curriculum).
B51. TECNICISMES: usa la terminologia normativa del camp (consulta el TERMCAT si cal).
B52. CALCS DEL CASTELLÀ: evita calcs sintàctics i lèxics (a nivell de → pel que fa a; en base a → basant-se en; de cara a → per a).
B53. REPETICIONS: evita repeticions lèxiques innecessàries; usa pronoms i sinònims contextuals.
B54. ECONOMIA LINGÜÍSTICA: no usar circumlocucions on hi ha un terme precís.
B55. COHERÈNCIA TERMINOLÒGICA: un sol terme per a cada concepte al llarg del document.
B56. PARAL·LELISME: manté estructura paral·lela en enumeracions i elements coordinats.
B57. PUNTUACIÓ: coma davant de connectors adversatius (però, tanmateix, sinó); punt i coma per a separar elements llargs d'una enumeració.
B58. COMETES: « » (angulars) en català; " " per a cites internes.
B59. GUIÓ/GUIONET: guió llarg (—) per a incisos; guionet (-) per a compostos i prefixos.
B60. PARÈNTESI/CLAUDÀTOR: parèntesi per a aclariments; claudàtor per a interpolacions en citació textual.

════════════════════════════════════════════════════════════════
BLOC C — MANUAL DE DOCUMENTS ADMINISTRATIUS (UV)
════════════════════════════════════════════════════════════════

C1. ENCAPÇALAMENTS. L'encapçalament del document inclou: emissor, destinatari, assumpte i data.
C2. SALUTACIÓ. En comunicació oficial: "Senyor/Senyora," o "Benvolgut/Benvolguda,".
    Evita fórmules col·loquials o estrangeres.
C3. COMIAT. Formes protocol·làries: "Atentament," "Amb respecte," "Cordialment,".
    Evita: "Quedant a la vostra disposició" (gal·licisme) → "Restant a la vostra disposició".
C4. VERB EN REGISTRE ADMINISTRATIU. Usa les formes plenes:
    · sol·licitar (NO demanar en context formal), manifestar (NO dir), comunicar (NO avisar).
C5. ESTRUCTURES FIXES DOCUMENTALS.
    · "Faig constar que…" (certificats)
    · "Expose / Sol·licite / …" (instàncies)
    · "…, i a tal efecte, RESOLC:" (resolucions)
C6. VOZ I PERSONA. Prefereix la primera persona del singular en documents personals;
    impersonal en documents institucionals.
C7. TRACTAMENT INSTITUCIONAL. La Universitat de València (no *l'Universitat);
    el Rectorat, la Junta de Govern, el Consell de Govern (majúscules en noms d'òrgans).
C8. CITACIÓ DE NORMATIVA. "D'acord amb l'article 5 de la Llei…" (amb article determinat davant «article»).

════════════════════════════════════════════════════════════════
BLOC D — NOVETATS IEC 2016-2023 (GIEC, OIEC, GEIEC, GBU)
════════════════════════════════════════════════════════════════

D1. ACCENT DIACRÍTIC REDUÏT (OIEC 2017):
    · Sistema diacrític REDUÏT a 15 mots: bé/be, déu/deu, és/es, mà/ma, més/mes, món/mon,
      pèl/pel, sé/se, sí/si, sòl/sol, són/son, té/te, ús/us, déus/deus, béns/bens, pèls/pels,
      sís/sis, sòls/sols — aplica'ls correctament
    · NO porten accent diacrític (elimina'ls si hi eren): bota, coc, dona (verb donar),
      feu (verb fer), fora, soc (verb ser), sec, seu (greix)

D2. ERRADICAR S'ESCRIU AMB ERR- (OIEC 2017):
    · eradicar → erradicar, eradicació → erradicació

D3. FORMES NO ACCEPTABLES RECENTS (IEC GBU):
    · "don Pere" / "dona Montserrat" → "el senyor Pere" / "la senyora Montserrat"
    · Verbs psicològics: "li va afectar" → "la va afectar"; "em donen por" → "em fan por"
    · Tenir de + infinitiu: "tenim de fer" → "hem de fer"
    · "pròpia/propi" en lloc de "mateixa/mateix": "la pròpia directora" → "la mateixa directora"

D4. ORACIONS COPULATIVES — SER vs ESTAR (recordatori):
    · En registres formals NO usar restar per a durada d'estat: "restarà tancat" → "estarà tancat"

D5. ORDRE DEL SINTAGMA NOMINAL (IEC GIEC) — ESPECIALMENT IMPORTANT:
    · Ordre habitual en català: SUBSTANTIU + ADJECTIU ("una reunió important", "un informe detallat")
    · Detecta i corregeix anteposicions calcades del castellà:
      "el proper curs" → "el curs vinent", "la present comunicació" → "la comunicació present",
      "la corresponent documentació" → "la documentació corresponent",
      "l'esmentada resolució" → "la resolució esmentada"

════════════════════════════════════════════════════════════════
INSTRUCCIONS DE RESPOSTA
════════════════════════════════════════════════════════════════

Quan rebis un text per corregir:
1. Retorna el text completament corregit aplicant totes les regles anteriors.
2. Proporciona una llista detallada de les correccions realitzades en format JSON dins de ```json ``` amb aquest esquema:
   [
     {
       "original": "text original amb l'error",
       "corregit": "text corregit",
       "tipus": "ortografia|morfologia|sintaxi|lèxic|estil|puntuació|registre",
       "regla": "codi de la regla (p. ex. A1, B52, C4)",
       "justificacio": "explicació breu de la correcció"
     },
     ...
   ]
3. Proporciona un resum breu (1-2 frases) de les principals àrees de millora detectades.

Format de resposta OBLIGATORI:
---TEXT CORREGIT---
[text complet corregit]
---FI TEXT---
---CORRECCIONS---
```json
[llista de correccions]
```
---FI CORRECCIONS---
---RESUM---
[resum breu]
---FI RESUM---
"""


class PeticioCorreccio(BaseModel):
    text: str = Field(
        ...,
        min_length  = 1,
        max_length  = 100_000,
        description = "Text en valencià a corregir.",
    )
    usar_languagetool: bool = Field(
        default     = True,
        description = "Si cal aplicar primer la capa de LanguageTool (ca-ES).",
    )
    usar_claude: bool = Field(
        default     = True,
        description = "Si cal aplicar la capa de correcció amb Claude Sonnet.",
    )


class ItemCorreccioLT(BaseModel):
    missatge:  str
    offset:    int
    longitud:  int
    original:  str
    suggerits: list[str]
    regla_id:  str


class ItemCorreccioC(BaseModel):
    original:     str
    corregit:     str
    tipus:        str
    regla:        str
    justificacio: str


class RespostaCorreccio(BaseModel):
    text_original:     str
    text_corregit:     str
    correccions_lt:    list[dict]
    correccions_claude: list[dict]
    resum:             str
    estat:             str


@app.post("/configura-anthropic", tags=["Configuració"],
          summary="[DEPRECAT] Usa /api-keys/desa en comptes d'aquest endpoint")
async def configura_anthropic(api_key: str = Body(..., embed=True)):
    """Mantingut per compatibilitat. Usa POST /api-keys/desa."""
    import os
    global ANTHROPIC_API_KEY_CORRECCIO
    if not api_key.startswith("sk-ant-"):
        raise HTTPException(status_code=400,
                            detail="Clau API d'Anthropic no vàlida (ha de començar per 'sk-ant-').")
    api_key_neta = api_key.strip().strip("'\"")
    os.environ["ANTHROPIC_API_KEY"] = api_key_neta
    os.environ["ANTHROPIC_API_KEY_CORRECCIO"] = api_key_neta
    ANTHROPIC_API_KEY_CORRECCIO = api_key_neta
    if _DOTENV_OK:
        _dotenv_set_key(str(_ENV_PATH), "ANTHROPIC_API_KEY_CORRECCIO", api_key_neta, quote_mode="never")
    log.info("Clau API d'Anthropic configurada (via endpoint antic).")
    return {"estat": "ok", "missatge": "Clau API d'Anthropic configurada correctament."}


# ─── Gestió centralitzada de claus API ───────────────────────────────────────

class PeticioClauAPI(BaseModel):
    servei: str = Field(..., description="'gemini' o 'anthropic'")
    clau:   str = Field(..., description="Valor de la clau API")


class RespostaClauAPI(BaseModel):
    servei:       str
    configurada:  bool
    clau_parcial: str


@app.post(
    "/api-keys/desa",
    response_model = RespostaClauAPI,
    summary        = "Desa una clau API al .env de forma persistent",
    tags           = ["Configuració"],
)
async def desa_clau_api(peticio: PeticioClauAPI) -> RespostaClauAPI:
    """
    Desa una clau API al fitxer .env de l'arrel del projecte de forma
    persistent (s'aplicarà en cada reinici del servidor).
    Actualitza també la sessió actual sense necessitat de reiniciar.
    """
    import os
    global ANTHROPIC_API_KEY_CORRECCIO

    servei = peticio.servei.lower().strip()
    clau   = peticio.clau.strip()

    if servei == "gemini":
        if not clau.startswith("AIza"):
            raise HTTPException(
                status_code=400,
                detail="Clau de Gemini no vàlida (ha de començar per 'AIza').",
            )
        nom_variable = "GEMINI_API_KEY"
        os.environ[nom_variable] = clau

    elif servei == "anthropic":
        if not clau.startswith("sk-ant-"):
            raise HTTPException(
                status_code=400,
                detail="Clau d'Anthropic no vàlida (ha de començar per 'sk-ant-').",
            )
        nom_variable = "ANTHROPIC_API_KEY_CORRECCIO"
        os.environ[nom_variable] = clau
        os.environ["ANTHROPIC_API_KEY"] = clau
        ANTHROPIC_API_KEY_CORRECCIO = clau

    else:
        raise HTTPException(status_code=400, detail=f"Servei desconegut: '{servei}'.")

    # Persisteix al .env (quote_mode="never" evita que python-dotenv afegisca cometes)
    if _DOTENV_OK:
        clau_neta = clau.strip().strip("'\"")
        _dotenv_set_key(str(_ENV_PATH), nom_variable, clau_neta, quote_mode="never")
        log.info("Clau '%s' desada al .env i a l'entorn.", nom_variable)
    else:
        log.warning(
            "python-dotenv no disponible: la clau '%s' s'ha establit per a "
            "la sessió actual però NO s'ha desat al .env.",
            nom_variable,
        )

    clau_parcial = ("••••••••" + clau[-4:]) if len(clau) > 4 else "••••"
    return RespostaClauAPI(servei=servei, configurada=True, clau_parcial=clau_parcial)


@app.get(
    "/api-keys/estat",
    summary = "Estat de configuració de les claus API",
    tags    = ["Configuració"],
)
async def estat_claus_api():
    """
    Retorna si cada clau API està configurada i els darrers 4 caràcters
    emmascarat. Mai retorna el valor real de la clau.
    """
    import os

    def parcial(clau: str) -> str:
        if not clau:
            return ""
        return "••••••••" + clau[-4:]

    clau_gemini    = _obte_api_key_gemini()
    clau_anthropic = _obte_api_key_anthropic()

    return {
        "gemini": {
            "configurada":  bool(clau_gemini),
            "clau_parcial": parcial(clau_gemini),
        },
        "anthropic": {
            "configurada":  bool(clau_anthropic),
            "clau_parcial": parcial(clau_anthropic),
        },
    }


@app.post(
    "/corregeix",
    response_model = RespostaCorreccio,
    summary        = "Corregeix un text en valencià (LanguageTool + Claude Sonnet)",
    tags           = ["Correcció"],
)
async def corregeix(peticio: PeticioCorreccio) -> RespostaCorreccio:
    """
    Corregeix un text en valencià aplicant dues capes:
    1. LanguageTool (API pública, ca-ES): errors ortogràfics i gramaticals bàsics.
    2. Claude Sonnet (claude-sonnet-4-6): normes específiques del valencià universitari.

    Retorna el text corregit, la llista de correccions de cada capa i un resum.
    """
    import os
    import httpx
    import anthropic as _anthropic
    import json as _json

    text_entrant = peticio.text.strip()
    if not text_entrant:
        raise HTTPException(status_code=422, detail="El camp 'text' no pot estar buit.")

    correccions_lt: list[dict] = []
    text_despres_lt = text_entrant

    # ── CAPA 1: LanguageTool ──────────────────────────────────────────────────
    if peticio.usar_languagetool:
        try:
            async with httpx.AsyncClient(timeout=30.0) as client:
                resp = await client.post(
                    "https://api.languagetool.org/v2/check",
                    data={
                        "text":     text_entrant,
                        "language": "ca-ES",
                    },
                    headers={"Accept": "application/json"},
                )
                resp.raise_for_status()
                dades_lt = resp.json()

            for match in dades_lt.get("matches", []):
                context  = match.get("context", {})
                original = context.get("text", "")[
                    context.get("offset", 0):
                    context.get("offset", 0) + context.get("length", 0)
                ]
                suggerits = [r.get("value", "") for r in match.get("replacements", [])[:3]]
                correccions_lt.append({
                    "missatge":  match.get("message", ""),
                    "offset":    match.get("offset", 0),
                    "longitud":  match.get("length", 0),
                    "original":  original,
                    "suggerits": suggerits,
                    "regla_id":  match.get("rule", {}).get("id", ""),
                })

            # Aplica el primer suggerit de cada correcció (ordre invers per no desplaçar offsets)
            text_despres_lt = text_entrant
            for c in sorted(correccions_lt, key=lambda x: x["offset"], reverse=True):
                if c["suggerits"]:
                    ini = c["offset"]
                    fi  = ini + c["longitud"]
                    text_despres_lt = (
                        text_despres_lt[:ini]
                        + c["suggerits"][0]
                        + text_despres_lt[fi:]
                    )
            log.info("LanguageTool — %d coincidències trobades", len(correccions_lt))

        except httpx.HTTPError as exc:
            log.warning("Error en la petició a LanguageTool: %s", exc)
            # Continua sense LanguageTool si l'API no respon
        except Exception as exc:
            log.warning("Error inesperat a LanguageTool: %s", exc)

    # ── CAPA 2: Claude Sonnet ─────────────────────────────────────────────────
    correccions_claude: list[dict] = []
    text_final = text_despres_lt
    resum = ""

    if peticio.usar_claude:
        api_key_anthropic = _obte_api_key_anthropic()
        if not api_key_anthropic:
            raise HTTPException(
                status_code=503,
                detail=(
                    "ANTHROPIC_API_KEY no configurada. "
                    "Introdueix-la al panell de configuració."
                ),
            )

        try:
            client_a = _anthropic.Anthropic(api_key=api_key_anthropic)

            missatge_usuari = (
                f"Corregeix el text següent en valencià aplicant totes les normes indicades:\n\n"
                f"{text_despres_lt}"
            )

            resposta = client_a.messages.create(
                model      = "claude-sonnet-4-6",
                max_tokens = 8192,
                system     = PROMPT_CORRECCIO_SISTEMA,
                messages   = [{"role": "user", "content": missatge_usuari}],
            )

            contingut_resposta = resposta.content[0].text

            # Corregeix text UTF-8 mal interpretat com Latin-1 (Ã© → é, etc.)
            def _neteja_enc(t: str) -> str:
                try:
                    return t.encode("latin-1").decode("utf-8")
                except (UnicodeEncodeError, UnicodeDecodeError):
                    return t

            if "Ã" in contingut_resposta or "â€" in contingut_resposta:
                contingut_resposta = _neteja_enc(contingut_resposta)

            # Extreu el text corregit
            m_text = re.search(
                r'---TEXT CORREGIT---\s*(.*?)\s*---FI TEXT---',
                contingut_resposta,
                re.DOTALL,
            )
            if m_text:
                text_final = m_text.group(1).strip()

            # Extreu la llista de correccions JSON
            m_corr = re.search(
                r'---CORRECCIONS---.*?```json\s*(.*?)\s*```.*?---FI CORRECCIONS---',
                contingut_resposta,
                re.DOTALL,
            )
            if m_corr:
                try:
                    correccions_claude = _json.loads(m_corr.group(1).strip())
                except _json.JSONDecodeError:
                    log.warning("No s'ha pogut parsejar el JSON de correccions de Claude.")

            # Extreu el resum
            m_resum = re.search(
                r'---RESUM---\s*(.*?)\s*---FI RESUM---',
                contingut_resposta,
                re.DOTALL,
            )
            if m_resum:
                resum = m_resum.group(1).strip()

            log.info(
                "Claude Sonnet — %d correccions aplicades",
                len(correccions_claude),
            )

        except _anthropic.AuthenticationError:
            raise HTTPException(
                status_code=401,
                detail="Clau API d'Anthropic invàlida o caducada.",
            )
        except Exception as exc:
            log.exception("Error en la correcció amb Claude: %s", exc)
            raise HTTPException(
                status_code=500,
                detail=f"Error en la correcció amb Claude Sonnet: {exc}",
            )

    return RespostaCorreccio(
        text_original      = text_entrant,
        text_corregit      = text_final,
        correccions_lt     = correccions_lt,
        correccions_claude = correccions_claude,
        resum              = resum,
        estat              = "ok",
    )


# ─── Correcció de documents DOCX / PPTX ───────────────────────────────────────

# Namespaces XML Word / PowerPoint (locals a aquest mòdul)
_NS_W_DOC    = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
_NS_A_DOC    = "http://schemas.openxmlformats.org/drawingml/2006/main"
# Namespaces de fórmules i equacions matemàtiques
_NS_MATH_DOC = "http://schemas.openxmlformats.org/officeDocument/2006/math"
_NS_VML_DOC  = "urn:schemas-microsoft-com:vml"

# Fitxers XML que cal processar dins un DOCX
_FITXERS_XML_DOCX_FIXES = {
    "word/document.xml",
    "word/footnotes.xml",
    "word/endnotes.xml",
    "word/comments.xml",
}
_PATRÓ_FITXERS_DOCX_EXTRA = re.compile(
    r"word/(header\d+|footer\d+)\.xml$"
)
# Fitxers XML que cal processar dins un PPTX
_PATRÓ_FITXERS_PPTX = re.compile(
    r"ppt/(slides/slide\d+|notesSlides/notesSlide\d+)\.xml$"
)
# Atribut xml:space="preserve"
_XML_SPACE = "{http://www.w3.org/XML/1998/namespace}space"


def _obte_text_paràgraf_doc(p_node, ns_t: str) -> str:
    """
    Concatena el text de tots els nodes <w:t>/<a:t> fills del paràgraf.
    Respecta xml:space='preserve': si un node té l'atribut però text None,
    contribueix un espai (pot ocórrer en alguns documents Word/PPTX).
    Usat tant per a DOCX (ns_t = NS_W) com per a PPTX (ns_t = NS_A).
    """
    parts = []
    for t in p_node.iter(f"{{{ns_t}}}t"):
        if t.text is not None:
            parts.append(t.text)
        elif t.get(_XML_SPACE) == "preserve":
            # Node marcat com a preserve però sense text: representa un espai
            parts.append(" ")
    return "".join(parts)


def _substitueix_text_runs(p_node, text_nou: str, ns_t: str) -> None:
    """
    Distribueix text_nou entre els nodes <w:t>/<a:t> del paràgraf preservant
    tots els atributs de format (rPr, pPr, etc.).
    Tot el text va al primer node t; la resta queden buits.
    Afegeix xml:space="preserve" a qualsevol node que continga espais per
    evitar que Word/PowerPoint elimine silenciosament espais inicials, finals
    o interns (p. ex. espai entre paraules quan tot el text va al primer run).
    """
    # Inclou nodes amb text no None I nodes buits sense fills (poden rebre text)
    nodes_t = [
        n for n in p_node.iter(f"{{{ns_t}}}t")
        if n.text is not None or len(n) == 0
    ]
    if not nodes_t:
        # Si no hi havia cap node t, crea'n un dins del primer run
        tag_r   = f"{{{ns_t}}}r" if ns_t else f"{{{_NS_W_DOC}}}r"
        primers = list(p_node.iter(tag_r))
        if primers:
            from lxml import etree as _et
            nou = _et.SubElement(primers[0], f"{{{ns_t}}}t")
            nou.text = text_nou
            nou.set(_XML_SPACE, "preserve")
        return

    # Posa tot el text al primer node
    primer = nodes_t[0]
    primer.text = text_nou
    # xml:space="preserve" cal sempre que hi haja qualsevol espai al text:
    # espais inicials/finals (obligatori per l'estàndard OOXML) però també
    # espais interns (defensiu: evita que alguns parsers XML col·lapsen
    # whitespace al reserialitzar el document).
    if text_nou and (' ' in text_nou or '\t' in text_nou or '\n' in text_nou):
        primer.set(_XML_SPACE, "preserve")
    elif _XML_SPACE in primer.attrib:
        del primer.attrib[_XML_SPACE]

    # Buida els nodes restants i neteja l'atribut preserve (no en necessiten)
    for node in nodes_t[1:]:
        node.text = ""
        if _XML_SPACE in node.attrib:
            del node.attrib[_XML_SPACE]


# ── Marcadors de run ⟦N⟧ ─────────────────────────────────────────────────────
# Usem els caràcters Unicode U+27E6/U+27E7 (⟦⟧) perquè no apareixen
# en text normal i Claude els preserva sense interpretar-los.

def _text_pla(text: str) -> str:
    """Elimina marcadors ⟦N⟧ del text (resultat pla sense indicadors de run)."""
    import re as _re
    return _re.sub(r'\u27e6\d+\u27e7', '', text)


def _runs_son_dins_formula(run_el, paragraph_el) -> bool:
    """
    Comprova si un run és descendent d'un element de fórmula dins del paràgraf.
    Usat per saltar runs de dins de <m:oMath>, <m:oMathPara>, <w:object> i <w:pict>
    en l'extracció de text, de forma que les fórmules queden intactes.
    """
    ns_m = _NS_MATH_DOC
    ns_w = _NS_W_DOC
    formules = (
        list(paragraph_el.iter(f"{{{ns_m}}}oMath")) +
        list(paragraph_el.iter(f"{{{ns_m}}}oMathPara")) +
        list(paragraph_el.iter(f"{{{ns_w}}}object")) +
        list(paragraph_el.iter(f"{{{ns_w}}}pict"))
    )
    for formula_el in formules:
        for desc in formula_el.iter():
            if desc is run_el:
                return True
    return False


def _extrau_text_amb_marcadors_docx(paragraph_el) -> tuple[str, list]:
    """
    Extreu el text del paràgraf DOCX inserint marcadors ⟦N⟧ al principi de cada run.

    Retorna (text_marcat, runs_info):
      text_marcat: "⟦0⟧text_run0⟦1⟧text_run1..." per enviar a Claude
      runs_info  : [{"idx": N, "node": element_run, "text": text_run}]

    Els runs que formen part d'elements de fórmula (<m:oMath>, <w:object>,
    <w:pict>) s'ometen: el seu text no s'envia a Claude i l'estructura
    de l'equació es conserva intacta.
    """
    ns_w = _NS_W_DOC
    runs_info = []
    parts = []
    for idx, run in enumerate(paragraph_el.iter(f"{{{ns_w}}}r")):
        # Salta els runs dins d'elements de fórmula (OMML, OLE, VML)
        if _runs_son_dins_formula(run, paragraph_el):
            continue
        text_run = "".join(
            t.text if t.text is not None
            else (" " if t.get(_XML_SPACE) == "preserve" else "")
            for t in run.iter(f"{{{ns_w}}}t")
        )
        runs_info.append({"idx": idx, "node": run, "text": text_run})
        parts.append(f"\u27e6{idx}\u27e7{text_run}")
    return "".join(parts), runs_info


def _runs_son_dins_formula_pptx(run_el, txBody_el) -> bool:
    """
    Comprova si un run de PPTX és descendent d'un element de fórmula OMML.
    Rep el <a:txBody> pare per poder cercar <m:oMath> dins d'ell.
    """
    ns_m = _NS_MATH_DOC
    for formula_el in txBody_el.iter(f"{{{ns_m}}}oMath"):
        for desc in formula_el.iter():
            if desc is run_el:
                return True
    return False


def _extrau_text_amb_marcadors_pptx(paragraph_el) -> tuple[str, list]:
    """
    Equivalent de _extrau_text_amb_marcadors_docx per a paràgrafs PPTX (<a:p>).
    Usa l'espai de noms DrawingML (_NS_A_DOC).

    Els runs dins d'elements de fórmula OMML (<m:oMath>) s'ometen
    per preservar les equacions intactes.
    """
    ns_a = _NS_A_DOC
    runs_info = []
    parts = []
    # Obté el <a:txBody> pare per poder comprovar fórmules OMML
    txBody_el = paragraph_el.getparent()
    for idx, run in enumerate(paragraph_el.iter(f"{{{ns_a}}}r")):
        # Salta runs dins de fórmules OMML (si el txBody és accessible)
        if txBody_el is not None and _runs_son_dins_formula_pptx(run, txBody_el):
            continue
        text_run = "".join(
            t.text if t.text is not None
            else (" " if t.get(_XML_SPACE) == "preserve" else "")
            for t in run.iter(f"{{{ns_a}}}t")
        )
        runs_info.append({"idx": idx, "node": run, "text": text_run})
        parts.append(f"\u27e6{idx}\u27e7{text_run}")
    return "".join(parts), runs_info


def _aplica_text_run_simple(run_info: dict, text_nou: str, ns_t: str) -> None:
    """
    Posa text_nou al primer <w:t>/<a:t> del run i buida la resta.
    Posa SEMPRE xml:space="preserve" si hi ha text (no sols quan hi ha espais
    marginals): això evita que alguns parsers XML col·lapsen espais interns
    quan el run conté text que comença o acaba sense espai però ve seguit/precedit
    per espais procedents d'un run adjacent.
    """
    from lxml import etree as _et
    run = run_info["node"]
    nodes_t = list(run.iter(f"{{{ns_t}}}t"))
    if not nodes_t:
        nou_t = _et.SubElement(run, f"{{{ns_t}}}t")
        nou_t.text = text_nou
        if text_nou:
            nou_t.set(_XML_SPACE, "preserve")
        return
    primer = nodes_t[0]
    primer.text = text_nou
    if text_nou:
        # Sempre preserve: protegeix espais inicials, finals I interns
        primer.set(_XML_SPACE, "preserve")
    else:
        # Text buit: elimina preserve (evita nodes preserve sense contingut)
        if _XML_SPACE in primer.attrib:
            del primer.attrib[_XML_SPACE]
    for node in nodes_t[1:]:
        node.text = ""
        if _XML_SPACE in node.attrib:
            del node.attrib[_XML_SPACE]


def _aplica_text_marcat(runs_info: list, text_corregit: str, ns_t: str) -> bool:
    """
    Intenta aplicar text_corregit (que pot contenir marcadors ⟦N⟧) als runs.
    Vàlid per a DOCX (ns_t=_NS_W_DOC) i PPTX (ns_t=_NS_A_DOC).

    Si el text conté marcadors ⟦N⟧ coherents amb runs_info, distribueix cada
    fragment al run corresponent (preservant la tipografia run a run).

    Retorna True si s'han aplicat els marcadors; False si cal usar el mètode
    de fallback (_substitueix_text_runs / text al primer run).
    """
    import re as _re
    if not runs_info:
        return False

    # Cerca marcadors ⟦N⟧ al text corregit
    marcadors = list(_re.finditer(r'\u27e6(\d+)\u27e7', text_corregit))
    if not marcadors:
        return False  # Claude no ha preservat els marcadors → usa fallback

    # Comprova coherència mínima: els índexs han d'estar presents
    índexs_marcats = {int(m.group(1)) for m in marcadors}
    índexs_esperats = {r["idx"] for r in runs_info}
    # Accepta si almenys la meitat dels runs esperats apareix com a marcador
    coincidències = índexs_marcats & índexs_esperats
    if len(coincidències) < max(1, len(índexs_esperats) // 2):
        return False

    # Mapa índex → text del run (de la posició del marcador fins al següent o fi)
    marcadors_ord = sorted(marcadors, key=lambda m: int(m.group(1)))
    textos_per_run: dict[int, str] = {}
    for i, m in enumerate(marcadors_ord):
        fi = marcadors_ord[i + 1].start() if i + 1 < len(marcadors_ord) else len(text_corregit)
        textos_per_run[int(m.group(1))] = text_corregit[m.end():fi]

    # Aplica a cada run (si un run no té marcador, conserva el text original)
    for run_info in runs_info:
        idx = run_info["idx"]
        text_nou = textos_per_run.get(idx, run_info["text"])
        _aplica_text_run_simple(run_info, text_nou, ns_t)

    return True


def _reconstrueix_text_marcat(
    text_marcat_orig: str,
    text_pla_orig: str,
    text_pla_corr: str,
) -> str:
    """
    Reconstrueix el text marcat corregit distribuint text_pla_corr entre els
    segments del text_marcat_orig de forma proporcional.

    Quan Claude NO preserva els marcadors ⟦N⟧, aquesta funció permet intentar
    projectar la correcció sobre el text original marcat perquè la distribució
    de text entre runs siga raonable (en lloc de posar tot al primer run).

    La distribució segueix la proporció de caràcters del segment original.
    Els espais inter-segments s'inclouen al segment actual per evitar que
    desapareguen quan es buiden els nodes t dels runs posteriors.
    """
    import re as _re

    if not text_marcat_orig or not text_pla_corr:
        return text_pla_corr

    patró_marc = _re.compile(r'(\u27e6\d+\u27e7)')
    parts_orig = patró_marc.split(text_marcat_orig)

    textos_orig: list[str] = []
    marcadors:   list[str] = []
    for part in parts_orig:
        if patró_marc.match(part):
            marcadors.append(part)
        else:
            textos_orig.append(part)

    if len(textos_orig) <= 1:
        return text_pla_corr

    total_orig = sum(len(t) for t in textos_orig)
    if total_orig == 0:
        return text_pla_corr

    textos_corr: list[str] = []
    restant = text_pla_corr

    for idx, text_seg_orig in enumerate(textos_orig):
        if idx == len(textos_orig) - 1:
            textos_corr.append(restant)
            break

        if not text_seg_orig:
            textos_corr.append("")
            continue

        proporcio = len(text_seg_orig) / total_orig
        n = max(0, round(len(text_pla_corr) * proporcio))

        # Ajusta al límit de paraula i inclou l'espai separador en el segment
        # actual per garantir que no es perd quan el run adjacent es buida
        if 0 < n < len(restant):
            while n < len(restant) and restant[n] not in (' ', '\t', '\n'):
                n += 1
            if n < len(restant) and restant[n] == ' ':
                n += 1   # Inclou l'espai en el fragment actual

        fragment = restant[:n]
        restant  = restant[n:]
        textos_corr.append(fragment)

    # Reconstrueix intercalant marcadors
    resultat = ""
    for idx, text in enumerate(textos_corr):
        resultat += text
        if idx < len(marcadors):
            resultat += marcadors[idx]
    return resultat


def _paràgraf_conté_formula(paragraph_el) -> bool:
    """
    Detecta si un paràgraf DOCX conté fórmules o equacions matemàtiques:
    - OMML (equacions natives Word 2007+): <m:oMath>, <m:oMathPara>
    - Objectes OLE (Equation Editor antic): <w:object>
    - Imatges VML (equacions com a imatges): <w:pict>
    - Camps de fórmula Word: <w:instrText> amb operadors EQ/\\F/\\I/\\R
    Aquests paràgrafs s'exclouen de la correcció automàtica per evitar
    que la substitució de text destruïsca l'estructura XML de l'equació.
    """
    ns_w = _NS_W_DOC
    ns_m = _NS_MATH_DOC
    # OMML — equacions natives de Word 2007+
    if paragraph_el.find(f".//{{{ns_m}}}oMath") is not None:
        return True
    if paragraph_el.find(f".//{{{ns_m}}}oMathPara") is not None:
        return True
    # Objectes OLE incrustats (Equation Editor 3.x)
    if paragraph_el.find(f".//{{{ns_w}}}object") is not None:
        return True
    # Imatges VML (equacions rasteritzades)
    if paragraph_el.find(f".//{{{ns_w}}}pict") is not None:
        return True
    # Camps de fórmula Word (instrText amb operadors matemàtics)
    for instr in paragraph_el.iter(f"{{{ns_w}}}instrText"):
        if instr.text and any(
            s in instr.text.upper() for s in ('EQ ', '\\F(', '\\I(', '\\R(')
        ):
            return True
    return False


def _shape_conté_formula_pptx(txBody_el) -> bool:
    """
    Detecta si un cos de text PPTX (<a:txBody>) conté fórmules matemàtiques.
    - OMML incrustat dins DrawingML: <m:oMath>
    - graphicData amb URI de matemàtiques
    """
    ns_m = _NS_MATH_DOC
    ns_a = _NS_A_DOC
    if txBody_el.find(f".//{{{ns_m}}}oMath") is not None:
        return True
    for gd in txBody_el.iter(f"{{{ns_a}}}graphicData"):
        uri = gd.get("uri") or ""
        if "math" in uri.lower():
            return True
    return False


def _extrau_paràgrafs_docx(xml_bytes: bytes) -> tuple[list[dict], object]:
    """
    Analitza el XML de Word i retorna (paràgrafs, arbre):
      paràgrafs: [{index, text, text_marcat, node}]
      arbre    : l'arbre lxml parsejat (per reutilitzar-lo si cal)

    text_marcat conté el text amb marcadors ⟦N⟧ entre runs, que es pot enviar
    a Claude per a que corregisca preservant els límits de run.
    NOTA: els nodes emmagatzemats a "node" pertanyen a aquest arbre (passada 1);
    en passada 2 se'n crea un de nou amb _et.fromstring(xml_original).
    """
    from lxml import etree as _et
    arbre    = _et.fromstring(xml_bytes)
    resultat = []
    for i, p in enumerate(arbre.iter(f"{{{_NS_W_DOC}}}p")):
        text        = _obte_text_paràgraf_doc(p, _NS_W_DOC)
        text_marcat, _ = _extrau_text_amb_marcadors_docx(p)
        té_formula  = _paràgraf_conté_formula(p)
        resultat.append({
            "index": i, "text": text, "text_marcat": text_marcat,
            "té_formula": té_formula, "node": p,
        })
    return resultat, arbre


def _obte_paraules_canviades(text_original: str, text_corregit: str) -> set[str]:
    """
    Compara dos textos i retorna el conjunt de paraules canviades al text corregit.
    Usa difflib amb comparació token a token per màxima precisió.
    Retorna un conjunt buit si els textos són idèntics o molt similars.
    """
    import difflib
    import re as _re

    if not text_original or not text_corregit:
        return set()

    # Normalitza espais per a la comparació
    orig_net = text_original.strip()
    corr_net = text_corregit.strip()

    # Si els textos són idèntics, no hi ha res a destacar
    if orig_net == corr_net:
        return set()

    # Calcula la similitud: si és massa alta, probablement és el mateix text
    similitud = difflib.SequenceMatcher(None, orig_net, corr_net).ratio()
    if similitud > 0.98:
        return set()

    # Tokenitza preservant la puntuació com a tokens separats
    def tokenitza(text: str) -> list[str]:
        return _re.findall(
            r"[A-Za-zÀ-ÿàáâãäåæçèéêëìíîïðñòóôõöùúûüýÿ·']+|[^\s\w]|\d+|\S",
            text,
        )

    tokens_orig = tokenitza(orig_net)
    tokens_corr = tokenitza(corr_net)

    if not tokens_orig or not tokens_corr:
        return set()

    paraules_canviades: set[str] = set()
    matcher = difflib.SequenceMatcher(None, tokens_orig, tokens_corr, autojunk=False)

    for tag, i1, i2, j1, j2 in matcher.get_opcodes():
        if tag in ('replace', 'insert'):
            for token in tokens_corr[j1:j2]:
                token_net = token.strip()
                # Inclou només paraules reals (no puntuació aïllada)
                if (
                    token_net
                    and len(token_net) > 1
                    and _re.search(
                        r'[A-Za-zÀ-ÿàáâãäåæçèéêëìíîïðñòóôõöùúûüýÿ]',
                        token_net,
                    )
                ):
                    paraules_canviades.add(token_net.lower())

    return paraules_canviades


def _aplica_highlight_groc_docx(
    paragraph_el,
    text_original: str,
    text_corregit: str,
    ns_w: str = _NS_W_DOC,
) -> None:
    """
    Aplica destacat groc SELECTIU: només als runs que contenen paraules que han
    canviat realment respecte de l'original.
    Usa comparació per paraula completa (\b) per evitar falsos positius
    de subcadena (p. ex. "el" no ha de destacar "el" dins "elegir").
    """
    import re as _re
    from lxml import etree as _et

    paraules_canviades = _obte_paraules_canviades(text_original, text_corregit)
    if not paraules_canviades:
        return

    # Patró per a paraules completes (no subcadenes)
    patró_paraules = _re.compile(
        r'\b(' + '|'.join(_re.escape(p) for p in paraules_canviades) + r')\b',
        _re.IGNORECASE,
    )

    for run in paragraph_el.iter(f"{{{ns_w}}}r"):
        text_run = "".join(t.text for t in run.iter(f"{{{ns_w}}}t") if t.text)
        if not text_run.strip():
            continue

        # Comprova si alguna paraula canviada coincideix com a paraula completa
        if not patró_paraules.search(text_run):
            continue

        rPr = run.find(f"{{{ns_w}}}rPr")
        if rPr is None:
            rPr = _et.SubElement(run, f"{{{ns_w}}}rPr")
            run.insert(0, rPr)
        hl_existent = rPr.find(f"{{{ns_w}}}highlight")
        if hl_existent is not None:
            rPr.remove(hl_existent)
        highlight = _et.SubElement(rPr, f"{{{ns_w}}}highlight")
        highlight.set(f"{{{ns_w}}}val", "yellow")


def _aplica_correccions_docx(arbre, paràgrafs: list[dict], segments_corregits: list[str]) -> bytes:
    """Substitueix el text de cada paràgraf amb la versió corregida i serialitza."""
    from lxml import etree as _et
    for p in paràgrafs:
        idx   = p["index"]
        nou   = segments_corregits[idx] if idx < len(segments_corregits) else ""
        antic = p["text"]
        if nou and nou != antic:
            _substitueix_text_runs(p["node"], nou, _NS_W_DOC)
            _aplica_highlight_groc_docx(p["node"], antic, nou, _NS_W_DOC)
    return _et.tostring(arbre, xml_declaration=True, encoding="UTF-8", standalone=True)


def _extrau_shapes_pptx(xml_bytes: bytes) -> tuple[list[dict], object]:
    """
    Analitza el XML d'una diapositiva PPTX i retorna:
    (shapes, arbre) on shapes = [{si, pi, text, text_marcat, node}]

    text_marcat conté el text amb marcadors ⟦N⟧ entre runs (igual que en DOCX).
    """
    from lxml import etree as _et
    arbre  = _et.fromstring(xml_bytes)
    shapes = []
    for si, txBody in enumerate(arbre.iter(f"{{{_NS_A_DOC}}}txBody")):
        té_formula_shape = _shape_conté_formula_pptx(txBody)
        for pi, p in enumerate(txBody.iter(f"{{{_NS_A_DOC}}}p")):
            text        = _obte_text_paràgraf_doc(p, _NS_A_DOC)
            text_marcat, _ = _extrau_text_amb_marcadors_pptx(p)
            shapes.append({
                "si": si, "pi": pi, "text": text, "text_marcat": text_marcat,
                "té_formula": té_formula_shape, "node": p,
            })
    return shapes, arbre


def _aplica_highlight_groc_pptx(
    paragraph_el,
    text_original: str,
    text_corregit: str,
    ns_a: str = _NS_A_DOC,
) -> None:
    """
    Aplica destacat groc SELECTIU en PPTX per paraules canviades exactes.
    En PPTX s'usa solidFill amb color groc (#FFFF00).
    Usa comparació per paraula completa (\b) per evitar falsos positius.
    """
    import re as _re
    from lxml import etree as _et

    paraules_canviades = _obte_paraules_canviades(text_original, text_corregit)
    if not paraules_canviades:
        return

    patró_paraules = _re.compile(
        r'\b(' + '|'.join(_re.escape(p) for p in paraules_canviades) + r')\b',
        _re.IGNORECASE,
    )

    for run in paragraph_el.iter(f"{{{ns_a}}}r"):
        text_run = "".join(t.text for t in run.iter(f"{{{ns_a}}}t") if t.text)
        if not text_run.strip():
            continue

        if not patró_paraules.search(text_run):
            continue

        rPr = run.find(f"{{{ns_a}}}rPr")
        if rPr is None:
            rPr = _et.SubElement(run, f"{{{ns_a}}}rPr")
            run.insert(0, rPr)
        hl_existent = rPr.find(f"{{{ns_a}}}highlight")
        if hl_existent is not None:
            rPr.remove(hl_existent)
        highlight = _et.SubElement(rPr, f"{{{ns_a}}}highlight")
        solidFill = _et.SubElement(highlight, f"{{{ns_a}}}solidFill")
        srgbClr   = _et.SubElement(solidFill, f"{{{ns_a}}}srgbClr")
        srgbClr.set("val", "FFFF00")


def _aplica_correccions_pptx(arbre, shapes: list[dict],
                               segments_corregits: list[str],
                               offset_inici: int) -> bytes:
    """Substitueix el text de cada shape/paràgraf amb la versió corregida."""
    from lxml import etree as _et
    for k, shape in enumerate(shapes):
        idx_global = offset_inici + k
        nou  = segments_corregits[idx_global] if idx_global < len(segments_corregits) else ""
        if nou and nou != shape["text"]:
            _substitueix_text_runs(shape["node"], nou, _NS_A_DOC)
            _aplica_highlight_groc_pptx(shape["node"], shape["text"], nou, _NS_A_DOC)
    return _et.tostring(arbre, xml_declaration=True, encoding="UTF-8", standalone=True)


def _aplica_highlight_groc_per_run_docx(
    runs_info: list,
    text_pla_orig: str,
    text_pla_corr: str,
) -> None:
    """
    Aplica destacat groc SELECTIU per run (DOCX): comprova el text actual de cada
    run i destaca únicament els que contenen paraules canviades respecte a l'original.
    Més precís que _aplica_highlight_groc_docx (paràgraf sencer) perquè conserva
    el format del text no modificat i minimitza els falsos positius.
    """
    import re as _re
    from lxml import etree as _et

    paraules_canviades = _obte_paraules_canviades(text_pla_orig, text_pla_corr)
    if not paraules_canviades:
        return

    patró = _re.compile(
        r'\b(' + '|'.join(_re.escape(p) for p in paraules_canviades) + r')\b',
        _re.IGNORECASE,
    )
    ns_w = _NS_W_DOC

    for run_info in runs_info:
        run = run_info["node"]
        # Llig el text actual del run (post-correcció)
        text_run = "".join(
            t.text if t.text else ""
            for t in run.iter(f"{{{ns_w}}}t")
        )
        if not text_run.strip() or not patró.search(text_run):
            continue
        rPr = run.find(f"{{{ns_w}}}rPr")
        if rPr is None:
            rPr = _et.SubElement(run, f"{{{ns_w}}}rPr")
            run.insert(0, rPr)
        hl = rPr.find(f"{{{ns_w}}}highlight")
        if hl is not None:
            rPr.remove(hl)
        hl = _et.SubElement(rPr, f"{{{ns_w}}}highlight")
        hl.set(f"{{{ns_w}}}val", "yellow")


def _aplica_highlight_groc_per_run_pptx(
    runs_info: list,
    text_pla_orig: str,
    text_pla_corr: str,
) -> None:
    """
    Aplica destacat groc SELECTIU per run (PPTX): usa solidFill #FFFF00 per als
    runs de DrawingML que contenen paraules canviades.
    """
    import re as _re
    from lxml import etree as _et

    paraules_canviades = _obte_paraules_canviades(text_pla_orig, text_pla_corr)
    if not paraules_canviades:
        return

    patró = _re.compile(
        r'\b(' + '|'.join(_re.escape(p) for p in paraules_canviades) + r')\b',
        _re.IGNORECASE,
    )
    ns_a = _NS_A_DOC

    for run_info in runs_info:
        run = run_info["node"]
        text_run = "".join(
            t.text if t.text else ""
            for t in run.iter(f"{{{ns_a}}}t")
        )
        if not text_run.strip() or not patró.search(text_run):
            continue
        rPr = run.find(f"{{{ns_a}}}rPr")
        if rPr is None:
            rPr = _et.SubElement(run, f"{{{ns_a}}}rPr")
            run.insert(0, rPr)
        hl = rPr.find(f"{{{ns_a}}}highlight")
        if hl is not None:
            rPr.remove(hl)
        hl = _et.SubElement(rPr, f"{{{ns_a}}}highlight")
        solidFill = _et.SubElement(hl, f"{{{ns_a}}}solidFill")
        srgbClr   = _et.SubElement(solidFill, f"{{{ns_a}}}srgbClr")
        srgbClr.set("val", "FFFF00")


def _neteja_encoding_valor(text: str) -> str:
    """
    Corregeix caràcters UTF-8 mal interpretats com Latin-1.
    Exemple: 'Ã©s' → 'és', 'Ã ' → 'à', etc.
    """
    if not text:
        return text
    try:
        return text.encode("latin-1").decode("utf-8")
    except (UnicodeEncodeError, UnicodeDecodeError):
        return text


async def _corregeix_segments_claude(
    segments: list[str],
    api_key: str,
    mida_lot: int = 3,
    segments_marcats: list[str] | None = None,
    segments_amb_formula: list[bool] | None = None,
) -> list[str]:
    """
    Envia segments de text a Claude Sonnet per a correcció normativa en lots.
    Usa PROMPT_CORRECCIO_SISTEMA com a system prompt per garantir l'aplicació
    de totes les normes del valencià universitari (UV).
    Retorna la llista de segments corregits conservant l'ordre i els índexs.

    segments_marcats    : si s'especifica, els valors JSON enviats a Claude
                          contenen marcadors ⟦N⟧ entre runs.
    segments_amb_formula: si s'especifica, els segments True s'exclouen de la
                          correcció automàtica per preservar equacions OMML/OLE.
    """
    import anthropic as _ant
    import json as _js

    client             = _ant.Anthropic(api_key=api_key)
    segments_corregits = list(segments)

    índexs_a_corregir = [
        i for i, s in enumerate(segments)
        if s and s.strip() and len(s.strip()) > 3
        # Exclou segments que contenen fórmules matemàtiques
        and not (
            segments_amb_formula
            and i < len(segments_amb_formula)
            and segments_amb_formula[i]
        )
    ]

    log.warning(
        "[CORRECCIO] Total segments: %d | A corregir: %d | Clau prefix: %s",
        len(segments), len(índexs_a_corregir),
        (api_key[:12] + "...") if api_key else "BUIDA",
    )

    if not índexs_a_corregir:
        log.warning("[CORRECCIO] Cap segment a corregir — tots filtrats")
        return segments_corregits

    errors_lots: list[str] = []

    for inici in range(0, len(índexs_a_corregir), mida_lot):
        lot_índexs = índexs_a_corregir[inici: inici + mida_lot]
        # Si tenim segments marcats (amb ⟦N⟧), els enviem a Claude per preservar runs
        usa_marcadors = segments_marcats is not None and len(segments_marcats) == len(segments)
        lot_textos = {
            str(i): (segments_marcats[i] if usa_marcadors else segments[i])
            for i in lot_índexs
        }
        num_lot    = inici // mida_lot + 1

        log.warning(
            "[CORRECCIO] Lot %d: %d segments (índexs %d–%d)",
            num_lot, len(lot_índexs), lot_índexs[0], lot_índexs[-1],
        )

        instruccions_marcadors = (
            "\n\nATENCIÓ — MARCADORS DE FORMAT ⟦N⟧:\n"
            "Els textos contenen marcadors ⟦0⟧, ⟦1⟧, ⟦2⟧... que delimiten fragments "
            "de text amb tipografia diferent (negreta, cursiva, color, etc.).\n"
            "REGLA CRÍTICA: PRESERVA els marcadors EXACTAMENT on estan, sense moure'ls, "
            "duplicar-los ni eliminar-los. Corregeix ÚNICAMENT el text que hi ha entre "
            "els marcadors. Exemple:\n"
            "  Entrada : ⟦0⟧Este ⟦1⟧text servix d'exemple\n"
            "  Sortida : ⟦0⟧Aquest ⟦1⟧text serveix d'exemple\n"
        ) if usa_marcadors else ""

        prompt_usuari = f"""Ets el corrector i posteditor lingüístic expert del SLPL de la Universitat de València. La teua tasca és aplicar una correcció i postedició EXHAUSTIVA, RIGOROSA i SISTEMÀTICA al text en valencià, basant-te en els quatre corpus normatius del teu sistema (Blocs A, B, C i D).{instruccions_marcadors}

METODOLOGIA OBLIGATÒRIA — ANALITZA EN ORDRE CADA CATEGORIA:

════════════════════════════════════════════════════════════
BLOC A — MORFOLOGIA (Criteris lingüístics UV + IEC GIEC/OIEC 2016-2023)
════════════════════════════════════════════════════════════

A1. DEMOSTRATIUS — usa SEMPRE les formes reforçades en registre formal escrit:
□ este/esta/estos/estes → aquest/aquesta/aquests/aquestes (1r grau)
□ eixe/eixa/eixos/eixes → aqueix/aqueixa/aqueixos/aqueixes (2n grau)
□ allò que és bo → (mai "lo bo") [vg. A8 Lo neutre]

A2. POSSESSIUS — formes amb -u- obligatòries:
□ seva/seves → seua/seues
□ meva/meves → meua/meues
□ teva/teves → teua/teues
□ Evita possessiu innecessari: "la biblioteca obre l'horari" (NO "el seu horari")
□ Evita possessiu amb parts del cos en CD agentiu: "va alçar el cap" (NO "el seu cap")
□ Evita possessiu amb noms de parentiu si el referent és clar: "acompanya el marit" (NO "el seu marit")

A3. VERBS INCOATIUS — present indicatiu i subjuntiu:
□ Indicatiu: servisc/serveixes/serveix/servim/serviu/serveixen (NO servix/servixen)
□ Subjuntiu: servisca/servisques/servisca/servim/serviu/servisquen (NO servixca/servixquen)

A4. PARTICIPIS — formes regulars preferibles:
□ complert→complit, ofert→oferit, establert→establit, omplert→omplit, sofert→sofrit, suplert→suplit
□ pertanyut→pertangut, planyut→plangut
□ Participi de ser: estat→sigut (preferible)

A5. INFINITIUS:
□ tindre→tenir, vindre→venir, caber→cabre, caler→caldre, doler→doldre, valer→valdre

A6. NUMERALS:
□ huit→vuit, díhuit→divuit, dèsset→disset, dènou→dinou
□ quint→cinquè, sext→sisè, dècim→desè
□ dos/dues: concorda en gènere

A7. PLURALS:
□ -ns→-s en registre formal: hòmens→homes, jóvens→joves, màrgens→marges, térmens→termes
□ -sc/-st/-xt/-ig: discos, gustos, textos, rojos (EXCEPCIÓ: raigs X, tests)
□ Femení professions amb -a: advocada, arquitecta, ministra, presidenta

A8. ACCENTUACIÓ — sistema GENERAL (no occidental):
□ anglés→anglès, francés→francès, interés→interès, permés→permès, compromés→compromès
□ ordinals: cinqué→cinquè, sisé→sisè
□ substantius: café→cafè, comité→comitè, mercé→mercè
□ infinitius: conéixer→conèixer, meréixer→merèixer, véncer→vèncer
□ imperfets ind.: féiem→fèiem, déiem→dèiem
□ 3a pl. 2a conj.: aprén→aprèn, comprén→comprèn, depén→depèn
□ EXCEPCIÓ accent agut: atén, entén, pretén, encén (precedits de t/c)
□ EXCEPCIÓ accent agut: congrés/congressos, exprés, procés, progrés (però interès/interessos)
□ Accents diacrítics (IEC OIEC 2017): sé/se, és/es, més/mes, bé/be, té/te, sí/si, mà/ma, món/mon, pèl/pel, déu/deu, sòl/sol, són/son, mén/men, bés/bes, mós/mos — aplica'ls correctament
□ NO porten accent diacrític: dona (verb donar), feu (verb fer), sec, fora, soc (verb ser)

A9. GRAFIES:
□ tl/tll: motle (NO motlle), espatla (NO espatlla), vetlar (NO vetllar)
□ EXCEPCIÓ sempre tll: bitllet, rotllo, butlletí, ratlla i derivats

A10. ALTERNANCES VOCÀLIQUES (Criteris UV):
□ a/e: nadar (NO nedar), nàixer (NO néixer), traure (NO treure), xarrar (NO xerrar)
□ e/o: fenoll (NO fonoll), redó (NO rodó), renyó (NO ronyó), arredonir (NO arrodonir)

════════════════════════════════════════════════════════════
BLOC B — SINTAXI (Gramàtica Zero + IEC GIEC 2016-2023)
════════════════════════════════════════════════════════════

B1. ORDRE DEL SINTAGMA NOMINAL (SN) — CRITERI IEC GIEC:
□ En català l'ordre habitual és SUBSTANTIU + ADJECTIU: "una reunió important", "un informe detallat"
□ L'ordre ADJECTIU + SUBSTANTIU (calc del castellà) és incorrecte en la majoria de contextos
□ Detecta i corregeix: "el proper curs"→"el curs vinent/proper", "la present comunicació"→"la comunicació present", "la corresponent documentació"→"la documentació corresponent", "l'esmentada resolució"→"la resolució esmentada"
□ Excepció: "el mateix informe" (demostratiu amb valor d'identitat), "certes persones" (cert en sentit d'indefinit)

B2. GERUNDI DE POSTERIORITAT/CONSEQÜÈNCIA — SEMPRE INCORRECTE:
□ "Va caure trencant-se una cama" → "Va caure i es va trencar una cama"
□ "Va dimitir, provocant una crisi" → "Va dimitir, la qual cosa va provocar una crisi"
□ GERUNDI DE CAUSA (IEC GIEC): "En ser tan alt" / "Al no tenir" → "Com que era tan alt" / "Com que no tenia"

B3. HAVER-HI — SEMPRE EN SINGULAR I AMB HI:
□ "hi han molts errors" → "hi ha molts errors"
□ "hi havien" → "hi havia"
□ "ha hagut" → "hi ha hagut"

B4. LO NEUTRE — SEMPRE INCORRECTE EN REGISTRES FORMALS:
□ "lo important" → "allò que és important" / "el més important"
□ "lo bo" → "allò que és bo" / "el bo"
□ "lo que fas" → "el que fas" / "allò que fas"

B5. ALGO — SEMPRE INCORRECTE:
□ algo (valor indeterminat) → "alguna cosa"
□ algo (valor quantitatiu) → "un poc" / "una mica"
□ algo (valor formal) → "quelcom"

B6. CAIGUDA DE PREPOSICIÓ DAVANT QUE:
□ "estic segur de que" → "estic segur que"
□ "confiem en que" → "confiem que"
□ "la idea de que" → "la idea que"
□ "el fet de que" → "el fet que"

B7. COMPLEMENT DIRECTE DE PERSONA AMB A:
□ "vaig veure al director" → "vaig veure el director"
□ "l'article cita a Marx" → "l'article cita Marx"

B8. PRONOM EN OBLIGATORI:
□ "no tinc" (referit a CD indeterminat) → "no en tinc"
□ "me vaig" → "me'n vaig"

B9. PRONOM HI OBLIGATORI:
□ "no he estat mai allà" → "no hi he estat mai"
□ "s'ha referit a ell" → "s'hi ha referit"

B10. DEGUT A (CAUSA) — INCORRECTE:
□ "degut a la pluja" → "a causa de la pluja" / "per la pluja"
□ "degut a que" → "perquè" / "a causa que"

B11. EN BASE A — EVITAR:
□ "en base al reglament" → "d'acord amb el reglament" / "a partir del reglament"

B12. ANAR A + INFINITIU PER A FUTUR (IEC GIEC):
□ "anem a comentar" → "comentarem"
□ "va a fer" → "farà"
□ "la sessió va a començar" → "la sessió està a punt de començar" / "la sessió començarà"

B13. CONDICIONAL INCORRECTE (IEC GIEC):
□ "el diagnòstic indica que la lesió seria greu" → "que la lesió és greu"

B14. FUTUR AMB VALOR DE PROBABILITAT (IEC GIEC):
□ "tindrà gana" (probabilitat) → "deu tenir gana" / "potser té gana"

B15. SUBJUNTIU INCORRECTE DESPRÉS D'ADVERBIS DE DUBTE (IEC GIEC):
□ "potser siga" → "potser és" (potser + indicatiu en registres formals)
□ EXCEPCIÓ: "pot ser que" + subjuntiu és correcte

B16. ÚS NO ACCEPTABLE DE MATEIX (IEC GIEC):
□ "els resultats de les mateixes" → "els seus resultats"
□ "la pròpia senyora Pia" → "la mateixa senyora Pia" / "la senyora Pia mateixa"

B17. HAGUÉS en comptes de HAURIA (principal condicional):
□ "no s'hagués perdut" → "no s'hauria perdut"

B18. MALGRAT + VERB sense QUE:
□ "malgrat plovia" → "malgrat que plovia"

B19. PER A QUÈ / PERQUÈ:
□ "us hem citat per a que digueu" → "perquè digueu"

B20. ORDRES/INSTRUCCIONS amb infinitiu:
□ "No fumar" → "No fumeu"; "Veure l'annex" → "Vegeu l'annex"

B21. SI NO / SINÓ:
□ "no era llest sino astut" → "sinó astut" (adversativa)

B22. VARIS — SEMPRE INCORRECTE:
□ "varis errors" → "diversos errors"

B23. LOCALITZACIÓ A / EN:
□ A + noms propis i article determinat: "visc a Alacant", "treballe al camp"
□ EN en la resta: "en alguna platja", "en una plaça"

B24. MÉS BÉ / MÉS AVIAT:
□ "és més bé un pobre xicot" → "és més aviat un pobre xicot"

════════════════════════════════════════════════════════════
BLOC C — LÈXIC I ESTIL ADMINISTRATIU (Manual de documents UV + Criteris UV)
════════════════════════════════════════════════════════════

C1. LÈXIC PREFERIT:
□ hui→avui, vore→veure, mentres→mentre, servici→servei, vacacions→vacances
□ desenrotllar→desenvolupar, ferramenta→eina (figurat), mitat→meitat, sendemà→endemà
□ juí→judici, perjuí→perjudici, vàries→diverses, xicotet/menut→petit (davant nom)

C2. TOPÒNIMS — forma tradicional valenciana:
□ Orihuela→Oriola, Zaragoza→Saragossa, Cádiz→Cadis, London→Londres

C3. TRACTAMENT PERSONAL — VÓS (preferible en documents formals):
□ Evitar vostè en documents administratius formals
□ NO "don Pere" / "dona Montserrat" → "el senyor Pere" / "la senyora Montserrat"

C4. FÓRMULES ADMINISTRATIVES — eliminar arcaismes i calcs:
□ "en base a" → "d'acord amb"; "a nivell de" → "a escala de"; "degut a" → "a causa de"
□ "al respecte" → "respecte a això"; "dur a terme" → "realitzar" / "fer"
□ "si escau", "d'ara endavant", "cal", "per endavant", "d'acord amb", "als efectes oportuns"
□ "la persona interessada" (NO "el/la interesado/a"); "quan pertoque" (NO "en su día")

C5. MAJÚSCULES I MINÚSCULES:
□ Càrrecs: MINÚSCULA (rector/a, conseller/a, president/a)
□ Institucions: MAJÚSCULA (la Universitat de València, la Generalitat)
□ Documents oficials: MAJÚSCULA inicial (la Llei orgànica 6/2001)
□ Dies, mesos, estacions: MINÚSCULA; Gentilicis: MINÚSCULA

C6. ESTIL:
□ Frases màx. 3 línies; eliminar redundàncies i parelles de sinònims innecessàries
□ Usar formes actives i directes en lloc de passives

C7. ABREVIACIONS:
□ Ordinals: 1r, 2n, 3r, 4t, 5è (NO 1º, 2º)
□ Sigles: NO apostrofar davant alfabètiques (la UPV); SÍ davant sil·làbiques (l'IVA)
□ Xifres: COMA decimal (43,3); PUNT de milers (2.076.000)

C8. LLENGUATGE IGUALITARI:
□ Col·lectius genèrics: l'alumnat, el professorat, la persona interessada
□ Càrrecs: adaptar al gènere real (la rectora, la ministra, la degana)

════════════════════════════════════════════════════════════
BLOC D — NOVETATS IEC 2016-2023 (GIEC, OIEC, GEIEC, GBU)
════════════════════════════════════════════════════════════

D1. ACCENT DIACRÍTIC REDUÏT (OIEC 2017):
□ Accent diacrític REDUÏT a 15 mots: bé/be, déu/deu, és/es, mà/ma, més/mes, món/mon, pèl/pel, sé/se, sí/si, sòl/sol, són/son, té/te, ús/us, déus/deus, béns/bens, pèls/pels, sís/sis, sòls/sols
□ NO porten accent diacrític (cal eliminar-los si hi eren): bota, coc, dona (verb), feu (verb), fora, soc (verb ser), sec

D2. ERRADICAR S'ESCRIU AMB ERR- (OIEC 2017):
□ eradicar → erradicar, eradicació → erradicació

D3. FORMES NO ACCEPTABLES RECENTS (IEC GBU):
□ "don Pere" / "dona Montserrat" → "el senyor Pere" / "la senyora Montserrat"
□ Verbs psicològics: "li va afectar" → "la va afectar"; "les aranyes em donen por" → "em fan por"
□ Tenir de + infinitiu: "tenim de fer" → "hem de fer"
□ "pròpia" en lloc de "mateixa": "la pròpia directora" → "la mateixa directora" / "la directora mateixa"

D4. ORACIONS COPULATIVES — SER vs ESTAR:
□ Ser: propietats inherents, identitat, classificació ("és alta", "és metge")
□ Estar: estats contingents, localitzacions ("està content", "està a casa")
□ En registres formals NO usar restar per a durada d'un estat: "restarà tancat" → "estarà tancat"

D5. CONCORDANÇA NOMINAL EN APOSICIÓ (IEC GIEC):
□ Noms fixats en aposició → sense concordança en registres formals: "les dates límit", "els episodis pilot"

════════════════════════════════════════════════════════════
INSTRUCCIONS DE RESPOSTA
════════════════════════════════════════════════════════════

IMPORTANT: Per a dubtes tipogràfics específics aplica sempre els criteris de l'Optimot i del Manual d'estil de les universitats valencianes.

REGLES:
- Retorna EXACTAMENT el mateix JSON amb les mateixes claus numèriques.
- Si un text ja és completament correcte, retorna'l IDÈNTIC sense cap canvi.
- NO afegeixis cap text fora del JSON.
- Preserva noms propis, sigles, xifres i puntuació estructural.
- NO canvies la longitud substancialment (±25% màxim).
- Aplica TOTES les correccions detectades, incloent les subtils (ordre SN, accentuació, possessius, majúscules/minúscules).
- Si detects un error però no n'estàs segur de la correcció exacta, aplica la forma més segura i normativa.

JSON d'entrada:
{_js.dumps(lot_textos, ensure_ascii=False, indent=2)}

Retorna ÚNICAMENT el JSON corregit:"""

        text_resp = ""
        try:
            resposta = client.messages.create(
                model      = "claude-sonnet-4-6",
                max_tokens = 4096,
                system     = PROMPT_CORRECCIO_SISTEMA,
                messages   = [{"role": "user", "content": prompt_usuari}],
            )
            text_resp = resposta.content[0].text.strip()
            log.warning(
                "[CORRECCIO] Lot %d resposta raw (primers 200 car.): %r",
                num_lot, text_resp[:200],
            )

            # Corregeix text UTF-8 mal interpretat com Latin-1 (Ã© → é, etc.)
            def _neteja_encoding(t: str) -> str:
                try:
                    return t.encode("latin-1").decode("utf-8")
                except (UnicodeEncodeError, UnicodeDecodeError):
                    return t

            if "Ã" in text_resp or "â€" in text_resp:
                text_resp = _neteja_encoding(text_resp)

            # Neteja possible markdown ```json ... ```
            if text_resp.startswith("```"):
                text_resp = re.sub(r"```(?:json)?\s*", "", text_resp)
                text_resp = text_resp.replace("```", "").strip()

            lot_corregit = _js.loads(text_resp)
            canvis = 0

            for i in lot_índexs:
                clau = str(i)
                if clau in lot_corregit and isinstance(lot_corregit[clau], str):
                    valor = _neteja_encoding_valor(lot_corregit[clau].strip())
                    if valor:
                        if valor != segments[i]:
                            canvis += 1
                        segments_corregits[i] = valor

            log.warning("[CORRECCIO] Lot %d OK — %d canvis aplicats", num_lot, canvis)

        except _ant.AuthenticationError as exc:
            msg = f"Lot {num_lot} ERROR AUTENTICACIÓ: {exc}"
            log.error("[CORRECCIO] %s", msg)
            errors_lots.append(msg)
            break  # Inútil continuar si la clau és incorrecta

        except _ant.APIError as exc:
            msg = f"Lot {num_lot} ERROR API Anthropic ({type(exc).__name__}): {exc}"
            log.error("[CORRECCIO] %s", msg)
            errors_lots.append(msg)

        except _js.JSONDecodeError as exc:
            msg = (
                f"Lot {num_lot} ERROR JSON: {exc} | "
                f"Text raw: {repr(text_resp[:300])}"
            )
            log.error("[CORRECCIO] %s", msg)
            errors_lots.append(msg)

        except Exception as exc:
            msg = f"Lot {num_lot} ERROR INESPERAT ({type(exc).__name__}): {exc}"
            log.error("[CORRECCIO] %s", msg)
            errors_lots.append(msg)

    if errors_lots:
        log.warning("[CORRECCIO] Errors totals: %d: %s", len(errors_lots), errors_lots)

    return segments_corregits


async def _processa_docx_correccio(fitxer_bytes: bytes, api_key: str) -> bytes:
    """Corregeix un DOCX i retorna el DOCX corregit preservant el format.

    Estratègia (CANVI 1+2+3):
    - Passada 1: extrau text pla + text marcat (⟦N⟧ per run) de cada paràgraf.
    - Envia lots de 3 segments a Claude amb instruccions de preservar marcadors.
    - Passada 2: per cada paràgraf, intenta reinjectar el text corregit run a run
      (si Claude ha preservat els marcadors). Si no, usa el mètode de fallback
      (_substitueix_text_runs). Destaca per run quan els marcadors s'han preservat.
    """
    tots_segments: list[str]         = []
    tots_segments_marcats: list[str] = []
    tots_segments_formula: list[bool] = []
    mapa_fitxer: dict[str, tuple]    = {}

    # ── Passada 1: extrau tots els segments de text ───────────────────────────
    with zipfile.ZipFile(io.BytesIO(fitxer_bytes)) as zin:
        noms = zin.namelist()
        for nom in noms:
            if nom in _FITXERS_XML_DOCX_FIXES or _PATRÓ_FITXERS_DOCX_EXTRA.match(nom):
                try:
                    xml_bytes = zin.read(nom)
                    paràgrafs, arbre = _extrau_paràgrafs_docx(xml_bytes)
                    inici = len(tots_segments)
                    for p in paràgrafs:
                        tots_segments.append(p["text"])
                        tots_segments_marcats.append(p["text_marcat"])
                        tots_segments_formula.append(p.get("té_formula", False))
                    mapa_fitxer[nom] = (paràgrafs, arbre, inici, xml_bytes)
                except Exception as exc:
                    log.warning("Error extraient paràgrafs de '%s': %s", nom, exc)
                    mapa_fitxer[nom] = None

    # ── Corregeix tots els segments en lots de 3 (CANVI 3) ────────────────────
    corregits = await _corregeix_segments_claude(
        tots_segments, api_key,
        segments_marcats=tots_segments_marcats,
        segments_amb_formula=tots_segments_formula,
    )

    # ── Passada 2: reconstrueix el ZIP ────────────────────────────────────────
    buf_eixida = io.BytesIO()
    with zipfile.ZipFile(io.BytesIO(fitxer_bytes)) as zin, \
         zipfile.ZipFile(buf_eixida, "w", zipfile.ZIP_DEFLATED) as zout:
        for nom in zin.namelist():
            if nom in mapa_fitxer and mapa_fitxer[nom] is not None:
                paràgrafs, arbre, inici, xml_original = mapa_fitxer[nom]
                segs_fitxer  = corregits[inici: inici + len(paràgrafs)]
                segs_per_idx = {p["index"]: segs_fitxer[k] for k, p in enumerate(paràgrafs)}

                from lxml import etree as _et
                arbre_treballat = _et.fromstring(xml_original)
                for i, p_node in enumerate(arbre_treballat.iter(f"{{{_NS_W_DOC}}}p")):
                    nou = segs_per_idx.get(i, "")
                    if i >= len(paràgrafs) or not nou:
                        continue
                    antic = paràgrafs[i]["text"]
                    # Compara el text pla (sense marcadors) amb l'original
                    if _text_pla(nou) == antic:
                        continue

                    # CANVI 1: Extrau runs_info del nou arbre (frescos, sense refs. mortes)
                    _, runs_info = _extrau_text_amb_marcadors_docx(p_node)
                    aplicat = _aplica_text_marcat(runs_info, nou, _NS_W_DOC)

                    if aplicat:
                        # CANVI 2: Destaca per run (més precís)
                        nou_pla = _text_pla(nou)
                        _aplica_highlight_groc_per_run_docx(runs_info, antic, nou_pla)
                    else:
                        # Fallback: tot el text al primer run
                        nou_pla = _text_pla(nou)
                        _substitueix_text_runs(p_node, nou_pla, _NS_W_DOC)
                        _aplica_highlight_groc_docx(p_node, antic, nou_pla, _NS_W_DOC)

                xml_corregit = _et.tostring(
                    arbre_treballat, xml_declaration=True, encoding="UTF-8", standalone=True
                )
                # Força l'idioma ca-ES
                xml_text = xml_corregit.decode("utf-8")
                xml_text = re.sub(r'w:lang\s+w:val="[^"]*"', 'w:lang w:val="ca-ES"', xml_text)
                zout.writestr(nom, xml_text.encode("utf-8"))
            else:
                zout.writestr(nom, zin.read(nom))

    buf_eixida.seek(0)
    return buf_eixida.read()


async def _processa_pptx_correccio(fitxer_bytes: bytes, api_key: str) -> bytes:
    """Corregeix un PPTX i retorna el PPTX corregit preservant el format.

    Estratègia (CANVI 1+2+3): equivalent a _processa_docx_correccio però per a PPTX.
    Usa marcadors ⟦N⟧ per reinjectar el text corregit run a run (DrawingML).
    """
    tots_segments: list[str]          = []
    tots_segments_marcats: list[str]  = []
    tots_segments_formula: list[bool] = []
    mapa_fitxer: dict[str, tuple]     = {}

    # ── Passada 1: extrau tots els segments ───────────────────────────────────
    with zipfile.ZipFile(io.BytesIO(fitxer_bytes)) as zin:
        noms = zin.namelist()
        for nom in noms:
            if _PATRÓ_FITXERS_PPTX.match(nom):
                try:
                    xml_bytes = zin.read(nom)
                    shapes, arbre = _extrau_shapes_pptx(xml_bytes)
                    inici = len(tots_segments)
                    for s in shapes:
                        tots_segments.append(s["text"])
                        tots_segments_marcats.append(s["text_marcat"])
                        tots_segments_formula.append(s.get("té_formula", False))
                    mapa_fitxer[nom] = (shapes, inici, xml_bytes)
                except Exception as exc:
                    log.warning("Error extraient shapes de '%s': %s", nom, exc)
                    mapa_fitxer[nom] = None

    # ── Corregeix en lots de 3 (CANVI 3) ──────────────────────────────────────
    corregits = await _corregeix_segments_claude(
        tots_segments, api_key,
        segments_marcats=tots_segments_marcats,
        segments_amb_formula=tots_segments_formula,
    )

    # ── Passada 2: reconstrueix el ZIP ────────────────────────────────────────
    buf_eixida = io.BytesIO()
    with zipfile.ZipFile(io.BytesIO(fitxer_bytes)) as zin, \
         zipfile.ZipFile(buf_eixida, "w", zipfile.ZIP_DEFLATED) as zout:
        for nom in zin.namelist():
            if nom in mapa_fitxer and mapa_fitxer[nom] is not None:
                shapes, inici, xml_original = mapa_fitxer[nom]
                from lxml import etree as _et
                arbre_treballat = _et.fromstring(xml_original)
                k = 0
                for txBody in arbre_treballat.iter(f"{{{_NS_A_DOC}}}txBody"):
                    for p_node in txBody.iter(f"{{{_NS_A_DOC}}}p"):
                        if k >= len(shapes):
                            break
                        idx_global = inici + k
                        nou  = corregits[idx_global] if idx_global < len(corregits) else ""
                        antic = shapes[k]["text"]

                        if nou and _text_pla(nou) != antic:
                            # CANVI 1: Extrau runs_info frescos i intenta aplicar marcadors
                            _, runs_info = _extrau_text_amb_marcadors_pptx(p_node)
                            aplicat = _aplica_text_marcat(runs_info, nou, _NS_A_DOC)

                            if aplicat:
                                # CANVI 2: Destaca per run
                                nou_pla = _text_pla(nou)
                                _aplica_highlight_groc_per_run_pptx(runs_info, antic, nou_pla)
                            else:
                                # Fallback
                                nou_pla = _text_pla(nou)
                                _substitueix_text_runs(p_node, nou_pla, _NS_A_DOC)
                                _aplica_highlight_groc_pptx(p_node, antic, nou_pla, _NS_A_DOC)
                        k += 1

                xml_corregit = _et.tostring(
                    arbre_treballat, xml_declaration=True, encoding="UTF-8", standalone=True
                )
                xml_text = xml_corregit.decode("utf-8")
                xml_text = re.sub(r'\blang="[^"]*"', 'lang="ca-ES"', xml_text)
                zout.writestr(nom, xml_text.encode("utf-8"))
            else:
                zout.writestr(nom, zin.read(nom))

    buf_eixida.seek(0)
    return buf_eixida.read()


@app.post(
    "/corregeix-document",
    summary = "Corregeix un document DOCX o PPTX en valencià preservant el format",
    tags    = ["Correcció"],
)
async def corregeix_document(fitxer: UploadFile = File(...)) -> Response:
    """
    Corregeix normativament un document .docx o .pptx en català/valencià
    preservant el format, la tipografia, els colors i l'estructura originals.

    Estratègia:
    1. Obre el ZIP intern del document
    2. Extreu el text de cada paràgraf (<w:p>/<a:p>)
    3. Envia tots els segments a Claude Sonnet (lots de 15)
    4. Reinjecta els textos corregits als nodes XML originals
    5. Retorna el document reconstruït
    """
    import os

    api_key = _obte_api_key_anthropic()
    if not api_key:
        raise HTTPException(
            status_code=503,
            detail=(
                "Clau API d'Anthropic no configurada. "
                "Introdueix-la al panell de configuració de la pestanya 'Correcció'."
            ),
        )

    nom       = fitxer.filename or "document"
    extensió  = Path(nom).suffix.lower()

    if extensió not in (".docx", ".pptx"):
        raise HTTPException(
            status_code=415,
            detail=(
                f"Extensió '{extensió}' no admesa. "
                "Només s'accepten fitxers .docx i .pptx."
            ),
        )

    contingut = await fitxer.read()
    if len(contingut) > MAX_MIDA_FITXER:
        raise HTTPException(
            status_code=413,
            detail=(
                f"El fitxer supera el límit de 150 MB "
                f"({len(contingut) / 1_048_576:.1f} MB rebuts)."
            ),
        )

    log.info(
        "POST /corregeix-document — fitxer='%s' mida=%.1f KB",
        nom, len(contingut) / 1024,
    )

    inici = time.perf_counter()
    try:
        if extensió == ".docx":
            resultat   = await _processa_docx_correccio(contingut, api_key)
            media_type = (
                "application/vnd.openxmlformats-officedocument"
                ".wordprocessingml.document"
            )
        else:
            resultat   = await _processa_pptx_correccio(contingut, api_key)
            media_type = (
                "application/vnd.openxmlformats-officedocument"
                ".presentationml.presentation"
            )

        temps_ms    = int((time.perf_counter() - inici) * 1000)
        nom_sortida = genera_nom_arxiu(nom, sufix="CORR_VAL")
        log.info(
            "Document corregit en %d ms → '%s' (%d bytes)",
            temps_ms, nom_sortida, len(resultat),
        )
        return Response(
            content    = resultat,
            media_type = media_type,
            headers    = {
                "Content-Disposition": f'attachment; filename="{nom_sortida}"',
                "Content-Length":      str(len(resultat)),
                "X-Temps-Ms":          str(temps_ms),
            },
        )

    except HTTPException:
        raise
    except Exception as exc:
        import traceback
        detall = f"{type(exc).__name__}: {exc}\n{traceback.format_exc()}"
        log.exception("Error corregint el document '%s': %s", nom, exc)
        raise HTTPException(
            status_code=500,
            detail=f"Error en la correcció del document: {detall}",
        )


# ═══════════════════════════════════════════════════════════════════════════════
# NOUS ENDPOINTS: Traducció amb Claude Sonnet + EN↔VA + Correcció millorada
# ═══════════════════════════════════════════════════════════════════════════════


async def _crida_claude_amb_cache(
    system_blocks: list[dict],
    missatge_usuari: str,
    api_key: str,
    max_tokens: int = 8192,
    model: str = "claude-sonnet-4-6",
) -> str:
    """
    Funció auxiliar per a cridar Claude amb prompt caching.

    Args:
        system_blocks: blocs del system prompt (amb cache_control)
        missatge_usuari: missatge de l'usuari
        api_key: clau API d'Anthropic
        max_tokens: tokens màxims de resposta
        model: model a usar

    Retorna:
        Text de la resposta de Claude
    """
    import anthropic as _ant

    client = _ant.Anthropic(api_key=api_key)

    resposta = client.messages.create(
        model=model,
        max_tokens=max_tokens,
        system=system_blocks,
        messages=[{"role": "user", "content": missatge_usuari}],
    )

    text_resposta = resposta.content[0].text

    # Neteja codificació UTF-8 mal interpretada
    if "Ã" in text_resposta or "â€" in text_resposta:
        try:
            text_resposta = text_resposta.encode("latin-1").decode("utf-8")
        except (UnicodeEncodeError, UnicodeDecodeError):
            pass

    log.info(
        "Claude [%s] — input: %d tokens, output: %d tokens, cache: %s/%s",
        model,
        resposta.usage.input_tokens if hasattr(resposta.usage, 'input_tokens') else 0,
        resposta.usage.output_tokens if hasattr(resposta.usage, 'output_tokens') else 0,
        getattr(resposta.usage, 'cache_creation_input_tokens', 'N/A'),
        getattr(resposta.usage, 'cache_read_input_tokens', 'N/A'),
    )

    return text_resposta


def _neteja_resposta_claude(text: str) -> str:
    """
    Elimina text parasitari que Claude pot afegir quan rep segments curts.
    Patrons coneguts: preguntes, ofertes d'ajuda, peticions de més text.
    """
    # Patrons parasitaris (anglés i català/valencià)
    patrons = [
        r"(?i)please provide.*?(?:full text|more context|complete text).*?[.?!]?\s*",
        r"(?i)could you (?:share|provide|send).*?[.?!]?\s*",
        r"(?i)it seems (?:only|like).*?(?:included|provided).*?[.?!]?\s*",
        r"(?i)I('d| would) (?:need|like).*?(?:full|complete|more).*?[.?!]?\s*",
        r"(?i)the (?:text|segment) (?:appears|seems).*?(?:incomplete|short|brief).*?[.?!]?\s*",
        r"(?i)(?:here is|here's) the translation:?\s*",
        r"(?i)si us plau.*?(?:proporcion|compart|envi).*?[.?!]?\s*",
        r"(?i)sembla que (?:només|nomé).*?[.?!]?\s*",
        r"(?i)necessit(?:e|o|aria).*?(?:text complet|més context).*?[.?!]?\s*",
    ]
    resultat = text
    for patro in patrons:
        resultat = re.sub(patro, '', resultat)
    # Neteja espais residuals
    resultat = re.sub(r'\n{3,}', '\n\n', resultat).strip()
    return resultat if resultat else text  # Si queda buit, retorna l'original


# ─── Endpoint: POST /tradueix-claude (traducció ES→VA amb Claude) ────────────

class PeticioTradueixClaude(BaseModel):
    text: str = Field(
        ...,
        min_length=1,
        max_length=100_000,
        description="Text en castellà a traduir.",
    )


@app.post(
    "/tradueix-claude",
    summary="Tradueix text castellà→valencià amb Claude Sonnet",
    tags=["Traducció"],
)
async def tradueix_claude(peticio: PeticioTradueixClaude):
    """
    Tradueix text del castellà al valencià estàndard universitari
    usant Claude Sonnet amb el prompt normatiu consolidat i prompt caching.
    """
    api_key = _obte_api_key_anthropic()
    if not api_key:
        raise HTTPException(
            status_code=503,
            detail="Clau API d'Anthropic no configurada.",
        )

    text = peticio.text.strip()
    if not text:
        raise HTTPException(status_code=422, detail="El camp 'text' no pot estar buit.")

    inici = time.perf_counter()

    try:
        system_blocks, prefix = construeix_prompt_traduccio_es_va()
        traduccio = await _crida_claude_amb_cache(
            system_blocks=system_blocks,
            missatge_usuari=prefix + text,
            api_key=api_key,
        )
        # Neteja possibles blocs markdown
        traduccio = traduccio.strip()
        if traduccio.startswith("```"):
            traduccio = re.sub(r'^```\w*\s*\n?', '', traduccio)
            traduccio = re.sub(r'\n?```\s*$', '', traduccio)

        temps_ms = int((time.perf_counter() - inici) * 1000)
        n_paraules = _compta_paraules(text)
        _stats["paraules_avui"] += n_paraules

        log.info("POST /tradueix-claude — %d ms, %d paraules", temps_ms, n_paraules)

        return {
            "translation": _neteja_resposta_claude(traduccio.strip()),
            "temps_ms": temps_ms,
            "motor": "claude-sonnet",
            "paraules": n_paraules,
        }

    except Exception as exc:
        log.exception("Error en la traducció amb Claude: %s", exc)
        raise HTTPException(
            status_code=500,
            detail=f"Error en la traducció amb Claude Sonnet: {exc}",
        )


# ─── Endpoint: POST /tradueix-document-claude (documents ES→VA amb Claude) ───

@app.post(
    "/tradueix-document-claude",
    summary="Tradueix un document .docx/.pptx del castellà al valencià amb Claude",
    tags=["Traducció"],
)
async def tradueix_document_claude(
    fitxer: UploadFile = File(..., description="Fitxer .docx o .pptx"),
    domini: str = Form(default="", description="Domini lingüístic per aplicar el glossari."),
) -> Response:
    """
    Tradueix un document .docx/.pptx del castellà al valencià estàndard
    universitari usant Claude Sonnet, preservant el format original.

    Procés en dues passades:
    1. Traducció de cada paràgraf
    2. Revisió profunda per Claude
    """
    api_key = _obte_api_key_anthropic()
    if not api_key:
        raise HTTPException(
            status_code=503,
            detail="Clau API d'Anthropic no configurada.",
        )

    nom = fitxer.filename or "document"
    extensio = Path(nom).suffix.lower()
    if extensio not in (".docx", ".pptx"):
        raise HTTPException(
            status_code=415,
            detail=f"Extensió '{extensio}' no admesa. Només .docx i .pptx.",
        )

    contingut = await fitxer.read()
    if len(contingut) > MAX_MIDA_FITXER:
        raise HTTPException(
            status_code=413,
            detail=f"Fitxer massa gran ({len(contingut) / 1_048_576:.1f} MB). Màx 150 MB.",
        )

    log.info(
        "POST /tradueix-document-claude — fitxer='%s' mida=%.1f KB",
        nom, len(contingut) / 1024,
    )

    inici = time.perf_counter()
    ext = extensio.lstrip('.')

    try:
        system_blocks, prefix = construeix_prompt_traduccio_es_va()

        # Si hi ha glossari del domini, afig-lo al prefix
        glossari_domini = carrega_glossari_com_diccionari(domini) if domini else {}
        if glossari_domini:
            log.info("Aplicant glossari '%s' (%d termes) a la traducció Claude",
                     domini, len(glossari_domini))
            termes_str = "\n".join(f"  {es} → {va}" for es, va in glossari_domini.items())
            prefix += (
                f"\n\nGLOSSARI OBLIGATORI del domini '{domini}':\n"
                f"Has d'usar EXACTAMENT aquests termes en la traducció:\n"
                f"{termes_str}\n\nTradueix el text següent:\n\n"
            )

        # Funció de traducció per lots (agrupats en blocs de ~800 paraules)
        _SEPARADOR_LOTS = '|||SEGMENT|||'
        _MAX_PARAULES_LOT = 800

        async def _tradueix_lots(textos):
            # Agrupa segments en lots per reduir crides a l'API
            lots = []          # cada lot = llista d'índexs dins de textos
            lot_actual = []
            paraules_lot = 0

            for i, text in enumerate(textos):
                if not text.strip() or len(text.strip()) < 4:
                    continue                    # es copiarà tal qual
                n_paraules = len(text.split())
                # Si afegir aquest segment supera el límit, tanca el lot
                if lot_actual and (paraules_lot + n_paraules) > _MAX_PARAULES_LOT:
                    lots.append(lot_actual)
                    lot_actual = []
                    paraules_lot = 0
                lot_actual.append(i)
                paraules_lot += n_paraules

            if lot_actual:
                lots.append(lot_actual)

            # Prepara resultats (per defecte = text original)
            resultats = list(textos)

            log.info("Traducció per lots: %d segments → %d lots", len(textos), len(lots))

            for num_lot, indexos in enumerate(lots, 1):
                if len(indexos) == 1:
                    # Lot d'un sol segment → crida directa
                    idx = indexos[0]
                    trad = await _crida_claude_amb_cache(
                        system_blocks=system_blocks,
                        missatge_usuari=prefix + textos[idx],
                        api_key=api_key,
                        max_tokens=4096,
                    )
                    resultats[idx] = _neteja_resposta_claude(trad.strip())
                else:
                    # Lot de múltiples segments → concatena amb separador
                    bloc = ('\n' + _SEPARADOR_LOTS + '\n').join(
                        textos[idx] for idx in indexos
                    )
                    instruccio = (
                        prefix
                        + "Tradueix cadascun dels segments següents. "
                        + f"Separa les traduccions amb «{_SEPARADOR_LOTS}» exactament com apareixen.\n\n"
                        + bloc
                    )
                    trad = await _crida_claude_amb_cache(
                        system_blocks=system_blocks,
                        missatge_usuari=instruccio,
                        api_key=api_key,
                        max_tokens=8192,
                    )
                    parts = trad.split(_SEPARADOR_LOTS)
                    # Assigna cada part al seu índex
                    for j, idx in enumerate(indexos):
                        if j < len(parts):
                            resultats[idx] = _neteja_resposta_claude(parts[j].strip())
                        # Si falten parts, manté l'original (ja copiat)

                    log.info("  Lot %d/%d: %d segments traduïts en 1 crida",
                             num_lot, len(lots), len(indexos))

            return resultats

        # Passada 1: traducció
        import asyncio
        loop = asyncio.get_event_loop()

        def _tradueix_sync(textos):
            future = asyncio.run_coroutine_threadsafe(_tradueix_lots(textos), loop)
            return future.result()

        resultat_bytes = await asyncio.to_thread(
            processa_document_traduccio,
            contingut, ext,
            _tradueix_sync,
        )

        # Passada 2: revisió per lots (sobre el document ja traduït)
        system_rev, _ = construeix_prompt_revisio("es_va")
        if ext == 'docx':
            editor_rev = DocxFormatPreservingEditor(resultat_bytes)
        else:
            editor_rev = PptxFormatPreservingEditor(resultat_bytes)

        try:
            paragrafs_rev = editor_rev.extrau_paragrafs()
            # Filtra paràgrafs amb prou text per revisar
            paragrafs_a_revisar = [p for p in paragrafs_rev if len(p['text'].strip()) >= 10]

            # Agrupa en lots de ~800 paraules
            lots_rev = []
            lot_actual_rev = []
            paraules_lot_rev = 0
            for p in paragrafs_a_revisar:
                n_p = len(p['text'].split())
                if lot_actual_rev and (paraules_lot_rev + n_p) > _MAX_PARAULES_LOT:
                    lots_rev.append(lot_actual_rev)
                    lot_actual_rev = []
                    paraules_lot_rev = 0
                lot_actual_rev.append(p)
                paraules_lot_rev += n_p
            if lot_actual_rev:
                lots_rev.append(lot_actual_rev)

            log.info("Revisió per lots: %d paràgrafs → %d lots", len(paragrafs_a_revisar), len(lots_rev))

            for num_lot, lot_paragrafs in enumerate(lots_rev, 1):
                if len(lot_paragrafs) == 1:
                    p = lot_paragrafs[0]
                    text_revisat = await _crida_claude_amb_cache(
                        system_blocks=system_rev,
                        missatge_usuari=f"Revisa aquest paràgraf traduït i corregeix si cal:\n\n{p['text']}",
                        api_key=api_key,
                        max_tokens=2048,
                    )
                    text_revisat = _neteja_resposta_claude(text_revisat.strip())
                    if text_revisat and text_revisat != p['text']:
                        editor_rev.substitueix_paragraf(p, text_revisat)
                else:
                    bloc_rev = ('\n' + _SEPARADOR_LOTS + '\n').join(
                        p['text'] for p in lot_paragrafs
                    )
                    instruccio_rev = (
                        "Revisa cadascun dels paràgrafs traduïts següents i corregeix si cal. "
                        f"Separa les revisions amb «{_SEPARADOR_LOTS}» exactament com apareixen.\n\n"
                        + bloc_rev
                    )
                    trad_rev = await _crida_claude_amb_cache(
                        system_blocks=system_rev,
                        missatge_usuari=instruccio_rev,
                        api_key=api_key,
                        max_tokens=8192,
                    )
                    parts_rev = trad_rev.split(_SEPARADOR_LOTS)
                    for j, p in enumerate(lot_paragrafs):
                        if j < len(parts_rev):
                            text_revisat = _neteja_resposta_claude(parts_rev[j].strip())
                            if text_revisat and text_revisat != p['text']:
                                editor_rev.substitueix_paragraf(p, text_revisat)

                    log.info("  Revisió lot %d/%d: %d paràgrafs revisats en 1 crida",
                             num_lot, len(lots_rev), len(lot_paragrafs))

            resultat_bytes = editor_rev.desa()
        finally:
            editor_rev.tanca()

        # Estableix la llengua predeterminada del document traduït (ES→VA = català)
        codi_llengua = "ca-ES"
        if ext == 'docx':
            editor_ll = DocxFormatPreservingEditor(resultat_bytes)
            editor_ll.estableix_llengua(codi_llengua)
            resultat_bytes = editor_ll.desa()
            editor_ll.tanca()
        elif ext == 'pptx':
            editor_ll = PptxFormatPreservingEditor(resultat_bytes)
            editor_ll.estableix_llengua(codi_llengua)
            resultat_bytes = editor_ll.desa()
            editor_ll.tanca()

        temps_ms = int((time.perf_counter() - inici) * 1000)
        nom_sortida = genera_nom_arxiu(nom, domini=domini, sufix="VAL")

        mime_types = {
            ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        }

        log.info("Document traduït amb Claude en %d ms → '%s'", temps_ms, nom_sortida)

        return Response(
            content=resultat_bytes,
            media_type=mime_types[extensio],
            headers={
                "Content-Disposition": f'attachment; filename="{nom_sortida}"',
                "Content-Length": str(len(resultat_bytes)),
                "X-Temps-Ms": str(temps_ms),
                "X-Motor": "claude-sonnet",
            },
        )

    except HTTPException:
        raise
    except Exception as exc:
        log.exception("Error traduint document amb Claude: %s", exc)
        raise HTTPException(
            status_code=500,
            detail=f"Error en la traducció del document amb Claude: {exc}",
        )


# ─── Endpoint: POST /corregeix-v2 (correcció millorada amb JSON estructurat) ─

class PeticioCorreccioV2(BaseModel):
    text: str = Field(
        ...,
        min_length=1,
        max_length=100_000,
        description="Text en valencià a corregir.",
    )
    usar_languagetool: bool = Field(default=True)
    usar_claude: bool = Field(default=True)
    domini: str = Field(default="", description="Domini lingüístic per aplicar el glossari.")


@app.post(
    "/corregeix-v2",
    summary="Correcció millorada amb JSON estructurat i estadístiques",
    tags=["Correcció"],
)
async def corregeix_v2(peticio: PeticioCorreccioV2):
    """
    Correcció millorada que retorna:
    - text_corregit
    - correccions en format JSON estructurat (num, paragraf, original, correccio, categoria, justificacio)
    - resum estadístic (total_errors, sint, morf, lex, orto, densitat, diagnostic)
    - correccions LanguageTool (si activat)

    Usa el prompt normatiu consolidat amb prompt caching.
    """
    import httpx
    import json as _json

    text_entrant = peticio.text.strip()
    if not text_entrant:
        raise HTTPException(status_code=422, detail="El camp 'text' no pot estar buit.")

    correccions_lt = []
    text_despres_lt = text_entrant

    # ── CAPA 1: LanguageTool ──
    if peticio.usar_languagetool:
        try:
            async with httpx.AsyncClient(timeout=30.0) as client:
                resp = await client.post(
                    "https://api.languagetool.org/v2/check",
                    data={"text": text_entrant, "language": "ca-ES"},
                    headers={"Accept": "application/json"},
                )
                resp.raise_for_status()
                dades_lt = resp.json()

            for match in dades_lt.get("matches", []):
                context = match.get("context", {})
                original = context.get("text", "")[
                    context.get("offset", 0):
                    context.get("offset", 0) + context.get("length", 0)
                ]
                suggerits = [r.get("value", "") for r in match.get("replacements", [])[:3]]
                correccions_lt.append({
                    "missatge": match.get("message", ""),
                    "offset": match.get("offset", 0),
                    "longitud": match.get("length", 0),
                    "original": original,
                    "suggerits": suggerits,
                    "regla_id": match.get("rule", {}).get("id", ""),
                })

            text_despres_lt = text_entrant
            for c in sorted(correccions_lt, key=lambda x: x["offset"], reverse=True):
                if c["suggerits"]:
                    ini = c["offset"]
                    fi = ini + c["longitud"]
                    text_despres_lt = text_despres_lt[:ini] + c["suggerits"][0] + text_despres_lt[fi:]

            log.info("LanguageTool — %d coincidències", len(correccions_lt))
        except Exception as exc:
            log.warning("Error LanguageTool: %s", exc)

    # ── CAPA 2: Claude Sonnet amb prompt caching ──
    correccions_claude = []
    text_corregit = text_despres_lt
    resum = {}

    if peticio.usar_claude:
        api_key = _obte_api_key_anthropic()
        if not api_key:
            raise HTTPException(
                status_code=503,
                detail="Clau API d'Anthropic no configurada.",
            )

        try:
            system_blocks, prefix = construeix_prompt_correccio()

            # Si hi ha glossari del domini, afig-lo al prefix
            glossari_domini = carrega_glossari_com_diccionari(peticio.domini) if peticio.domini else {}
            if glossari_domini:
                log.info("Aplicant glossari '%s' (%d termes) a la correcció de text",
                         peticio.domini, len(glossari_domini))
                termes_str = "\n".join(f"  {es} → {va}" for es, va in glossari_domini.items())
                prefix += (
                    f"\n\nGLOSSARI D'ESPECIALITAT ({peticio.domini}):\n"
                    f"Si trobes algun d'aquests termes, assegura't que la forma "
                    f"valenciana correcta és la indicada:\n{termes_str}\n\n"
                )

            resposta_text = await _crida_claude_amb_cache(
                system_blocks=system_blocks,
                missatge_usuari=prefix + text_despres_lt,
                api_key=api_key,
            )

            # Intentar parsejar JSON
            try:
                clean_json = resposta_text.strip()
                if clean_json.startswith("```"):
                    clean_json = re.sub(r'^```\w*\s*\n?', '', clean_json)
                    clean_json = re.sub(r'\n?```\s*$', '', clean_json)
                dades_claude = _json.loads(clean_json)

                text_corregit = dades_claude.get("text_corregit", text_despres_lt)
                correccions_claude = dades_claude.get("correccions", [])
                resum = dades_claude.get("resum", {})

            except _json.JSONDecodeError:
                log.warning("No s'ha pogut parsejar el JSON de Claude. Resposta crua guardada.")
                # Intentar extraure amb el format antic ---TEXT CORREGIT---
                m_text = re.search(
                    r'---TEXT CORREGIT---\s*(.*?)\s*---FI TEXT---',
                    resposta_text, re.DOTALL,
                )
                if m_text:
                    text_corregit = m_text.group(1).strip()

            log.info("Claude v2 — %d correccions", len(correccions_claude))

        except Exception as exc:
            log.exception("Error Claude v2: %s", exc)
            raise HTTPException(
                status_code=500,
                detail=f"Error en la correcció amb Claude: {exc}",
            )

    return {
        "text_original": text_entrant,
        "text_corregit": text_corregit,
        "correccions_lt": correccions_lt,
        "correccions_claude": correccions_claude,
        "resum": resum,
        "estat": "ok",
    }


# ─── Endpoint: POST /corregeix-document-v2 (doc + JSON + ressaltat) ──────────

@app.post(
    "/corregeix-document-v2",
    summary="Corregeix un document .docx/.pptx i retorna fitxer + correccions JSON",
    tags=["Correcció"],
)
async def corregeix_document_v2(
    fitxer: UploadFile = File(..., description="Fitxer .docx o .pptx"),
    domini: str = Form(default="", description="Domini lingüístic per aplicar el glossari."),
):
    """
    Corregeix un document preservant format:
    1. Extrau text per paràgrafs
    2. Envia a Claude per correcció (JSON estructurat)
    3. Aplica correccions amb ressaltat groc (docx)
    4. Retorna JSON amb el fitxer codificat en base64 + correccions + resum
    """
    import json as _json

    api_key = _obte_api_key_anthropic()
    if not api_key:
        raise HTTPException(status_code=503, detail="Clau API d'Anthropic no configurada.")

    nom = fitxer.filename or "document"
    extensio = Path(nom).suffix.lower()
    if extensio not in (".docx", ".pptx"):
        raise HTTPException(status_code=415, detail="Només .docx i .pptx.")

    contingut = await fitxer.read()
    if len(contingut) > MAX_MIDA_FITXER:
        raise HTTPException(status_code=413, detail="Fitxer massa gran.")

    log.info("POST /corregeix-document-v2 — '%s' %.1f KB", nom, len(contingut) / 1024)

    inici = time.perf_counter()

    try:
        ext = extensio.lstrip('.')
        if ext == 'docx':
            editor = DocxFormatPreservingEditor(contingut)
        else:
            editor = PptxFormatPreservingEditor(contingut)

        paragrafs = editor.extrau_paragrafs()
        text_complet = "\n\n".join(
            f"§{i+1}: {p['text']}" for i, p in enumerate(paragrafs) if p['text'].strip()
        )

        # Enviar tot el text a Claude per correcció
        system_blocks, prefix = construeix_prompt_correccio()

        # Si hi ha glossari del domini, afig-lo al prefix
        glossari_domini = carrega_glossari_com_diccionari(domini) if domini else {}
        if glossari_domini:
            log.info("Aplicant glossari '%s' (%d termes) a la correcció de document",
                     domini, len(glossari_domini))
            termes_str = "\n".join(f"  {es} → {va}" for es, va in glossari_domini.items())
            prefix += (
                f"\n\nGLOSSARI D'ESPECIALITAT ({domini}):\n"
                f"Si trobes algun d'aquests termes, assegura't que la forma "
                f"valenciana correcta és la indicada:\n{termes_str}\n\n"
            )

        resposta_text = await _crida_claude_amb_cache(
            system_blocks=system_blocks,
            missatge_usuari=prefix + text_complet,
            api_key=api_key,
        )

        # Parsejar resposta
        correccions = []
        resum = {}
        text_corregit_complet = ""

        try:
            clean_json = resposta_text.strip()
            if clean_json.startswith("```"):
                clean_json = re.sub(r'^```\w*\s*\n?', '', clean_json)
                clean_json = re.sub(r'\n?```\s*$', '', clean_json)
            dades = _json.loads(clean_json)
            correccions = dades.get("correccions", [])
            resum = dades.get("resum", {})
            text_corregit_complet = dades.get("text_corregit", "")
        except _json.JSONDecodeError:
            log.warning("JSON no parsejable en correcció de document.")

        # Aplicar correccions al document amb ressaltat
        if correccions:
            resultat_bytes = processa_document_correccio(
                contingut, ext, correccions
            )
        else:
            resultat_bytes = contingut

        editor.tanca()

        temps_ms = int((time.perf_counter() - inici) * 1000)
        nom_sortida = genera_nom_arxiu(nom, domini=domini, sufix="CORR_VAL")

        # Retornar tot en JSON
        return {
            "fitxer_base64": base64.b64encode(resultat_bytes).decode('utf-8'),
            "nom_fitxer": nom_sortida,
            "correccions": correccions,
            "resum": resum,
            "temps_ms": temps_ms,
            "estat": "ok",
        }

    except HTTPException:
        raise
    except Exception as exc:
        log.exception("Error corregint document v2: %s", exc)
        raise HTTPException(
            status_code=500,
            detail=f"Error en la correcció del document: {exc}",
        )


# ─── Endpoint: POST /tradueix-angles (EN↔VA) ────────────────────────────────

class PeticioTradueixAngles(BaseModel):
    text: str = Field(
        ...,
        min_length=1,
        max_length=100_000,
        description="Text a traduir.",
    )
    direccio: str = Field(
        default="en_va",
        description="Direcció: 'en_va' (anglés→valencià) o 'va_en' (valencià→anglés).",
    )


@app.post(
    "/tradueix-angles",
    summary="Tradueix entre anglés britànic i valencià",
    tags=["Traducció"],
)
async def tradueix_angles(peticio: PeticioTradueixAngles):
    """
    Traducció EN↔VA (anglés britànic ↔ valencià estàndard universitari)
    usant Claude Sonnet amb prompt normatiu i prompt caching.
    """
    api_key = _obte_api_key_anthropic()
    if not api_key:
        raise HTTPException(status_code=503, detail="Clau API d'Anthropic no configurada.")

    text = peticio.text.strip()
    if not text:
        raise HTTPException(status_code=422, detail="El camp 'text' no pot estar buit.")

    inici = time.perf_counter()

    try:
        system_blocks, prefix = construeix_prompt_traduccio_en_va(peticio.direccio)
        traduccio = await _crida_claude_amb_cache(
            system_blocks=system_blocks,
            missatge_usuari=prefix + text,
            api_key=api_key,
        )
        traduccio = traduccio.strip()
        if traduccio.startswith("```"):
            traduccio = re.sub(r'^```\w*\s*\n?', '', traduccio)
            traduccio = re.sub(r'\n?```\s*$', '', traduccio)

        temps_ms = int((time.perf_counter() - inici) * 1000)
        n_paraules = _compta_paraules(text)
        _stats["paraules_avui"] += n_paraules

        log.info(
            "POST /tradueix-angles [%s] — %d ms, %d paraules",
            peticio.direccio, temps_ms, n_paraules,
        )

        return {
            "translation": _neteja_resposta_claude(traduccio.strip()),
            "temps_ms": temps_ms,
            "motor": "claude-sonnet",
            "direccio": peticio.direccio,
            "paraules": n_paraules,
        }

    except Exception as exc:
        log.exception("Error traducció EN↔VA: %s", exc)
        raise HTTPException(
            status_code=500,
            detail=f"Error en la traducció anglés↔valencià: {exc}",
        )


# ─── Endpoint: POST /tradueix-document-angles (documents EN↔VA) ─────────────

@app.post(
    "/tradueix-document-angles",
    summary="Tradueix un document .docx/.pptx entre anglés i valencià",
    tags=["Traducció"],
)
async def tradueix_document_angles(
    fitxer: UploadFile = File(..., description="Fitxer .docx o .pptx"),
    direccio: str = Form(default="en_va", description="'en_va' o 'va_en'"),
) -> Response:
    """
    Tradueix un document entre anglés britànic i valencià estàndard universitari,
    preservant el format original.
    """
    api_key = _obte_api_key_anthropic()
    if not api_key:
        raise HTTPException(status_code=503, detail="Clau API d'Anthropic no configurada.")

    nom = fitxer.filename or "document"
    extensio = Path(nom).suffix.lower()
    if extensio not in (".docx", ".pptx"):
        raise HTTPException(status_code=415, detail="Només .docx i .pptx.")

    contingut = await fitxer.read()
    if len(contingut) > MAX_MIDA_FITXER:
        raise HTTPException(status_code=413, detail="Fitxer massa gran.")

    log.info(
        "POST /tradueix-document-angles [%s] — '%s' %.1f KB",
        direccio, nom, len(contingut) / 1024,
    )

    inici = time.perf_counter()
    ext = extensio.lstrip('.')

    try:
        system_blocks, prefix = construeix_prompt_traduccio_en_va(direccio)

        async def _tradueix_lots_angles(textos):
            resultats = []
            for text in textos:
                if not text.strip() or len(text.strip()) < 4:
                    resultats.append(text)
                    continue
                trad = await _crida_claude_amb_cache(
                    system_blocks=system_blocks,
                    missatge_usuari=prefix + text,
                    api_key=api_key,
                    max_tokens=4096,
                )
                resultats.append(_neteja_resposta_claude(trad.strip()))
            return resultats

        import asyncio
        loop = asyncio.get_event_loop()

        def _tradueix_sync(textos):
            future = asyncio.run_coroutine_threadsafe(_tradueix_lots_angles(textos), loop)
            return future.result()

        resultat_bytes = await asyncio.to_thread(
            processa_document_traduccio,
            contingut, ext,
            _tradueix_sync,
        )

        # Estableix la llengua predeterminada del document traduït
        codi_llengua = "ca-ES" if direccio == "en_va" else "en-GB"
        if ext == 'docx':
            editor_ll = DocxFormatPreservingEditor(resultat_bytes)
            editor_ll.estableix_llengua(codi_llengua)
            resultat_bytes = editor_ll.desa()
            editor_ll.tanca()
        elif ext == 'pptx':
            editor_ll = PptxFormatPreservingEditor(resultat_bytes)
            editor_ll.estableix_llengua(codi_llengua)
            resultat_bytes = editor_ll.desa()
            editor_ll.tanca()

        temps_ms = int((time.perf_counter() - inici) * 1000)
        sufix = "EN_VA" if direccio == "en_va" else "VA_EN"
        nom_sortida = genera_nom_arxiu(nom, sufix=sufix)

        mime_types = {
            ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        }

        log.info("Document EN↔VA traduït en %d ms → '%s'", temps_ms, nom_sortida)

        return Response(
            content=resultat_bytes,
            media_type=mime_types[extensio],
            headers={
                "Content-Disposition": f'attachment; filename="{nom_sortida}"',
                "Content-Length": str(len(resultat_bytes)),
                "X-Temps-Ms": str(temps_ms),
                "X-Motor": "claude-sonnet",
                "X-Direccio": direccio,
            },
        )

    except HTTPException:
        raise
    except Exception as exc:
        log.exception("Error traduint document EN↔VA: %s", exc)
        raise HTTPException(
            status_code=500,
            detail=f"Error en la traducció del document anglés↔valencià: {exc}",
        )


# ─── Gestió global d'errors ───────────────────────────────────────────────────

@app.exception_handler(404)
async def no_trobat(request, exc):
    from fastapi.responses import JSONResponse
    return JSONResponse(
        status_code=404,
        content={"error": "Recurs no trobat", "detall": str(exc.detail)},
    )


@app.exception_handler(500)
async def error_intern(request, exc):
    from fastapi.responses import JSONResponse
    log.error("Error intern no controlat: %s", exc)
    return JSONResponse(
        status_code=500,
        content={"error": "Error intern del servidor", "detall": "Consulta els registres de l'API."},
    )


# ─── Inici del servidor (execució directa) ───────────────────────────────────

if __name__ == "__main__":
    import uvicorn
    log.info("Iniciant servidor Uvicorn en mode directe...")
    uvicorn.run(app, host="0.0.0.0", port=8000, reload=True)