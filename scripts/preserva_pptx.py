# -*- coding: utf-8 -*-
"""
preserva_pptx.py
----------------
Tradueix presentacions .pptx preservant el format original run a run.

Correccions aplicades:
  - neteja_traduccio(): versió robusta (equivalent a preserva_docx.py).
  - substitueix_text_para_pptx(): nova estratègia que BUIDA el text de
    tots els runs (sense eliminar-los del XML) i restaura el format
    explícitament.
  - _estableix_llengua_catalana_pptx(): marca lang="ca-ES" al a:rPr.
"""

import io
import re
import unicodedata

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


# ── Helpers globals ────────────────────────────────────────────────────────────

# Detecta URLs per protegir-les abans de traduir (el model les trenca afegint espais)
URL_PATTERN = re.compile(
    r'https?://[^\s<>"\']+|www\.[^\s<>"\']+',
    re.IGNORECASE
)


def protegeix_urls(text: str) -> tuple[str, dict]:
    """Substitueix les URLs per marcadors i retorna el text modificat i el mapa.

    Usa el format URLTOKEN{n}XEND per evitar que el model el processe com
    a paraula normal (els guions baixos i els dobles guions sí que els modifica).
    """
    urls: dict = {}

    def reemplaca(m):
        key = f'URLTOKEN{len(urls)}XEND'
        urls[key] = m.group(0)
        return key

    text_protegit = URL_PATTERN.sub(reemplaca, text)
    return text_protegit, urls


def restaura_urls(text: str, urls: dict) -> str:
    """Restaura les URLs originals als marcadors."""
    for key, url in urls.items():
        text = text.replace(key, url)
    return text


def te_format_explicit(fmt: dict) -> bool:
    """Retorna True si el run tenia almenys un atribut de format definit explícitament."""
    return any([
        fmt['bold'] is not None,
        fmt['italic'] is not None,
        fmt['underline'] is not None,
        fmt['font_name'] is not None,
        fmt['font_size'] is not None,
        fmt['color'] is not None,
    ])


def distribueix_text_entre_runs(paraules: list, formats: list) -> list:
    """
    Distribueix les paraules proporcionalment entre els runs assegurant
    que no es perden espais entre paraules per arrodoniment.
    """
    total_original = sum(f['text_len'] for f in formats)
    total_paraules = len(paraules)
    fragments = []
    paraula_idx = 0

    for i, fmt in enumerate(formats):
        if i == len(formats) - 1:
            fragment = ' '.join(paraules[paraula_idx:])
        else:
            if total_original > 0:
                proporcio = fmt['text_len'] / total_original
                n_paraules = max(1, round(total_paraules * proporcio))
            else:
                n_paraules = 1
            fragment = ' '.join(paraules[paraula_idx:paraula_idx + n_paraules])
            paraula_idx = min(paraula_idx + n_paraules, total_paraules)
        fragments.append(fragment)

    return fragments


def corregeix_paraules_enganxades(text: str) -> str:
    """
    Detecta paraules enganxades (minúscula seguida directament de majúscula
    o lletres seguides de números sense espai) i afegeix l'espai que falta.
    Aplica NOMÉS als casos clars per evitar falsos positius.
    """
    # Cas: lletra minúscula enganxada amb majúscula (ex: "serveiEl" → "servei El")
    text = re.sub(r'([a-záéíóúàèìòùïüç])([A-ZÁÉÍÓÚÀÈÌÒÙÏÜÇ])', r'\1 \2', text)
    # Cas: número enganxat amb lletra (ex: "2023la" → "2023 la")
    text = re.sub(r'(\d)([A-ZÁÉÍÓÚÀÈÌÒÙÏÜÇa-záéíóúàèìòùïüç])', r'\1 \2', text)
    # Cas: lletra enganxada amb número (ex: "any2023" → "any 2023")
    # EXCEPCIÓ: no se separen URLs ni codis tècnics (ja protegits amb marcadors)
    text = re.sub(r'([a-záéíóúàèìòùïüç])(\d)', r'\1 \2', text)
    return text


def corregeix_apostrofacions(text: str) -> str:
    """
    Corregeix els espais incorrectes en apostrofacions catalanes/valencianes.
    Casos: l' , d' , m' , t' , s' , n' , c' , j' , qu'
    Exemples: "d' exemple" → "d'exemple", "l' ús" → "l'ús"
    """
    # Elimina espai entre apòstrof i paraula següent
    text = re.sub(r"([dlmtsncjq]u?)'(\s+)([^\s])", r"\1'\3", text, flags=re.IGNORECASE)

    # "de " + vocal/h → "d'" (prepositional elision)
    # NOTA: pot produir falsos positius en noms propis; descomenta si cal
    text = re.sub(
        r'\bde\s+([aeiouàèéíïòóúüh])',
        lambda m: "d'" + m.group(1),
        text,
        flags=re.IGNORECASE,
    )

    # Normalitza espais múltiples que puguen haver quedat
    text = re.sub(r' {2,}', ' ', text)

    return text.strip()


def neteja_traduccio(text: str) -> str:
    """Elimina tots els artefactes Markdown i caràcters espuris del model."""
    if not text:
        return text

    # Elimina negreta i cursiva Markdown (**text**, *text*, __text__)
    text = re.sub(r'\*{1,3}(.*?)\*{1,3}', r'\1', text, flags=re.DOTALL)
    text = re.sub(r'_{1,2}(.*?)_{1,2}', r'\1', text, flags=re.DOTALL)

    # Elimina capçaleres Markdown
    text = re.sub(r'^#{1,6}\s+', '', text, flags=re.MULTILINE)

    # Elimina guions de llista al principi de línia
    text = re.sub(r'^[\-\–\—\•]\s+', '', text, flags=re.MULTILINE)

    # Elimina numeració al principi de línia AMB puntuació (1. text, 1) text)
    text = re.sub(r'^\d+[\.\)\-]\s+', '', text, flags=re.MULTILINE)

    # Elimina número sol al principi de línia SENSE puntuació (el cas "1 Segons")
    text = re.sub(r'^\d+\s+(?=[A-ZÁÉÍÓÚÀÈÌÒÙÏÜÇ·])', '', text, flags=re.MULTILINE)

    # Elimina número sol al principi de text (sense salts de línia)
    text = re.sub(r'^\d+\s+', '', text)

    # Elimina asteriscos solts que hagin quedat
    text = re.sub(r'\*+', '', text)

    # Normalitza espais de no-ruptura i caràcters de control Unicode
    text = re.sub(r'[\u00a0\u200b\u200c\u200d\ufeff]', ' ', text)
    text = re.sub(r' {2,}', ' ', text)

    # Elimina espais davant de puntuació
    text = re.sub(r' ([.,;:!?»\)\]])', r'\1', text)

    # Elimina línies buides múltiples
    text = re.sub(r'\n{3,}', '\n\n', text)

    return text.strip()


def _estableix_llengua_catalana_pptx(run) -> None:
    """
    Estableix l'atribut lang="ca-ES" al rPr (a:rPr) del run DrawingML.
    Sobreescriu qualsevol valor de llengua preexistent.
    """
    try:
        rPr = run._r.get_or_add_rPr()
        rPr.set('lang', 'ca-ES')
    except Exception:
        pass


def _aplica_format_pptx(run, fmt: dict) -> None:
    """
    Aplica el format guardat a un run PPTX.
    Escriu NOMÉS els valors explícits (no None) per evitar sobreescriure
    format heretat del tema o del marcador de posició.
    """
    if fmt['bold'] is not None:
        run.font.bold = fmt['bold']
    if fmt['italic'] is not None:
        run.font.italic = fmt['italic']
    if fmt['underline'] is not None:
        run.font.underline = fmt['underline']
    if fmt['font_name']:
        run.font.name = fmt['font_name']
    if fmt['font_size']:
        run.font.size = fmt['font_size']
    if fmt['color']:
        try:
            run.font.color.rgb = fmt['color']
        except Exception:
            pass


def substitueix_text_para_pptx(para, text_traduit: str) -> None:
    """
    Substitueix el text del paràgraf PPTX distribuint el text traduït
    proporcionalment entre els runs originals i preservant el format de cada un.

    Estratègia:
      1. Captura el format de TOTS els runs.
      2. Distribueix el text traduït amb distribueix_text_entre_runs().
      3. Aplica _aplica_format_pptx() NOMÉS als runs que tenien format explícit
         (te_format_explicit) per evitar afegir negreta/cursiva fantasma.
      4. Estableix la llengua ca-ES a tots els runs modificats.
    """
    if not para.runs:
        return

    # ── Captura el format de TOTS els runs ────────────────────────────────────
    formats = []
    for run in para.runs:
        fmt = {
            'bold':      run.font.bold,
            'italic':    run.font.italic,
            'underline': run.font.underline,
            'font_name': run.font.name,
            'font_size': run.font.size,
        }
        try:
            fmt['color'] = run.font.color.rgb
        except Exception:
            fmt['color'] = None
        fmt['text_len'] = len(run.text)
        formats.append(fmt)

    total_original = sum(f['text_len'] for f in formats)

    # ── Cas simple: un sol run o tots els runs buits ───────────────────────────
    if total_original == 0 or len(para.runs) == 1:
        para.runs[0].text = text_traduit
        if te_format_explicit(formats[0]):
            _aplica_format_pptx(para.runs[0], formats[0])
        _estableix_llengua_catalana_pptx(para.runs[0])
        return

    # ── Buida el text de TOTS els runs (preserva l'estructura XML) ────────────
    for run in para.runs:
        run.text = ''

    # ── Distribueix i aplica ───────────────────────────────────────────────────
    paraules = text_traduit.split()
    fragments = distribueix_text_entre_runs(paraules, formats)

    for run, fmt, fragment in zip(para.runs, formats, fragments):
        run.text = fragment
        if te_format_explicit(fmt):
            _aplica_format_pptx(run, fmt)
        _estableix_llengua_catalana_pptx(run)


# ── Classe principal ───────────────────────────────────────────────────────────

class PreservadorPptx:
    def __init__(self, traductor):
        self.traductor = traductor

    def tradueix_segment(self, text: str) -> str:
        """
        Tradueix un segment de text aplicant tota la canonada de processament:
          1. Retorna sense canvis si el text és buit o és únicament una URL.
          2. Protegeix les URLs amb marcadors URLTOKEN{n}XEND.
          3. Tradueix amb el model.
          4. Neteja artefactes Markdown (neteja_traduccio).
          5. Restaura les URLs originals.
          6. Corregeix paraules enganxades (corregeix_paraules_enganxades).
          7. Corregeix apostrofacions incorrectes (corregeix_apostrofacions).
        """
        text = text.strip()
        if not text:
            return text
        # Si el segment és únicament una URL, no el tradueixis
        if URL_PATTERN.fullmatch(text):
            return text
        # Protegeix les URLs dins del text
        text_protegit, urls = protegeix_urls(text)
        # Traducció + neteja
        traduit_raw = self.traductor(text_protegit)
        traduit = neteja_traduccio(traduit_raw)
        # Restaura les URLs originals
        traduit = restaura_urls(traduit, urls)
        # Postprocessament
        traduit = corregeix_paraules_enganxades(traduit)
        traduit = corregeix_apostrofacions(traduit)
        return traduit

    def tradueix_text_frame(self, tf):
        """
        Tradueix cada paràgraf del marc de text delegant a tradueix_segment i
        substitueix_text_para_pptx.
        """
        for para in tf.paragraphs:
            text = para.text.strip()
            if not text or not para.runs:
                continue
            traduit = self.tradueix_segment(text)
            substitueix_text_para_pptx(para, traduit)

    def tradueix_shape(self, shape):
        try:
            if shape.has_text_frame:
                self.tradueix_text_frame(shape.text_frame)
            if shape.has_table:
                for fila in shape.table.rows:
                    for cel in fila.cells:
                        self.tradueix_text_frame(cel.text_frame)
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                for s in shape.shapes:
                    self.tradueix_shape(s)
        except Exception:
            pass

    def tradueix_document(self, entrada, sortida=None):
        if isinstance(entrada, bytes):
            entrada = io.BytesIO(entrada)
        prs = Presentation(entrada)
        for diap in prs.slides:
            for shape in diap.shapes:
                self.tradueix_shape(shape)
            if diap.has_notes_slide:
                try:
                    self.tradueix_text_frame(
                        diap.notes_slide.notes_text_frame)
                except Exception:
                    pass
        buf = io.BytesIO()
        prs.save(buf)
        if sortida:
            prs.save(sortida)
        return buf.getvalue()
