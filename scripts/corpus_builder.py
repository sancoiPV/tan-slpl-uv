# -*- coding: utf-8 -*-
"""
corpus_builder.py
-----------------
Processa fitxers Word i PowerPoint paral·lels (_CAS/_VAL) i genera
un corpus d'entrenament i validació per al motor TAN SLPL·UV.

Ús:
  python corpus_builder.py --input <carpeta_corpus> --output <carpeta_sortida>
  python corpus_builder.py --input <carpeta> --output <sortida> --max-delta 20

Arguments:
  --input      Carpeta amb fitxers _CAS.docx, _VAL.docx, _CAS.pptx, _VAL.pptx
  --output     Carpeta on es desaran els fitxers resultants
  --max-delta  Exclou parells de documents amb diferència de paràgrafs superior
               a aquest valor (defecte: 999, és a dir sense límit)

Fitxers generats a --output:
  train.es        Frases castellà d'entrenament (90%)
  train.ca        Frases valencià d'entrenament (90%)
  val.es          Frases castellà de validació (10%)
  val.ca          Frases valencià de validació (10%)
  corpus_net.tsv  Tots els parells vàlids en format TSV per a revisió
  informe.txt     Resum del processament + secció de desalineaments

Dependències: python-docx, python-pptx, nltk, bertalign (opcional, recomanat)

Alineació semàntica (Bertalign):
  Quan el nombre de paràgrafs difereix entre _CAS i _VAL, s'usa Bertalign
  (model LaBSE) per trobar la correspondència semàntica correcta.
  Si Bertalign no és disponible, s'aplica fallback a truncació posicional.
  Instal·lació: pip install bertalign
"""

import argparse
import logging
import random
import re
import sys
from pathlib import Path

# ── Dependències opcionals amb missatge d'error clar ──────────────────────────
try:
    from docx import Document as DocxDocument
except ImportError:
    sys.exit('[ERROR] python-docx no instal·lat. Executa: pip install python-docx')

try:
    from pptx import Presentation
except ImportError:
    sys.exit('[ERROR] python-pptx no instal·lat. Executa: pip install python-pptx')

try:
    import nltk
    try:
        nltk.data.find('tokenizers/punkt_tab')
    except LookupError:
        nltk.download('punkt_tab', quiet=True)
    try:
        nltk.data.find('tokenizers/punkt')
    except LookupError:
        nltk.download('punkt', quiet=True)
    from nltk.tokenize import sent_tokenize
    _NLTK_OK = True
except ImportError:
    _NLTK_OK = False

# ── Configuració del logging ──────────────────────────────────────────────────
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s  %(levelname)-7s  %(message)s',
    datefmt='%H:%M:%S',
)
log = logging.getLogger(__name__)


# ═══════════════════════════════════════════════════════════════════════════════
# Extracció de text
# ═══════════════════════════════════════════════════════════════════════════════

def extrau_paragrafs_docx(path: Path) -> list[str]:
    """Retorna la llista de paràgrafs no buits d'un fitxer .docx."""
    doc = DocxDocument(path)
    paragrafs = []
    # Cos principal
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            paragrafs.append(text)
    # Taules
    for taula in doc.tables:
        for fila in taula.rows:
            for cel in fila.cells:
                for para in cel.paragraphs:
                    text = para.text.strip()
                    if text:
                        paragrafs.append(text)
    return paragrafs


def extrau_paragrafs_pptx(path: Path) -> list[str]:
    """Retorna la llista de paràgrafs no buits d'un fitxer .pptx."""
    prs = Presentation(path)
    paragrafs = []
    for diapositiva in prs.slides:
        for forma in diapositiva.shapes:
            if not forma.has_text_frame:
                continue
            for para in forma.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    paragrafs.append(text)
    return paragrafs


def segmenta_frases(text: str, llengua: str = 'spanish') -> list[str]:
    """
    Segmenta un text en frases individuals.
    Usa nltk si disponible; si no, divideix per punts/salts de línia.
    """
    if _NLTK_OK:
        try:
            frases = sent_tokenize(text, language=llengua)
            return [f.strip() for f in frases if f.strip()]
        except Exception:
            pass
    # Fallback: divideix per punts finals i salts de línia
    frases = re.split(r'(?<=[.!?])\s+|\n+', text)
    return [f.strip() for f in frases if f.strip()]


# ═══════════════════════════════════════════════════════════════════════════════
# Alineació semàntica amb Bertalign
# ═══════════════════════════════════════════════════════════════════════════════

def alinea_amb_bertalign(text_cas: str, text_val: str) -> list[tuple[str, str]]:
    """
    Alinea dos textos paral·lels usant Bertalign (embeddings semàntics LaBSE).
    Gestiona correctament fusió i divisió de frases entre CAS i VAL
    (alineacions 1-a-molts, molts-a-1 i molts-a-molts).

    Retorna una llista de parells (text_castellà, text_valencià) alineats.
    Si Bertalign no està disponible o falla, retorna llista buida
    perquè la funció cridadora aplique el fallback per posició.
    """
    try:
        from bertalign import Bertalign  # type: ignore

        aligner = Bertalign(text_cas, text_val, src_lang='es', tgt_lang='ca')
        aligner.align_sents()

        parells = []
        for src_idx, tgt_idx in aligner.result:
            src_sents = [aligner.src_sents[i] for i in src_idx]
            tgt_sents = [aligner.tgt_sents[i] for i in tgt_idx]

            src_text = ' '.join(src_sents).strip()
            tgt_text = ' '.join(tgt_sents).strip()

            if src_text and tgt_text:
                parells.append((src_text, tgt_text))

        log.debug('Bertalign: %d parells alineats', len(parells))
        return parells

    except ImportError:
        log.warning('Bertalign no disponible — usa alineació per posició com a fallback')
        return []
    except Exception as e:
        log.warning('Error amb Bertalign: %r — usa alineació per posició com a fallback', e)
        return []


# ═══════════════════════════════════════════════════════════════════════════════
# Neteja i validació de parells
# ═══════════════════════════════════════════════════════════════════════════════

def similitud_bigrames(text1: str, text2: str) -> float:
    """
    Calcula la similitud entre dos textos basada en bigrames de caràcters.
    Retorna un valor entre 0 (cap similitud) i 1 (idèntics).
    Útil per detectar parells ES/CA mal alineats: una traducció real
    compartirà noms propis, números, sigles i paraules comunes.
    """
    def bigrames(text):
        text = text.lower().strip()
        return set(text[i:i+2] for i in range(len(text) - 1))

    bg1 = bigrames(text1)
    bg2 = bigrames(text2)

    if not bg1 or not bg2:
        return 0.0

    interseccio = len(bg1 & bg2)
    unio = len(bg1 | bg2)
    return interseccio / unio if unio > 0 else 0.0


def es_valid(
    src: str,
    tgt: str,
    min_tok: int = 3,
    max_tok: int = 200,
    min_similitud: float = 0.08,
) -> bool:
    """
    Comprova si un parell és vàlid per al corpus.
    Filtra frases massa curtes, massa llargues, idèntiques, poc alfabètiques
    o amb baixa similitud de bigrames (probable mal alineament ES/CA).
    """
    s_tok = src.split()
    t_tok = tgt.split()
    if not (min_tok <= len(s_tok) <= max_tok):
        return False
    if not (min_tok <= len(t_tok) <= max_tok):
        return False
    if src.strip() == tgt.strip():
        return False
    # Mínim 40% de caràcters alfabètics (filtra taules de números, codis, etc.)
    if sum(c.isalpha() for c in src) / max(len(src), 1) < 0.4:
        return False
    if sum(c.isalpha() for c in tgt) / max(len(tgt), 1) < 0.4:
        return False
    # Filtre de similitud: descarta parells probablement mal alineats.
    # Una traducció ES→CA real sempre compartirà almenys un 8% de bigrames
    # (noms propis, números, sigles, paraules cognades).
    if similitud_bigrames(src, tgt) < min_similitud:
        return False
    return True


# ═══════════════════════════════════════════════════════════════════════════════
# Processament de parells de fitxers
# ═══════════════════════════════════════════════════════════════════════════════

def extrau_paragrafs_slide(slide) -> list[str]:
    """
    Extreu els paràgrafs de text d'una diapositiva com una llista.
    Recorre totes les formes de text i les taules.
    """
    fragments = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    fragments.append(text)
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    text = cell.text.strip()
                    if text:
                        fragments.append(text)
    return fragments


def processa_parell_docx(
    path_cas: Path,
    path_val: Path,
    desalineaments: list,
    max_delta: int = 999,
) -> list[tuple[str, str]]:
    """
    Processa un parell de fitxers DOCX.

    Estratègia d'alineació:
    - Delta == 0 (paràgrafs iguals): alineació posicional directa (ràpida i precisa).
    - Delta > 0  (paràgrafs desiguals): Bertalign (embeddings semàntics LaBSE)
      per trobar la correspondència correcta fins i tot quan el traductor ha
      fusionat o dividit paràgrafs. Si Bertalign no és disponible, fallback
      a truncació posicional.
    """
    try:
        doc_cas = DocxDocument(path_cas)
        doc_val = DocxDocument(path_val)
    except Exception as e:
        log.warning('Error obrint %s: %s', path_cas.name, e)
        return []

    def extrau_paragrafs(doc):
        paras = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                paras.append(text)
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    text = cell.text.strip()
                    if text and text not in paras:
                        paras.append(text)
        return paras

    paras_cas = extrau_paragrafs(doc_cas)
    paras_val = extrau_paragrafs(doc_val)

    n_cas = len(paras_cas)
    n_val = len(paras_val)
    delta = abs(n_cas - n_val)

    if delta > 0:
        desalineaments.append({
            'tipus':           'DOCX',
            'fitxer':          path_cas.name,
            'paras_cas':       n_cas,
            'paras_val':       n_val,
            'delta':           delta,
            'parells_extrets': min(n_cas, n_val),
            'exclòs':          delta > max_delta,
        })
        if delta > max_delta:
            log.warning(
                'EXCLÒS Δ=%d > max_delta=%d: %s', delta, max_delta, path_cas.name,
            )
            return []
        log.warning(
            'Paràgrafs diferent (%d vs %d): %s — intentant Bertalign',
            n_cas, n_val, path_cas.name,
        )
        # Intenta alineació semàntica amb Bertalign
        text_cas = '\n'.join(paras_cas)
        text_val = '\n'.join(paras_val)
        parells = alinea_amb_bertalign(text_cas, text_val)
        if parells:
            log.info('  Bertalign: %d parells alineats (Δ=%d)', len(parells), delta)
            return parells
        # Fallback: truncació posicional
        log.warning('  Fallback posicional: truncant a %d', min(n_cas, n_val))
        n = min(n_cas, n_val)
        return list(zip(paras_cas[:n], paras_val[:n]))

    # Delta == 0: alineació posicional directa (cas habitual, ràpid)
    return list(zip(paras_cas, paras_val))


def processa_parell_pptx(
    path_cas: Path,
    path_val: Path,
    desalineaments: list,
    max_delta: int = 999,
) -> list[tuple[str, str]]:
    """
    Processa un parell de fitxers PPTX amb estratègia per diapositiva:

    Per a cada parell de diapositives N:
      - Paràgrafs iguals → alineació posicional directa (cas habitual).
      - Paràgrafs desiguals → Bertalign (embeddings semàntics) per trobar
        la correspondència correcta. Si Bertalign no és disponible, fallback
        a un sol parell concatenant tot el text de la diapositiva.
    """
    try:
        prs_cas = Presentation(path_cas)
        prs_val = Presentation(path_val)
    except Exception as e:
        log.warning('Error obrint %s: %s', path_cas.name, e)
        return []

    n_slides_cas = len(prs_cas.slides)
    n_slides_val = len(prs_val.slides)
    delta = abs(n_slides_cas - n_slides_val)

    if delta > 0:
        desalineaments.append({
            'tipus':           'PPTX',
            'fitxer':          path_cas.name,
            'paras_cas':       n_slides_cas,
            'paras_val':       n_slides_val,
            'delta':           delta,
            'parells_extrets': 0,
            'exclòs':          delta > max_delta,
        })
        if delta > max_delta:
            log.warning(
                'EXCLÒS Δ=%d > max_delta=%d: %s', delta, max_delta, path_cas.name,
            )
            return []
        log.warning(
            'Diapositives diferent (%d vs %d): %s — truncant a %d',
            n_slides_cas, n_slides_val, path_cas.name, min(n_slides_cas, n_slides_val),
        )

    parells = []
    n_slides = min(n_slides_cas, n_slides_val)

    for i in range(n_slides):
        paras_cas = extrau_paragrafs_slide(prs_cas.slides[i])
        paras_val = extrau_paragrafs_slide(prs_val.slides[i])

        if not paras_cas or not paras_val:
            continue

        if len(paras_cas) == len(paras_val):
            # Alineació posicional directa dins la diapositiva (cas habitual)
            for s, t in zip(paras_cas, paras_val):
                s, t = s.strip(), t.strip()
                if s and t:
                    parells.append((s, t))
        else:
            # Paràgrafs desequilibrats dins la diapositiva: prova Bertalign
            text_cas = '\n'.join(paras_cas)
            text_val = '\n'.join(paras_val)
            slide_parells = alinea_amb_bertalign(text_cas, text_val)
            if slide_parells:
                parells.extend(slide_parells)
            else:
                # Fallback: un sol parell amb tot el text de la diapositiva
                text_cas_complet = ' '.join(paras_cas).strip()
                text_val_complet = ' '.join(paras_val).strip()
                if text_cas_complet and text_val_complet:
                    parells.append((text_cas_complet, text_val_complet))

    return parells


# ═══════════════════════════════════════════════════════════════════════════════
# Funció principal
# ═══════════════════════════════════════════════════════════════════════════════

def main():
    parser = argparse.ArgumentParser(
        description='Genera corpus TAN des de fitxers _CAS/_VAL (docx/pptx)',
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog=__doc__,
    )
    parser.add_argument(
        '--input', required=True,
        help='Carpeta amb fitxers _CAS.docx/_VAL.docx i _CAS.pptx/_VAL.pptx',
    )
    parser.add_argument(
        '--output', required=True,
        help='Carpeta on es desaran els fitxers resultants',
    )
    parser.add_argument(
        '--max-delta', type=int, default=999,
        dest='max_delta',
        help='Exclou parells de documents amb diferència de paràgrafs superior '
             'a aquest valor (defecte: 999, sense límit pràctic)',
    )
    parser.add_argument(
        '--min-similitud', type=float, default=0.08,
        dest='min_similitud',
        help='Similitud mínima de bigrames de caràcters entre ES i CA '
             '(defecte: 0.08). Valors entre 0.0 (desactivat) i 1.0.',
    )
    args = parser.parse_args()

    input_dir  = Path(args.input)
    output_dir = Path(args.output)
    max_delta  = args.max_delta

    if not input_dir.exists():
        sys.exit(f'[ERROR] La carpeta d\'entrada no existeix: {input_dir}')

    output_dir.mkdir(parents=True, exist_ok=True)
    # Comprova disponibilitat de Bertalign
    try:
        from bertalign import Bertalign as _bt  # type: ignore  # noqa: F401
        _bertalign_ok = True
        log.info('Bertalign:        disponible (alineació semàntica LaBSE)')
    except ImportError:
        _bertalign_ok = False
        log.warning('Bertalign:        NO disponible — fallback posicional per a desalineaments')

    log.info('Carpeta entrada:  %s', input_dir)
    log.info('Carpeta sortida:  %s', output_dir)
    log.info('Max delta:        %d', max_delta)
    log.info('Min similitud:    %.2f', args.min_similitud)

    # ── Cerca tots els fitxers _CAS ────────────────────────────────────────────
    fitxers_cas_docx = sorted(input_dir.glob('**/*_CAS.docx'))
    fitxers_cas_pptx = sorted(input_dir.glob('**/*_CAS.pptx'))

    total_fitxers    = 0
    total_extrets    = 0
    total_eliminats  = 0
    total_exclosos   = 0       # documents exclosos per --max-delta
    registre_fitxers = []
    desalineaments   = []      # acumula tots els desalineaments detectats

    all_pairs = []

    # ── Processa fitxers DOCX ──────────────────────────────────────────────────
    for path_cas in fitxers_cas_docx:
        path_val = path_cas.parent / path_cas.name.replace('_CAS', '_VAL')
        if not path_val.exists():
            log.warning('No s\'ha trobat el fitxer VAL per a: %s', path_cas.name)
            continue
        log.info('DOCX: %s + %s', path_cas.name, path_val.name)
        try:
            parells = processa_parell_docx(
                path_cas, path_val, desalineaments, max_delta
            )
            if parells is None:
                # Exclòs per --max-delta
                total_fitxers += 2
                total_exclosos += 1
                registre_fitxers.append(
                    f'DOCX  {path_cas.name} + {path_val.name}'
                    f'  →  EXCLÒS (Δ>{max_delta})'
                )
                continue
            all_pairs.extend(parells)
            log.info('  → %d parells extrets', len(parells))
            total_fitxers += 2
            total_extrets += len(parells)
            registre_fitxers.append(
                f'DOCX  {path_cas.name} + {path_val.name}  →  {len(parells)} parells'
            )
        except Exception as exc:
            log.error('Error processant %s: %s', path_cas.name, exc)

    # ── Processa fitxers PPTX ──────────────────────────────────────────────────
    for path_cas in fitxers_cas_pptx:
        path_val = path_cas.parent / path_cas.name.replace('_CAS', '_VAL')
        if not path_val.exists():
            log.warning('No s\'ha trobat el fitxer VAL per a: %s', path_cas.name)
            continue
        log.info('PPTX: %s + %s', path_cas.name, path_val.name)
        try:
            parells = processa_parell_pptx(
                path_cas, path_val, desalineaments, max_delta
            )
            if parells is None:
                total_fitxers += 2
                total_exclosos += 1
                registre_fitxers.append(
                    f'PPTX  {path_cas.name} + {path_val.name}'
                    f'  →  EXCLÒS (Δ>{max_delta})'
                )
                continue
            all_pairs.extend(parells)
            log.info('  → %d parells extrets', len(parells))
            total_fitxers += 2
            total_extrets += len(parells)
            registre_fitxers.append(
                f'PPTX  {path_cas.name} + {path_val.name}  →  {len(parells)} parells'
            )
        except Exception as exc:
            log.error('Error processant %s: %s', path_cas.name, exc)

    if total_fitxers == 0:
        log.warning('No s\'ha trobat cap parell _CAS/_VAL a: %s', input_dir)

    # ── Neteja i deduplicació ──────────────────────────────────────────────────
    clean_pairs = []
    seen = set()
    eliminats_similitud = 0
    for src, tgt in all_pairs:
        src, tgt = src.strip(), tgt.strip()
        # Comptabilitza específicament els eliminats pel filtre de similitud:
        # passa tots els altres filtres però falla el de bigrames.
        if (
            es_valid(src, tgt, min_similitud=0.0)
            and similitud_bigrames(src, tgt) < args.min_similitud
        ):
            eliminats_similitud += 1
        if not es_valid(src, tgt, min_similitud=args.min_similitud):
            total_eliminats += 1
            continue
        clau = f'{src.lower()}|||{tgt.lower()}'
        if clau in seen:
            total_eliminats += 1
            continue
        seen.add(clau)
        clean_pairs.append((src, tgt))

    total_valid = len(clean_pairs)

    # ── Divisió train/val ──────────────────────────────────────────────────────
    random.seed(42)
    random.shuffle(clean_pairs)
    tall = int(len(clean_pairs) * 0.9)
    train_pairs = clean_pairs[:tall]
    val_pairs   = clean_pairs[tall:]

    # ── Guardat dels fitxers de sortida ───────────────────────────────────────
    def escriu_parells(pairs: list, prefix: str) -> None:
        path_es = output_dir / f'{prefix}.es'
        path_ca = output_dir / f'{prefix}.ca'
        with open(path_es, 'w', encoding='utf-8') as fe, \
             open(path_ca, 'w', encoding='utf-8') as fc:
            for src, tgt in pairs:
                fe.write(src + '\n')
                fc.write(tgt + '\n')

    escriu_parells(train_pairs, 'train')
    escriu_parells(val_pairs, 'val')

    # Fitxer TSV complet per a revisió
    tsv_path = output_dir / 'corpus_net.tsv'
    with open(tsv_path, 'w', encoding='utf-8') as f:
        f.write('castellà\tvalencià\n')
        for src, tgt in clean_pairs:
            f.write(f'{src}\t{tgt}\n')

    # ── Construeix la secció de desalineaments ─────────────────────────────────
    total_parells_doc = total_fitxers // 2
    n_desalineats  = len(desalineaments)
    n_alineats     = total_parells_doc - n_desalineats - total_exclosos
    n_exclosos_des = sum(1 for d in desalineaments if d['exclòs'])

    des_lines = [
        '',
        '  ' + '═' * 51,
        '  DESALINEAMENTS DETECTATS',
        '  ' + '═' * 51,
        f'  Documents processats (parells):    {total_parells_doc}',
        f'  Documents perfectament alineats:   {n_alineats}',
        f'  Documents amb desalineament (Δ>0): {n_desalineats}',
        f'  Documents exclosos (Δ>{max_delta}):' + ' ' * max(1, 16 - len(str(max_delta))) +
            f'{n_exclosos_des}',
        '',
    ]

    if desalineaments:
        # Capçalera de la taula
        des_lines += [
            '  Llista de documents desalineats (ordenats per Δ descendent):',
            '',
            f'  {"Tipus":<5}  {"CAS":>5}  {"VAL":>5}  {"Δ":>5}  {"Usats":>6}  {"Estat":<8}  Fitxer',
            '  ' + '─' * 100,
        ]
        ordenats = sorted(desalineaments, key=lambda d: d['delta'], reverse=True)
        for d in ordenats:
            estat = 'EXCLÒS' if d.get('exclòs') else 'truncat'
            nom = d.get('fitxer', d.get('fitxer_cas', ''))
            usats = d.get('parells_extrets', d.get('parells_usats', 0))
            des_lines.append(
                f'  {d["tipus"]:<5}  {d["paras_cas"]:>5}  {d["paras_val"]:>5}'
                f'  {d["delta"]:>5}  {usats:>6}  {estat:<8}  {nom}'
            )

        # Top 10
        des_lines += [
            '',
            '  Top 10 desalineaments majors:',
            '',
            f'  {"Tipus":<5}  {"CAS":>5}  {"VAL":>5}  {"Δ":>5}  {"Usats":>6}  Fitxer',
            '  ' + '─' * 80,
        ]
        for d in ordenats[:10]:
            nom = d.get('fitxer', d.get('fitxer_cas', ''))
            usats = d.get('parells_extrets', d.get('parells_usats', 0))
            des_lines.append(
                f'  {d["tipus"]:<5}  {d["paras_cas"]:>5}  {d["paras_val"]:>5}'
                f'  {d["delta"]:>5}  {usats:>6}  {nom}'
            )
    else:
        des_lines.append('  ✓ Cap desalineament detectat. Tots els documents estan perfectament alineats.')

    des_lines += ['', '  ' + '═' * 51]

    # ── Informe de processament complet ───────────────────────────────────────
    informe_lines = [
        '=' * 55,
        '  INFORME DE PROCESSAMENT — corpus_builder.py',
        '  Servei de Llengües i Política Lingüística · UV',
        '=' * 55,
        '',
        f'  Carpeta entrada:  {input_dir}',
        f'  Carpeta sortida:  {output_dir}',
        f'  Max delta:        {max_delta}',
        f'  Min similitud:    {args.min_similitud:.2f}',
        f'  Bertalign:        {"disponible (LaBSE)" if _bertalign_ok else "NO disponible (fallback posicional)"}',
        '',
        '  FITXERS PROCESSATS:',
    ]
    for linia in registre_fitxers:
        informe_lines.append(f'    {linia}')
    if not registre_fitxers:
        informe_lines.append('    (cap fitxer processat)')
    informe_lines += [
        '',
        '  ESTADÍSTIQUES:',
        f'    Fitxers processats:  {total_fitxers}',
        f'    Documents exclosos:  {total_exclosos}',
        f'    Parells extrets:     {total_extrets}',
        f'    Parells eliminats:   {total_eliminats}',
        f'      (baixa similitud): {eliminats_similitud}',
        f'    Parells vàlids:      {total_valid}',
        f'    Parells train (90%): {len(train_pairs)}',
        f'    Parells val   (10%): {len(val_pairs)}',
        '',
        '  FITXERS GENERATS:',
        f'    train.es / train.ca  ({len(train_pairs)} parells)',
        f'    val.es   / val.ca    ({len(val_pairs)} parells)',
        f'    corpus_net.tsv       ({total_valid} parells + capçalera)',
        '',
    ]
    if total_valid < 500:
        informe_lines += [
            '  ⚠ ADVERTÈNCIA: corpus molt petit (< 500 parells).',
            '    Recomana acumular més textos abans de fer l\'afinament.',
            '',
        ]
    informe_lines.append('=' * 55)

    # Afegeix la secció de desalineaments
    informe_lines.extend(des_lines)

    informe_text = '\n'.join(informe_lines)
    informe_path = output_dir / 'informe.txt'
    with open(informe_path, 'w', encoding='utf-8') as f:
        f.write(informe_text + '\n')

    # ── Mostra resum per consola ───────────────────────────────────────────────
    print()
    print(informe_text)
    print()
    log.info('Fitxers desats a: %s', output_dir)

    if total_valid < 500:
        sys.exit(1)


if __name__ == '__main__':
    main()
